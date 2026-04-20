#!/usr/bin/env python3
"""
Create sample documents to test the enhanced preview system.
"""

import os
from pathlib import Path

def create_sample_documents():
    """Create sample PDF, Excel, and Word documents for testing."""
    
    test_dir = Path("preview_test_files")
    test_dir.mkdir(exist_ok=True)
    
    print("📄 Creating sample documents for enhanced preview...")
    
    # Create a sample Excel file
    try:
        import openpyxl
        from openpyxl import Workbook
        
        wb = Workbook()
        ws = wb.active
        ws.title = "Sample Data"
        
        # Add sample data
        headers = ["Name", "Age", "City", "Salary"]
        ws.append(headers)
        
        data = [
            ["Alice Johnson", 28, "New York", 75000],
            ["Bob Smith", 35, "Los Angeles", 82000],
            ["Carol Davis", 31, "Chicago", 68000],
            ["David Wilson", 42, "Houston", 91000],
            ["Eva Brown", 26, "Phoenix", 62000],
            ["Frank Miller", 39, "Philadelphia", 78000],
            ["Grace Taylor", 33, "San Antonio", 71000],
            ["Henry Anderson", 45, "San Diego", 95000],
            ["Ivy Thomas", 29, "Dallas", 67000],
            ["Jack Jackson", 37, "San Jose", 88000],
            ["Kate White", 32, "Austin", 73000],
            ["Liam Harris", 41, "Jacksonville", 86000],
            ["Mia Martin", 27, "Fort Worth", 64000],
            ["Noah Thompson", 36, "Columbus", 79000],
            ["Olivia Garcia", 30, "Charlotte", 70000]
        ]
        
        for row in data:
            ws.append(row)
        
        # Add some formatting
        for col in range(1, 5):
            ws.cell(row=1, column=col).font = openpyxl.styles.Font(bold=True)
        
        excel_file = test_dir / "sample_data.xlsx"
        wb.save(excel_file)
        print(f"✅ Created Excel file: {excel_file}")
        
    except ImportError:
        print("❌ openpyxl not available - skipping Excel file creation")
    
    # Create a sample Word document
    try:
        from docx import Document
        
        doc = Document()
        
        # Add title
        doc.add_heading('File Scout Enhanced Preview Report', 0)
        
        # Add introduction
        doc.add_paragraph('This is a sample Word document created to demonstrate the enhanced preview capabilities of File Scout 3.2.')
        
        # Add sections
        doc.add_heading('Features Demonstrated', level=1)
        doc.add_paragraph('The enhanced preview system supports multiple document types including:')
        
        features = [
            'Word documents with paragraph extraction',
            'Excel spreadsheets with data grid preview',
            'PDF files with text content extraction',
            'PowerPoint presentations with slide text',
            'Audio files with metadata display',
            'Code files with syntax highlighting',
            'Archive files with content listing'
        ]
        
        for feature in features:
            doc.add_paragraph(f'• {feature}', style='List Bullet')
        
        # Add technical details
        doc.add_heading('Technical Implementation', level=1)
        doc.add_paragraph('The preview system uses a modular handler architecture with lazy loading of dependencies. Each file type has a dedicated handler that extracts relevant information for display.')
        
        # Add conclusion
        doc.add_heading('Conclusion', level=1)
        doc.add_paragraph('This enhanced preview system transforms File Scout from a basic file finder into a comprehensive file inspection tool, providing users with immediate insights into file contents without requiring external applications.')
        
        word_file = test_dir / "sample_document.docx"
        doc.save(word_file)
        print(f"✅ Created Word file: {word_file}")
        
    except ImportError:
        print("❌ python-docx not available - skipping Word file creation")
    
    # Create a sample PowerPoint presentation
    try:
        from pptx import Presentation
        
        prs = Presentation()
        
        # Title slide
        slide_layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "File Scout 3.2"
        subtitle.text = "Enhanced Preview System Demonstration"
        
        # Overview slide
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Preview Features"
        content = slide.placeholders[1]
        content.text = """• Document Preview (PDF, Word, Excel, PowerPoint)
• Code Syntax Highlighting
• Media Information (Audio/Video)
• Archive Content Listing
• Binary File Hex View
• Enhanced Properties Tab"""
        
        # Technical details slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Technical Architecture"
        content = slide.placeholders[1]
        content.text = """• Modular Handler System
• Lazy Loading Dependencies
• Performance Optimizations
• Memory Management
• Error Handling & Fallbacks"""
        
        # Benefits slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "User Benefits"
        content = slide.placeholders[1]
        content.text = """• No External Applications Needed
• Instant File Content Preview
• 50+ File Type Support
• Professional-Grade Features
• Large Directory Performance"""
        
        pptx_file = test_dir / "sample_presentation.pptx"
        prs.save(pptx_file)
        print(f"✅ Created PowerPoint file: {pptx_file}")
        
    except ImportError:
        print("❌ python-pptx not available - skipping PowerPoint file creation")
    
    print(f"\n🎉 Sample documents created in: {test_dir.absolute()}")
    print("\nTo test the enhanced preview:")
    print("1. Run File Scout 3.2")
    print("2. Search in the preview_test_files directory")
    print("3. Select the document files to see enhanced previews!")
    
    return test_dir

if __name__ == "__main__":
    create_sample_documents()
