import csv
import zipfile
from pathlib import Path

# Optional preview dependencies — gated per-handler, fail gracefully
try:
    from pygments import highlight
    from pygments.lexers import get_lexer_for_filename
    from pygments.formatters import HtmlFormatter
    PYGMENTS_AVAILABLE = True
except ImportError:
    PYGMENTS_AVAILABLE = False

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import mutagen
    MUTAGEN_AVAILABLE = True
except ImportError:
    MUTAGEN_AVAILABLE = False

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    import xlrd
    XLRD_AVAILABLE = True
except ImportError:
    XLRD_AVAILABLE = False


class PreviewHandler:
    """Base class for file preview handlers."""

    def __init__(self, name, extensions=None):
        self.name = name
        self.extensions = extensions or []

    def can_handle(self, file_path):
        file_path = Path(file_path)
        return file_path.suffix.lower() in self.extensions

    def generate_preview(self, file_path, max_size=1024*1024):
        raise NotImplementedError


class TextPreviewHandler(PreviewHandler):

    def __init__(self):
        super().__init__("Text", ['.txt', '.log', '.cfg', '.ini', '.conf', '.md', '.rst'])

    def generate_preview(self, file_path, max_size=1024*1024):
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read(min(max_size, 5000))

                if PYGMENTS_AVAILABLE:
                    try:
                        lexer = get_lexer_for_filename(str(file_path))
                        formatter = HtmlFormatter(style='default', linenos=True)
                        html_content = highlight(content, lexer, formatter)
                        return ("html", html_content, {"syntax_highlighted": True})
                    except:
                        return ("text", content, {"syntax_highlighted": False})
                else:
                    return ("text", content, {"syntax_highlighted": False})
        except Exception as e:
            return ("error", str(e), {})


class CodePreviewHandler(PreviewHandler):

    def __init__(self):
        super().__init__("Code", [
            '.py', '.js', '.html', '.css', '.java', '.cpp', '.c', '.h', '.hpp',
            '.cs', '.php', '.rb', '.go', '.rs', '.swift', '.kt', '.scala',
            '.sh', '.bat', '.ps1', '.sql', '.xml', '.json', '.yaml', '.yml',
            '.toml', '.dockerfile', '.gitignore', '.env', '.htaccess'
        ])

    def generate_preview(self, file_path, max_size=1024*1024):
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read(min(max_size, 5000))

                if PYGMENTS_AVAILABLE:
                    try:
                        lexer = get_lexer_for_filename(str(file_path))
                        formatter = HtmlFormatter(style='monokai', linenos=True)
                        html_content = highlight(content, lexer, formatter)
                        return ("html", html_content, {"syntax_highlighted": True, "theme": "monokai"})
                    except:
                        return ("text", content, {"syntax_highlighted": False})
                else:
                    return ("text", content, {"syntax_highlighted": False})
        except Exception as e:
            return ("error", str(e), {})


class PDFPreviewHandler(PreviewHandler):

    def __init__(self):
        super().__init__("PDF", ['.pdf'])

    def generate_preview(self, file_path, max_size=1024*1024):
        if not PYMUPDF_AVAILABLE:
            return ("error", "PyMuPDF not installed. Install with: pip install PyMuPDF", {})

        try:
            doc = fitz.open(str(file_path))
            if doc.page_count > 0:
                all_text = ""
                formatted_pages = []
                page_count = doc.page_count

                for page_num in range(min(3, page_count)):
                    page = doc[page_num]
                    text = page.get_text()
                    if text.strip():
                        all_text += f"\n--- Page {page_num + 1} ---\n{text}"

                        formatted_text = f"═══════════════════════════════════════\n"
                        formatted_text += f"          PAGE {page_num + 1} of {page_count}\n"
                        formatted_text += f"═══════════════════════════════════════\n\n"
                        formatted_text += text
                        formatted_pages.append(formatted_text)

                doc.close()

                if all_text.strip():
                    simple_text = all_text[:3000]
                    formatted_text = '\n\n'.join(formatted_pages)

                    return ("pdf_dual", simple_text, {
                        "formatted": formatted_text,
                        "pages": page_count,
                        "preview_pages": min(3, page_count),
                        "file_path": str(file_path)
                    })
                else:
                    return ("pdf_dual", "PDF contains no extractable text (image-based PDF)\n\nUse the PDF Viewer tab to see the visual content.", {
                        "formatted": "Image-based PDF - No extractable text available\n\nUse the PDF Viewer tab to see the visual content.",
                        "pages": page_count,
                        "preview_pages": min(3, page_count),
                        "file_path": str(file_path)
                    })
            else:
                doc.close()
                return ("error", "PDF appears to be empty", {})
        except Exception as e:
            return ("error", f"PDF Error: {e}", {})


class ExcelPreviewHandler(PreviewHandler):

    def __init__(self):
        super().__init__("Excel", ['.xlsx', '.xlsm', '.xls'])

    def generate_preview(self, file_path, max_size=1024*1024):
        file_path = Path(file_path)
        extension = file_path.suffix.lower()

        try:
            if extension == '.xls':
                if not XLRD_AVAILABLE:
                    return ("error", "xlrd not installed for .xls files. Install with: pip install xlrd", {})
                return self._read_xls_file(file_path)
            else:
                if not OPENPYXL_AVAILABLE:
                    return ("error", "openpyxl not installed for .xlsx/.xlsm files. Install with: pip install openpyxl", {})
                return self._read_xlsx_file(file_path)

        except Exception as e:
            return ("error", str(e), {})

    def _read_xls_file(self, file_path):
        try:
            wb = xlrd.open_workbook(str(file_path))
            sheet = wb.sheet_by_index(0)

            content = []
            for row_idx in range(min(10, sheet.nrows)):
                row_data = []
                for col_idx in range(min(10, sheet.ncols)):
                    cell = sheet.cell(row_idx, col_idx)

                    if cell.ctype == xlrd.XL_CELL_NUMBER:
                        if 0 < cell.value < 1 and sheet.book.datemode:
                            try:
                                date_tuple = xlrd.xldate_as_tuple(cell.value, sheet.book.datemode)
                                if date_tuple[0] > 0:
                                    cell_value = f"{date_tuple[0]}-{date_tuple[1]:02d}-{date_tuple[2]:02d}"
                                else:
                                    cell_value = str(cell.value)
                            except:
                                cell_value = str(cell.value)
                        else:
                            cell_value = str(cell.value)
                    elif cell.ctype == xlrd.XL_CELL_TEXT:
                        cell_value = str(cell.value)
                    elif cell.ctype == xlrd.XL_CELL_BLANK:
                        cell_value = ''
                    else:
                        cell_value = str(cell.value)

                    row_data.append(cell_value)

                if any(cell.strip() for cell in row_data if cell):
                    content.append('\t'.join(row_data))

            if content:
                return ("text", '\n'.join(content), {"sheets": wb.nsheets, "active_sheet": sheet.name, "format": "XLS"})
            else:
                return ("text", "Excel file appears to be empty", {"sheets": wb.nsheets, "format": "XLS"})

        except Exception as e:
            return ("error", f"Error reading .xls file: {e}", {})

    def _read_xlsx_file(self, file_path):
        try:
            wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
            if wb.sheetnames:
                sheet = wb.active
                content = []
                for row in sheet.iter_rows(max_row=10, max_col=10, values_only=True):
                    if any(cell is not None for cell in row):
                        content.append('\t'.join(str(cell) if cell is not None else '' for cell in row))

                wb.close()
                if content:
                    return ("text", '\n'.join(content), {"sheets": len(wb.sheetnames), "active_sheet": sheet.title, "format": "XLSX"})
                else:
                    return ("text", "Excel file appears to be empty", {"sheets": len(wb.sheetnames), "format": "XLSX"})
        except Exception as e:
            return ("error", f"Error reading .xlsx/.xlsm file: {e}", {})


class CSVPreviewHandler(PreviewHandler):

    def __init__(self):
        super().__init__("CSV", ['.csv'])

    def generate_preview(self, file_path, max_size=1024*1024):
        try:
            content = []
            row_count = 0
            max_rows = 20
            max_cols = 15

            with open(str(file_path), 'r', encoding='utf-8', errors='replace') as f:
                sample = f.read(1024)
                f.seek(0)

                sniffer = csv.Sniffer()
                try:
                    delimiter = sniffer.sniff(sample).delimiter
                except:
                    delimiter = ','

                reader = csv.reader(f, delimiter=delimiter)

                for row in reader:
                    if row_count >= max_rows:
                        break

                    formatted_row = []
                    for i, cell in enumerate(row[:max_cols]):
                        cell_str = str(cell).strip()
                        if len(cell_str) > 50:
                            cell_str = cell_str[:47] + "..."
                        formatted_row.append(cell_str)

                    if len(row) > max_cols:
                        formatted_row.append("...")

                    if any(cell.strip() for cell in formatted_row):
                        content.append('\t'.join(formatted_row))
                        row_count += 1

            if content:
                metadata = {
                    "delimiter": delimiter,
                    "rows_shown": row_count,
                    "max_columns": max_cols,
                    "format": "CSV"
                }

                if content:
                    header_info = f"CSV Preview (delimiter: '{delimiter}') - First {row_count} rows, max {max_cols} columns\n"
                    header_info += "\u2500" * 60 + "\n"
                    return ("text", header_info + '\n'.join(content), metadata)
                else:
                    return ("text", '\n'.join(content), metadata)
            else:
                return ("text", "CSV file appears to be empty or contains only empty rows", {"format": "CSV"})

        except Exception as e:
            return ("error", f"Error reading CSV file: {e}", {})


class WordPreviewHandler(PreviewHandler):

    def __init__(self):
        super().__init__("Word", ['.docx'])

    def generate_preview(self, file_path, max_size=1024*1024):
        if not DOCX_AVAILABLE:
            return ("error", "python-docx not installed. Install with: pip install python-docx", {})

        try:
            doc = Document(str(file_path))
            paragraphs = []
            for para in doc.paragraphs[:20]:
                if para.text.strip():
                    paragraphs.append(para.text)

            if paragraphs:
                return ("text", '\n\n'.join(paragraphs), {"paragraphs": len(doc.paragraphs)})
            else:
                return ("text", "Document appears to be empty", {"paragraphs": len(doc.paragraphs)})
        except Exception as e:
            return ("error", str(e), {})


class PowerPointPreviewHandler(PreviewHandler):

    def __init__(self):
        super().__init__("PowerPoint", ['.pptx'])

    def generate_preview(self, file_path, max_size=1024*1024):
        if not PPTX_AVAILABLE:
            return ("error", "python-pptx not installed. Install with: pip install python-pptx", {})

        try:
            prs = Presentation(str(file_path))
            slide_texts = []
            for i, slide in enumerate(prs.slides[:10]):
                text_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text)
                if text_content:
                    slide_texts.append(f"Slide {i+1}:\n" + '\n'.join(text_content))

            if slide_texts:
                return ("text", '\n\n---\n\n'.join(slide_texts), {"slides": len(prs.slides)})
            else:
                return ("text", f"Presentation has {len(prs.slides)} slides but no extractable text", {"slides": len(prs.slides)})
        except Exception as e:
            return ("error", str(e), {})


class AudioPreviewHandler(PreviewHandler):

    def __init__(self):
        super().__init__("Audio", ['.mp3', '.flac', '.ogg', '.wav', '.m4a', '.aac'])

    def generate_preview(self, file_path, max_size=1024*1024):
        if not MUTAGEN_AVAILABLE:
            return ("error", "mutagen not installed. Install with: pip install mutagen", {})

        try:
            audio = mutagen.File(str(file_path))
            if audio is not None:
                metadata = {}
                for key, value in audio.items():
                    if isinstance(value, list) and value:
                        metadata[key] = str(value[0])
                    else:
                        metadata[key] = str(value)

                info_lines = [f"Audio File Metadata:"]
                info_lines.extend([f"{k}: {v}" for k, v in metadata.items()])
                return ("text", '\n'.join(info_lines), {"metadata_keys": len(metadata)})
            else:
                return ("text", "Unable to read audio metadata", {})
        except Exception as e:
            return ("error", str(e), {})


class VideoPreviewHandler(PreviewHandler):

    def __init__(self):
        super().__init__("Video", ['.mp4', '.avi', '.mkv', '.mov', '.wmv', '.flv', '.webm'])

    def generate_preview(self, file_path, max_size=1024*1024):
        try:
            file_path = Path(file_path)
            size_mb = file_path.stat().st_size / (1024 * 1024)

            info_lines = [
                f"Video File Information:",
                f"File: {file_path.name}",
                f"Size: {size_mb:.2f} MB",
                f"Extension: {file_path.suffix}",
                "",
                "Note: Detailed video metadata requires ffmpeg or similar tools."
            ]

            return ("text", '\n'.join(info_lines), {"size_mb": size_mb})
        except Exception as e:
            return ("error", str(e), {})


class ArchivePreviewHandler(PreviewHandler):

    def __init__(self):
        super().__init__("Archive", ['.zip', '.rar', '.7z', '.tar', '.gz'])

    def generate_preview(self, file_path, max_size=1024*1024):
        try:
            if Path(file_path).suffix.lower() == '.zip':
                with zipfile.ZipFile(str(file_path), 'r') as zf:
                    file_list = zf.namelist()[:50]
                    total_files = len(zf.namelist())
                    total_size = sum(f.file_size for f in zf.filelist)

                    info_lines = [
                        f"ZIP Archive Contents:",
                        f"Total files: {total_files}",
                        f"Total size: {total_size / (1024*1024):.2f} MB",
                        "",
                        "First 50 files:"
                    ]
                    info_lines.extend(file_list)

                    return ("text", '\n'.join(info_lines), {"total_files": total_files, "total_size_mb": total_size / (1024*1024)})
            else:
                return ("text", f"Archive preview supported for ZIP files only. File: {Path(file_path).name}", {})
        except Exception as e:
            return ("error", str(e), {})


class HexPreviewHandler(PreviewHandler):

    def __init__(self):
        super().__init__("Hex", ['.exe', '.dll', '.bin', '.dat', '.img', '.iso'])

    def generate_preview(self, file_path, max_size=1024*1024):
        try:
            with open(file_path, 'rb') as f:
                data = f.read(min(max_size, 512))

                hex_lines = []
                for i in range(0, len(data), 16):
                    chunk = data[i:i+16]
                    hex_part = ' '.join(f'{b:02x}' for b in chunk)
                    ascii_part = ''.join(chr(b) if 32 <= b <= 126 else '.' for b in chunk)
                    hex_lines.append(f'{i:08x}  {hex_part:<48} |{ascii_part}|')

                return ("text", '\n'.join(hex_lines), {"bytes_previewed": len(data)})
        except Exception as e:
            return ("error", str(e), {})
