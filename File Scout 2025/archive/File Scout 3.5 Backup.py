"""
File Scout
----------
A powerful and intuitive utility for auditing and finding files across directories.
Features advanced filtering, duplicate finding, interactive results, and profile management.
"""

import os
import sys
import re
import datetime
import json
import csv
import subprocess
import time
from pathlib import Path
import mimetypes
import hashlib
import argparse
from collections import defaultdict
import concurrent.futures
import shutil # <-- ADDED IMPORT
import zipfile
import xml.etree.ElementTree as ET
import base64
import struct

# Enhanced preview imports (lazy loaded)
try:
    import pygments
    from pygments import highlight
    from pygments.lexers import get_lexer_for_filename, TextLexer
    from pygments.formatters import HtmlFormatter
    PYGMENTS_AVAILABLE = True
except ImportError:
    PYGMENTS_AVAILABLE = False

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import mutagen
    MUTAGEN_AVAILABLE = True
except ImportError:
    MUTAGEN_AVAILABLE = False

# Core dependencies (always available)
try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    import xlrd
    XLRD_AVAILABLE = True
except ImportError:
    XLRD_AVAILABLE = False
try:
    from send2trash import send2trash
except Exception:
    send2trash = None

# Quiet benign Qt logs and stabilize HiDPI behavior before importing PyQt6 modules
os.environ.setdefault("QT_LOGGING_RULES", "qt.qpa.screen=false;qt.svg.warning=false")
os.environ.setdefault("QT_SCALE_FACTOR_ROUNDING_POLICY", "RoundPreferFloor")

from PyQt6.QtWidgets import (QApplication, QMainWindow, QWidget, QVBoxLayout,
                             QHBoxLayout, QLabel, QLineEdit, QPushButton,
                             QFileDialog, QProgressBar, QTableWidget, QComboBox,
                             QCheckBox, QGroupBox, QGridLayout, QMessageBox,
                             QListWidget, QSplitter, QDateEdit, QMenu,
                             QInputDialog, QTableWidgetItem, QHeaderView,
                             QDialog, QDialogButtonBox, QListWidgetItem, QWidgetAction,
                             QStatusBar, QTextEdit, QTabWidget, QSystemTrayIcon, QScrollArea)
from PyQt6.QtCore import Qt, QThread, pyqtSignal, QSize, QDate, QSettings, QObject, QTimer, QEvent, QPoint, QRect
from PyQt6.QtGui import QIcon, QColor, QAction, QActionGroup, QBrush, QPixmap, QGuiApplication, QFont, QImage, QPainter, QPen

# File Audit Dialog (Google Drive integration)
try:
    from file_audit_dialog import FileAuditDialog
    FILE_AUDIT_AVAILABLE = True
except ImportError:
    FILE_AUDIT_AVAILABLE = False

# --- Constants and Utilities
APP_NAME = "File Scout"
APP_VERSION = "3.2" # Version bump for new feature
SETTINGS_ORG = "WindsurfAI"
SETTINGS_APP = "FileScout"

# Performance limits for large directories
MAX_RESULTS = 50000  # Maximum number of results to prevent memory issues
MAX_SCAN_FILES = 1000000  # Maximum files to scan before stopping
LARGE_DIR_THRESHOLD = 100000  # Threshold for considering a directory "large"

class DropLineEdit(QLineEdit):
    """A QLineEdit that accepts drag-and-drop for directory paths."""
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setAcceptDrops(True)

    def dragEnterEvent(self, event):
        if event.mimeData().hasUrls() and len(event.mimeData().urls()) == 1:
            url = event.mimeData().urls()[0]
            if url.isLocalFile() and os.path.isdir(url.toLocalFile()):
                event.acceptProposedAction()

    def dropEvent(self, event):
        url = event.mimeData().urls()[0]
        self.setText(url.toLocalFile())

class ExcelExporter:
    """Enhanced Excel exporter with formatting and multiple export options."""
    def __init__(self, theme='light'):
        self.theme = theme
        if theme == 'dark':
            self.header_bg = '1F1F1F'
            self.header_fg = 'FFFFFF'
            self.alt_row_bg = '2A2A2A'
        else:
            self.header_bg = '4472C4'
            self.header_fg = 'FFFFFF'
            self.alt_row_bg = 'E6E6E6'

    def export_data(self, file_path, headers, data, sheet_name="File Scout Results"):
        """Export data to Excel with formatting."""
        try:
            import openpyxl
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            from openpyxl.utils import get_column_letter

            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = sheet_name

            for col, header in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col, value=header)
                cell.font = Font(bold=True, color=self.header_fg)
                cell.fill = PatternFill(start_color=self.header_bg, end_color=self.header_bg, fill_type="solid")
                cell.alignment = Alignment(horizontal='center', vertical='center')

            thin_border = Border(left=Side(style='thin'), right=Side(style='thin'), top=Side(style='thin'), bottom=Side(style='thin'))

            for row_idx, row_data in enumerate(data, 2):
                row_fill = PatternFill(start_color=self.alt_row_bg, end_color=self.alt_row_bg, fill_type="solid") if row_idx % 2 == 0 else None
                for col_idx, header in enumerate(headers, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=row_data.get(header, ""))
                    cell.border = thin_border
                    if row_fill:
                        cell.fill = row_fill

            for col in ws.columns:
                max_length = 0
                column_letter = get_column_letter(col[0].column)
                for cell in col:
                    if cell.value:
                        max_length = max(max_length, len(str(cell.value)))
                ws.column_dimensions[column_letter].width = min(max_length + 2, 50)

            ws.freeze_panes = 'A2'
            wb.save(file_path)
            return True
        except ImportError:
            raise ImportError("The 'openpyxl' library is required for Excel export.\nPlease install it: pip install openpyxl")
        except Exception as e:
            raise Exception(f"Failed to export to Excel: {str(e)}")

ICON_SEARCH = "PHN2ZyB4bWxucz0iaHR0cDovL3d3dy53My5vcmcvMjAwMC9zdmciIHZpZXdCb3g9IjAgMCAyNCAyNCIgZmlsbD0iI2ZmZmZmZiI+PHBhdGggZD0iTTE1LjUgMTRoLS43OWwtLjI4LS4yN0E2LjQ3MSA2LjQ3MSAwIDAgMCAxNiA5LjVDMTYgNS45MSAxMy4wOSAzIDkuNSAzUzMgNS45MSAzIDkuNSA1LjkxIDE2IDkuNSAxNmMxLjQzIDAgMi43Ni0uNDcgMy45MS0xLjI1bC4yNy4yOHYuNzlsNSA0Ljk5TDIwLjQ9IDE5bC00Ljk5LTV6bS02IDBDNy4wMSAxNCA1IDExLjk5IDUgOS41UzcuMDEgNSA5LjUgNSAxNCA3LjAxIDE0IDkuNSAxMS45OSAxNCA5LjUgMTR6Ii8+PC9zdmc+"
ICON_STOP = "PHN2ZyB4bWxucz0iaHR0cDovL3d3dy53My5vcmcvMjAwMC9zdmciIHZpZXdCb3g9IjAgMCAyNCAyNCIgZmlsbD0iI2ZmZmZmZiI+PHBhdGggZD0iTTYgNmgxMnYxMkg2eiIvPjwvc3ZnPg=="
ICON_EXPORT = "PHN2ZyB4bWxucz0iaHR0cDovL3d3dy53My5vcmcvMjAwMC9zdmciIHZpZXdCb3g9IjAgMCAyNCAyNCIgZmlsbD0iI2ZmZmZmZiI+PHBhdGggZD0iTTE5IDEydjdINVYzdjloMlY5aDEwbDd2NGgtMnpNOCA1djhoOFY5SDh6bS0yIDE2aDEydjJIMFYxMmg0djlaIi8+PC9zdmc+"

THEME_COLORS = {
    'light': {
        'primary': '#4a90e2',
        'secondary': '#7ed321',
        'accent': '#f5a623',
        'background': '#f0f0f0',
        'surface': '#ffffff',
        'text': '#000000',
        'text_secondary': '#666666',
        'border': '#cccccc',
        'header': '#e8e8e8',
        'grid': '#dddddd',
        'menu': '#f8f8f8',
        'hover': '#e6f3ff'
    },
    'dark': {
        'primary': '#4a90e2',
        'secondary': '#7ed321',
        'accent': '#f5a623',
        'background': '#2d2d30',
        'surface': '#3e3e42',
        'text': '#f1f1f1',
        'text_secondary': '#cccccc',
        'border': '#555555',
        'header': '#3e3e42',
        'grid': '#4a4a4a',
        'menu': '#3e3e42',
        'hover': '#555555'
    },
    'cyberpunk': {
        'primary': '#00ffff',
        'secondary': '#ff00ff',
        'accent': '#ffff00',
        'background': '#0a0a0a',
        'surface': '#1a1a1a',
        'text': '#00ff41',
        'text_secondary': '#00ccff',
        'border': '#ff0080',
        'header': '#2a2a2a',
        'grid': '#333333',
        'menu': '#1a1a1a',
        'hover': '#003d4d'
    },
    'halloween': {
        'primary': '#ff6600',
        'secondary': '#9900cc',
        'accent': '#ffaa00',
        'background': '#1a0d00',
        'surface': '#2d1a00',
        'text': '#ff9933',
        'text_secondary': '#cc6600',
        'border': '#663300',
        'header': '#331a00',
        'grid': '#4d2600',
        'menu': '#2d1a00',
        'hover': '#4d2600'
    },
    'ocean': {
        'primary': '#0077be',
        'secondary': '#00a693',
        'accent': '#40e0d0',
        'background': '#f0f8ff',
        'surface': '#ffffff',
        'text': '#003366',
        'text_secondary': '#0066cc',
        'border': '#87ceeb',
        'header': '#e6f3ff',
        'grid': '#b0e0e6',
        'menu': '#f0f8ff',
        'hover': '#cce6ff'
    },
    'forest': {
        'primary': '#228b22',
        'secondary': '#32cd32',
        'accent': '#90ee90',
        'background': '#f5f5dc',
        'surface': '#ffffff',
        'text': '#2f4f2f',
        'text_secondary': '#556b2f',
        'border': '#9acd32',
        'header': '#f0fff0',
        'grid': '#d3d3d3',
        'menu': '#f5f5dc',
        'hover': '#e6ffe6'
    },
    'sunset': {
        'primary': '#ff4500',
        'secondary': '#ff69b4',
        'accent': '#ffd700',
        'background': '#fff8dc',
        'surface': '#ffffff',
        'text': '#8b0000',
        'text_secondary': '#cd5c5c',
        'border': '#dda0dd',
        'header': '#ffe4e1',
        'grid': '#f0e68c',
        'menu': '#fff8dc',
        'hover': '#ffe4b5'
    },
    'retro': {
        'primary': '#8b4513',
        'secondary': '#d2691e',
        'accent': '#daa520',
        'background': '#f5deb3',
        'surface': '#fff8dc',
        'text': '#654321',
        'text_secondary': '#a0522d',
        'border': '#cd853f',
        'header': '#f4e4bc',
        'grid': '#deb887',
        'menu': '#f5deb3',
        'hover': '#f0e68c'
    },
    'pastel': {
        'primary': '#ffb6c1',
        'secondary': '#98fb98',
        'accent': '#f0e68c',
        'background': '#fafafa',
        'surface': '#ffffff',
        'text': '#696969',
        'text_secondary': '#808080',
        'border': '#e6e6fa',
        'header': '#f5f5f5',
        'grid': '#f0f0f0',
        'menu': '#fafafa',
        'hover': '#f8f8ff'
    },
    'high_contrast': {
        'primary': '#ffff00',
        'secondary': '#00ff00',
        'accent': '#ff00ff',
        'background': '#000000',
        'surface': '#1a1a1a',
        'text': '#ffffff',
        'text_secondary': '#ffff00',
        'border': '#ffffff',
        'header': '#333333',
        'grid': '#666666',
        'menu': '#1a1a1a',
        'hover': '#333333'
    },
    'corporate_blue': {
        'primary': '#003f7f',
        'secondary': '#0066cc',
        'accent': '#4da6ff',
        'background': '#f7f9fc',
        'surface': '#ffffff',
        'text': '#1a1a1a',
        'text_secondary': '#666666',
        'border': '#c0d6e8',
        'header': '#e8f1fa',
        'grid': '#dde8f0',
        'menu': '#f7f9fc',
        'hover': '#e1ecf7'
    },
    'crimson': {
        'primary': '#dc143c',
        'secondary': '#b22222',
        'accent': '#ff6347',
        'background': '#fff5f5',
        'surface': '#ffffff',
        'text': '#4b0000',
        'text_secondary': '#8b0000',
        'border': '#f08080',
        'header': '#ffe4e1',
        'grid': '#f5c2c7',
        'menu': '#fff5f5',
        'hover': '#ffe4e1'
    },
    'mint': {
        'primary': '#00a86b',
        'secondary': '#40e0d0',
        'accent': '#7fffd4',
        'background': '#f0fff0',
        'surface': '#ffffff',
        'text': '#006400',
        'text_secondary': '#2e8b57',
        'border': '#90ee90',
        'header': '#f0fff0',
        'grid': '#e0ffe0',
        'menu': '#f0fff0',
        'hover': '#e6ffe6'
    },
    'space': {
        'primary': '#4b0082',
        'secondary': '#9370db',
        'accent': '#da70d6',
        'background': '#0b0b2f',
        'surface': '#1a1a3a',
        'text': '#e6e6fa',
        'text_secondary': '#b19cd9',
        'border': '#483d8b',
        'header': '#2e2e5e',
        'grid': '#3a3a6a',
        'menu': '#1a1a3a',
        'hover': '#3e3e7e'
    },
    'nord': {
        'primary': '#81a1c1',
        'secondary': '#88c0d0',
        'accent': '#8fbcbb',
        'background': '#2e3440',
        'surface': '#3b4252',
        'text': '#eceff4',
        'text_secondary': '#d8dee9',
        'border': '#4c566a',
        'header': '#3b4252',
        'grid': '#434c5e',
        'menu': '#2e3440',
        'hover': '#4c566a'
    },
    'mocha': {
        'primary': '#a67b5b',
        'secondary': '#856d4d',
        'accent': '#c89f7f',
        'background': '#2a1d13',
        'surface': '#3d2c1d',
        'text': '#fff5e6',
        'text_secondary': '#d4bda5',
        'border': '#543c2a',
        'header': '#3d2c1d',
        'grid': '#543c2a',
        'menu': '#2a1d13',
        'hover': '#543c2a'
    },
    'solarized_dark': {
        'primary': '#268bd2',
        'secondary': '#2aa198',
        'accent': '#b58900',
        'background': '#002b36',
        'surface': '#073642',
        'text': '#93a1a1',
        'text_secondary': '#839496',
        'border': '#586e75',
        'header': '#073642',
        'grid': '#073642',
        'menu': '#002b36',
        'hover': '#073642'
    },
    'matrix': {
        'primary': '#00ff41',
        'secondary': '#00c732',
        'accent': '#ffffff',
        'background': '#000000',
        'surface': '#0d0d0d',
        'text': '#00ff41',
        'text_secondary': '#008f2c',
        'border': '#005c1c',
        'header': '#0d0d0d',
        'grid': '#1a1a1a',
        'menu': '#0d0d0d',
        'hover': '#003b12'
    },
    'dracula': {
        'primary': '#ff79c6',
        'secondary': '#8be9fd',
        'accent': '#bd93f9',
        'background': '#282a36',
        'surface': '#44475a',
        'text': '#f8f8f2',
        'text_secondary': '#6272a4',
        'border': '#6272a4',
        'header': '#44475a',
        'grid': '#3a3c4a',
        'menu': '#282a36',
        'hover': '#44475a'
    },
    'gruvbox_dark': {
        'primary': '#458588',
        'secondary': '#98971a',
        'accent': '#689d6a',
        'background': '#282828',
        'surface': '#3c3836',
        'text': '#ebdbb2',
        'text_secondary': '#a89984',
        'border': '#504945',
        'header': '#3c3836',
        'grid': '#504945',
        'menu': '#282828',
        'hover': '#504945'
    },
    'blossom': {
        'primary': '#feacc6',
        'secondary': '#a6d8a8',
        'accent': '#d4a373',
        'background': '#fffbff',
        'surface': '#ffffff',
        'text': '#5b5353',
        'text_secondary': '#8f8a8a',
        'border': '#fde2e8',
        'header': '#fff0f3',
        'grid': '#fde2e8',
        'menu': '#fffbff',
        'hover': '#fff0f3'
    },

    # --- New high-contrast additions ---
    'vampire': {
        'primary': '#e60073',
        'secondary': '#ff1744',
        'accent': '#8c00ff',
        'background': '#0d001a',
        'surface': '#1a001f',
        'text': '#f5f5f5',
        'text_secondary': '#a64d79',
        'border': '#330033',
        'header': '#1f0026',
        'grid': '#260033',
        'menu': '#1a001f',
        'hover': '#330033'
    },
    'frostbyte': {
        'primary': '#00eaff',
        'secondary': '#00b7ff',
        'accent': '#ffffff',
        'background': '#0d1117',
        'surface': '#161b22',
        'text': '#c9d1d9',
        'text_secondary': '#8b949e',
        'border': '#30363d',
        'header': '#161b22',
        'grid': '#21262d',
        'menu': '#0d1117',
        'hover': '#1f6feb'
    },
    'catppuccin_mocha': {
        'primary': '#89b4fa',
        'secondary': '#f38ba8',
        'accent': '#a6e3a1',
        'background': '#1e1e2e',
        'surface': '#313244',
        'text': '#cdd6f4',
        'text_secondary': '#a6adc8',
        'border': '#45475a',
        'header': '#313244',
        'grid': '#45475a',
        'menu': '#1e1e2e',
        'hover': '#585b70'
    },
    'neon_noir': {
        'primary': '#ff007f',
        'secondary': '#00f0ff',
        'accent': '#faff00',
        'background': '#0a0a0a',
        'surface': '#1b1b1b',
        'text': '#e0e0e0',
        'text_secondary': '#808080',
        'border': '#222222',
        'header': '#1b1b1b',
        'grid': '#2b2b2b',
        'menu': '#0a0a0a',
        'hover': '#ff007f33'
    },
    'midnight_aurora': {
        'primary': '#7fffd4',
        'secondary': '#40e0d0',
        'accent': '#9370db',
        'background': '#081221',
        'surface': '#122033',
        'text': '#f0f8ff',
        'text_secondary': '#87cefa',
        'border': '#1e3a5f',
        'header': '#122033',
        'grid': '#1c2e4a',
        'menu': '#081221',
        'hover': '#2a4064'
    },
    'monokai_pro': {
        'primary': '#a6e22e',
        'secondary': '#f92672',
        'accent': '#fd971f',
        'background': '#272822',
        'surface': '#3e3d32',
        'text': '#f8f8f2',
        'text_secondary': '#75715e',
        'border': '#49483e',
        'header': '#3e3d32',
        'grid': '#49483e',
        'menu': '#272822',
        'hover': '#75715e'
    },
    'inferno': {
        'primary': '#ff6b00',
        'secondary': '#ff0000',
        'accent': '#ffff00',
        'background': '#1a0000',
        'surface': '#330000',
        'text': '#ffcc00',
        'text_secondary': '#ff6600',
        'border': '#660000',
        'header': '#330000',
        'grid': '#4d0000',
        'menu': '#1a0000',
        'hover': '#660000'
    },
    'arctic_dawn': {
        'primary': '#5dade2',
        'secondary': '#48c9b0',
        'accent': '#f4d03f',
        'background': '#eaf2f8',
        'surface': '#ffffff',
        'text': '#154360',
        'text_secondary': '#1f618d',
        'border': '#aed6f1',
        'header': '#d6eaf8',
        'grid': '#d4e6f1',
        'menu': '#eaf2f8',
        'hover': '#d1f2eb'
    },
    'steampunk': {
        'primary': '#b87333',
        'secondary': '#c0a080',
        'accent': '#ffd700',
        'background': '#2b1b0f',
        'surface': '#3c2615',
        'text': '#f5deb3',
        'text_secondary': '#d2b48c',
        'border': '#8b5a2b',
        'header': '#3c2615',
        'grid': '#4d331f',
        'menu': '#2b1b0f',
        'hover': '#5c3a24'
    },
    'synthwave': {
        'primary': '#ff00c8',
        'secondary': '#00fff9',
        'accent': '#fcee0c',
        'background': '#1b0033',
        'surface': '#2d004d',
        'text': '#f5f5f5',
        'text_secondary': '#a64dff',
        'border': '#7f00ff',
        'header': '#2d004d',
        'grid': '#3a0066',
        'menu': '#1b0033',
        'hover': '#5a00b3'
    },
    'terminal_green': {
        'primary': '#00ff00',
        'secondary': '#00cc00',
        'accent': '#99ff99',
        'background': '#000000',
        'surface': '#0a0a0a',
        'text': '#00ff00',
        'text_secondary': '#00cc00',
        'border': '#003300',
        'header': '#0a0a0a',
        'grid': '#1a1a1a',
        'menu': '#000000',
        'hover': '#003300'
    },
    'inkwell': {
        'primary': '#1a1a1a',
        'secondary': '#2e2e2e',
        'accent': '#4a4a4a',
        'background': '#f8f8f8',
        'surface': '#ffffff',
        'text': '#121212',
        'text_secondary': '#4a4a4a',
        'border': '#cccccc',
        'header': '#e6e6e6',
        'grid': '#dcdcdc',
        'menu': '#f8f8f8',
        'hover': '#dddddd'
    }
}
EXCLUDED_EXTENSIONS = {'.bak', '.log', '.ini', '.dwl', '.dwl2', '.tmp', '.lnk', '.db'}
FILE_TYPE_MAPPINGS = {
    "All Files": [], "Documents": ["doc", "docx", "pdf", "txt", "rtf", "odt", "pages"],
    "Spreadsheets": ["xls", "xlsx", "csv", "xlsm", "xlsb", "ods", "numbers"],
    "Presentations": ["ppt", "pptx", "pptm", "key"], "Images": ["jpg", "jpeg", "png", "gif", "bmp", "tiff", "webp", "svg", "heic"],
    "CAD Files": ["dwg", "dxf", "dwf", "rvt", "rfa", "dgn"], "Archives": ["zip", "rar", "7z", "tar", "gz", "bz2"]
}

class SearchEngine(QObject):
    """Core logic for finding files and duplicates. Can be used by GUI or CLI."""
    progress_update = pyqtSignal(int, str)

    def __init__(self):
        super().__init__()
        self.stopped = False
        self.match_count = 0
        self.group_count = 0
        self.scan_start_time = None
        self.last_progress_time = 0
        self.last_processed_count = 0

    def stop(self):
        self.stopped = True

    def _hash_file(self, path):
        """Calculate SHA256 hash of a file."""
        h = hashlib.sha256()
        try:
            with open(path, 'rb') as f:
                while True:
                    chunk = f.read(h.block_size * 1024)
                    if not chunk or self.stopped:
                        break
                    h.update(chunk)
            return h.hexdigest() if not self.stopped else None
        except (IOError, OSError):
            return None

    def find_duplicates(self, params):
        """Generator that finds and yields groups of duplicate files."""
        self.stopped = False
        self.match_count = 0
        self.group_count = 0
        
        self.progress_update.emit(0, "Stage 1/3: Grouping files by size...")
        size_map = defaultdict(list)
        total_files, processed_count = self._pre_scan(params)

        for root, dirs, files in os.walk(params['search_dir']):
            if self.stopped: return
            dirs[:] = [d for d in dirs if Path(root, d).as_posix() not in params['exclude_dirs']]
            for file in files:
                if self.stopped: return
                processed_count += 1
                try:
                    path = Path(root, file)
                    size = path.stat().st_size
                    if size > params.get('min_size_bytes', 1024):
                        size_map[size].append(path)
                except OSError:
                    continue
                if total_files > 0:
                    self.progress_update.emit(int(processed_count / total_files * 33), f"Scanned {processed_count}/{total_files} files...")

        self.progress_update.emit(33, "Stage 2/3: Hashing potential duplicates...")
        hash_map = defaultdict(list)
        potential_dupes = {size: paths for size, paths in size_map.items() if len(paths) > 1}
        
        files_to_hash = sum(len(paths) for paths in potential_dupes.values())
        hashed_count = 0

        with concurrent.futures.ThreadPoolExecutor(max_workers=os.cpu_count()) as executor:
            future_to_path = {executor.submit(self._hash_file, path): path for size, paths in potential_dupes.items() for path in paths}
            for future in concurrent.futures.as_completed(future_to_path):
                if self.stopped:
                    executor.shutdown(wait=False)
                    return
                path = future_to_path[future]
                try:
                    file_hash = future.result()
                    if file_hash:
                        hash_map[file_hash].append(path)
                except Exception:
                    pass
                hashed_count += 1
                if files_to_hash > 0:
                    self.progress_update.emit(33 + int(hashed_count / files_to_hash * 33), f"Hashed {hashed_count}/{files_to_hash} files...")

        self.progress_update.emit(66, "Stage 3/3: Finalizing groups...")
        duplicate_groups = {h: p for h, p in hash_map.items() if len(p) > 1}
        self.group_count = len(duplicate_groups)
        
        for file_hash, paths in duplicate_groups.items():
            if self.stopped: return
            file_group = []
            for path in paths:
                self.match_count += 1
                info = self._get_file_info(path, {'hash': file_hash})
                if info:
                    file_group.append(info)
            if file_group:
                yield file_group

    def find_files(self, params):
        """Generator that finds and yields files matching criteria with multi-threaded scanning."""
        self.stopped = False
        self.match_count = 0
        self._compile_patterns(params)
        total_files, processed_count = self._pre_scan(params)
        self.scan_start_time = time.time()
        self.last_progress_time = self.scan_start_time
        self.last_processed_count = 0

        # Check if this is a large directory and warn user
        if total_files > LARGE_DIR_THRESHOLD:
            self.progress_update.emit(0, f"⚠️ Large directory detected: {total_files:,} files. Scan limited to {MAX_RESULTS:,} results.")

        # Use multi-threaded scanning for better performance
        if params.get('use_multithreading', True):
            for file_info in self._find_files_multithreaded(params, total_files):
                if self.stopped: break
                if self.match_count >= MAX_RESULTS:
                    self.progress_update.emit(100, f"⚠️ Reached maximum results limit ({MAX_RESULTS:,}). Use more specific search criteria.")
                    break
                yield file_info
        else:
            # Fallback to single-threaded for compatibility
            for file_info in self._find_files_singlethreaded(params, total_files, 0):
                if self.stopped: break
                if self.match_count >= MAX_RESULTS:
                    self.progress_update.emit(100, f"⚠️ Reached maximum results limit ({MAX_RESULTS:,}). Use more specific search criteria.")
                    break
                yield file_info
    
    def _find_files_singlethreaded(self, params, total_files, processed_count):
        """Original single-threaded implementation"""
        for root, dirs, files in os.walk(params['search_dir']):
            if self.stopped: break
            dirs[:] = [d for d in dirs if Path(root, d).as_posix() not in params['exclude_dirs']]
            
            # Update progress with current folder name
            current_folder = Path(root).name or str(root)
            if len(current_folder) > 50:
                current_folder = '...' + current_folder[-47:]
            
            for file in files:
                if self.stopped: break
                processed_count += 1
                
                # Stop scanning if we've hit the maximum scan limit
                if processed_count >= MAX_SCAN_FILES:
                    self.progress_update.emit(100, f"⚠️ Scan limit reached ({MAX_SCAN_FILES:,} files). Use more specific search criteria.")
                    break
                
                # Enhanced progress feedback with speed calculation
                if total_files > 0 and processed_count % 100 == 0:
                    speed_info = self._calculate_speed(processed_count)
                    progress_msg = f"📂 {current_folder} | {processed_count:,}/{total_files:,} files{speed_info}"
                    self.progress_update.emit(int(processed_count / total_files * 100), progress_msg)
                elif total_files <= 0 and processed_count % 500 == 0:
                    speed_info = self._calculate_speed(processed_count)
                    progress_msg = f"📂 {current_folder} | {processed_count:,} files{speed_info}"
                    self.progress_update.emit(-1, progress_msg)

                file_path = Path(root, file)
                if self._is_file_match(file_path, params):
                    self.match_count += 1
                    info = self._get_file_info(file_path)
                    if info:
                        yield info
    
    def _find_files_multithreaded(self, params, total_files):
        """Multi-threaded implementation using ThreadPoolExecutor for parallel directory scanning"""
        import queue
        import threading
        
        # Thread-safe queue for results
        result_queue = queue.Queue(maxsize=1000)
        processed_count = [0]  # Use list for thread-safe counter
        counter_lock = threading.Lock()
        
        def scan_directory(dir_path):
            """Scan a single directory and its subdirectories"""
            local_results = []
            try:
                for root, dirs, files in os.walk(dir_path):
                    if self.stopped: break
                    
                    # Filter excluded dirs
                    dirs[:] = [d for d in dirs if Path(root, d).as_posix() not in params['exclude_dirs']]
                    
                    current_folder = Path(root).name or str(root)
                    if len(current_folder) > 50:
                        current_folder = '...' + current_folder[-47:]
                    
                    for file in files:
                        if self.stopped: break
                        
                        # Thread-safe counter increment
                        with counter_lock:
                            processed_count[0] += 1
                            count = processed_count[0]
                        
                        # Stop scanning if we've hit the maximum scan limit
                        if count >= MAX_SCAN_FILES:
                            self.progress_update.emit(100, f"⚠️ Scan limit reached ({MAX_SCAN_FILES:,} files). Use more specific search criteria.")
                            break
                        
                        # Progress update (less frequent to reduce lock contention)
                        if count % 200 == 0:
                            speed_info = self._calculate_speed(count)
                            if total_files > 0:
                                progress_msg = f"📂 {current_folder} | {count:,}/{total_files:,} files{speed_info}"
                                self.progress_update.emit(int(count / total_files * 100), progress_msg)
                            else:
                                progress_msg = f"📂 {current_folder} | {count:,} files{speed_info}"
                                self.progress_update.emit(-1, progress_msg)
                        
                        file_path = Path(root, file)
                        if self._is_file_match(file_path, params):
                            info = self._get_file_info(file_path)
                            if info:
                                local_results.append(info)
                                # Add to queue immediately for streaming
                                if len(local_results) >= 50:  # Batch for efficiency
                                    result_queue.put(local_results.copy())
                                    local_results.clear()
            except Exception as e:
                pass  # Skip inaccessible directories
            
            # Add remaining results
            if local_results:
                result_queue.put(local_results)
        
        # Get top-level directories to parallelize
        search_dir = Path(params['search_dir'])
        subdirs = []
        try:
            for item in search_dir.iterdir():
                if item.is_dir() and item.as_posix() not in params['exclude_dirs']:
                    subdirs.append(str(item))
        except Exception:
            subdirs = [str(search_dir)]  # Fallback to single dir
        
        # If too few subdirs, scan the root directly
        if len(subdirs) < 2:
            subdirs = [str(search_dir)]
        
        # Start parallel scanning with ThreadPoolExecutor
        max_workers = min(os.cpu_count() or 4, len(subdirs), 8)  # Limit to 8 threads max
        
        def producer():
            with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
                futures = [executor.submit(scan_directory, subdir) for subdir in subdirs]
                concurrent.futures.wait(futures)
            result_queue.put(None)  # Sentinel to signal completion
        
        # Start producer thread
        producer_thread = threading.Thread(target=producer, daemon=True)
        producer_thread.start()
        
        # Consumer: yield results as they come in
        while True:
            batch = result_queue.get()
            if batch is None:  # Sentinel value
                break
            for file_info in batch:
                self.match_count += 1
                yield file_info
                if self.stopped:
                    return

    def _pre_scan(self, params):
        if params.get('count_files', True):
            total_files = self._count_files(params)
            return total_files, 0
        return -1, 0

    def _count_files(self, params):
        self.progress_update.emit(0, "Counting total files...")
        count = 0
        for root, dirs, files in os.walk(params['search_dir']):
            if self.stopped: return -1
            dirs[:] = [d for d in dirs if Path(root, d).as_posix() not in params['exclude_dirs']]
            count += len(files)
        self.progress_update.emit(0, f"Found {count} total files to process.")
        return count

    def _get_file_info(self, file_path, extra_data=None):
        try:
            stat = file_path.stat()
            info = {
                'filename': file_path.name,
                'path': str(file_path.parent),
                'full_path': str(file_path),
                'extension': file_path.suffix.lower(),
                'size_kb': round(stat.st_size / 1024, 2),
                'size_bytes': stat.st_size,
                'modified_date': datetime.datetime.fromtimestamp(stat.st_mtime),
                'created_date': datetime.datetime.fromtimestamp(stat.st_ctime),
            }
            if extra_data:
                info.update(extra_data)
            return info
        except OSError:
            return None

    def _calculate_speed(self, processed_count):
        """Calculate and format scanning speed"""
        current_time = time.time()
        time_diff = current_time - self.last_progress_time
        
        if time_diff > 0.5:  # Update speed every 0.5 seconds
            files_diff = processed_count - self.last_processed_count
            speed = files_diff / time_diff
            self.last_progress_time = current_time
            self.last_processed_count = processed_count
            
            # Calculate elapsed time
            elapsed = current_time - self.scan_start_time
            elapsed_str = f"{int(elapsed)}s"
            if elapsed >= 60:
                elapsed_str = f"{int(elapsed//60)}m {int(elapsed%60)}s"
            
            return f" | ⚡ {speed:.0f}/s | ⏱️ {elapsed_str}"
        return ""
    
    def _compile_patterns(self, params):
        self.keywords = [k.lower().strip() for k in params['keywords'].split(',') if k.strip()]
        self.exclusion_keywords = [k.lower().strip() for k in params['exclusion_keywords'].split(',') if k.strip()]
        if params['use_regex']:
            self.patterns = [re.compile(kw, re.IGNORECASE) for kw in self.keywords]
            self.exclusion_patterns = [re.compile(kw, re.IGNORECASE) for kw in self.exclusion_keywords]
        elif params['whole_words']:
            self.patterns = [re.compile(r'\b' + re.escape(kw) + r'\b', re.IGNORECASE) for kw in self.keywords]
            self.exclusion_patterns = [re.compile(r'\b' + re.escape(kw) + r'\b', re.IGNORECASE) for kw in self.exclusion_keywords]
        else:
            self.patterns = self.keywords
            self.exclusion_patterns = self.exclusion_keywords

    def _is_file_match(self, file_path, params):
        file_ext = file_path.suffix.lower()
        if not file_ext or file_ext in EXCLUDED_EXTENSIONS: return False
        if params['allowed_extensions'] and file_ext.lstrip('.') not in params['allowed_extensions']: return False
        try:
            file_size_kb = file_path.stat().st_size / 1024
            if params['min_size_kb'] is not None and file_size_kb < params['min_size_kb']: return False
            if params['max_size_kb'] is not None and file_size_kb > params['max_size_kb']: return False
        except OSError: return False
        if params['date_filter']:
            try:
                ts = file_path.stat().st_mtime if params['date_filter_type'] == 'modified' else file_path.stat().st_ctime
                file_date = datetime.datetime.fromtimestamp(ts)
                if params['min_date'] and file_date < params['min_date']: return False
                if params['max_date'] and file_date > params['max_date']: return False
            except OSError: return False
        if params.get('content_search'):
            mime_type, _ = mimetypes.guess_type(str(file_path))
            if mime_type and mime_type.startswith('text/'):
                try:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        content = f.read(1024 * 1024)  # Read up to 1MB
                        if params['content_search'].lower() not in content.lower():
                            return False
                except Exception:
                    return False
        return self._matches_keywords(file_path.stem)
    
    def _matches_keywords(self, file_name):
        if not self.keywords and not self.exclusion_keywords: return True
        file_name_lower = file_name.lower()
        if self.exclusion_keywords:
            if self.exclusion_patterns and isinstance(self.exclusion_patterns[0], re.Pattern):
                if any(p.search(file_name_lower) for p in self.exclusion_patterns): return False
            else:
                if any(kw in file_name_lower for kw in self.exclusion_patterns): return False
        if not self.keywords: return True
        if self.patterns and isinstance(self.patterns[0], re.Pattern):
            return any(p.search(file_name_lower) for p in self.patterns)
        else:
            return any(kw in file_name_lower for kw in self.patterns)

class FileSearchWorker(QThread):
    """Worker thread that runs the core search logic and emits signals."""
    progress_update = pyqtSignal(int, str)
    search_complete = pyqtSignal(bool, str)
    live_result = pyqtSignal(dict)
    duplicate_group_found = pyqtSignal(list)

    def __init__(self, params):
        super().__init__()
        self.params = params
        self.engine = SearchEngine()
        self.engine.progress_update.connect(self.progress_update)
        self.stopped = False

    def run(self):
        try:
            if self.params['search_mode'] == 'duplicates':
                for group in self.engine.find_duplicates(self.params):
                    if self.stopped: break
                    self.duplicate_group_found.emit(group)
            else:
                for file_info in self.engine.find_files(self.params):
                    if self.stopped: break
                    self.live_result.emit(file_info)
            
            if self.stopped:
                self.search_complete.emit(False, "Search stopped by user.")
            else:
                msg = f"Search complete. Found {self.engine.match_count} files."
                if self.params['search_mode'] == 'duplicates':
                    msg = f"Search complete. Found {self.engine.match_count} duplicate files in {self.engine.group_count} groups."
                self.search_complete.emit(True, msg)
        except Exception as e:
            self.search_complete.emit(False, f"An error occurred: {e}")

    def stop(self):
        self.stopped = True
        self.engine.stop()

class ProfileManagerDialog(QDialog):
    """A dialog for managing saved search profiles."""
    def __init__(self, parent=None, zoom_level=100):
        super().__init__(parent)
        self.setWindowTitle("Profile Manager")
        self.setMinimumSize(400, 300)
        self.settings = QSettings(SETTINGS_ORG, SETTINGS_APP)
        self.zoom_level = zoom_level
        
        layout = QVBoxLayout(self)
        self.profile_list = QListWidget()
        self.profile_list.itemDoubleClicked.connect(self.accept)
        layout.addWidget(self.profile_list)
        button_layout = QHBoxLayout()
        self.rename_btn = QPushButton("Rename")
        self.rename_btn.clicked.connect(self.rename_profile)
        self.delete_btn = QPushButton("Delete")
        self.delete_btn.clicked.connect(self.delete_profile)
        button_layout.addWidget(self.rename_btn)
        button_layout.addWidget(self.delete_btn)
        button_layout.addStretch()
        self.dialog_buttons = QDialogButtonBox(QDialogButtonBox.StandardButton.Open | QDialogButtonBox.StandardButton.Cancel)
        self.dialog_buttons.accepted.connect(self.accept)
        self.dialog_buttons.rejected.connect(self.reject)
        button_layout.addWidget(self.dialog_buttons)
        layout.addLayout(button_layout)
        self.load_profiles()
        
        # Apply zoom from parent window AFTER UI is created
        if zoom_level != 100:
            self._apply_zoom()
    
    def _apply_zoom(self):
        """Apply zoom level to all dialog elements"""
        # Calculate new font size based on zoom
        base_font_size = 9
        new_font_size = int(base_font_size * (self.zoom_level / 100))
        
        # Create font for the dialog
        app_font = self.font()
        app_font.setPointSize(new_font_size)
        
        # Apply to dialog and all widgets
        self.setFont(app_font)
        
        # Update all child widgets
        for widget in self.findChildren(QWidget):
            widget.setFont(app_font)

    def load_profiles(self):
        self.profile_list.clear()
        profiles = self.settings.value("profiles", {})
        for name in sorted(profiles.keys()):
            self.profile_list.addItem(name)

    def rename_profile(self):
        current_item = self.profile_list.currentItem()
        if not current_item: return
        old_name = current_item.text()
        new_name, ok = QInputDialog.getText(self, "Rename Profile", "New profile name:", text=old_name)
        if ok and new_name and new_name != old_name:
            profiles = self.settings.value("profiles", {})
            if new_name in profiles:
                QMessageBox.warning(self, "Name Exists", "A profile with this name already exists.")
                return
            profiles[new_name] = profiles.pop(old_name)
            self.settings.setValue("profiles", profiles)
            self.load_profiles()

    def delete_profile(self):
        current_item = self.profile_list.currentItem()
        if not current_item: return
        reply = QMessageBox.question(self, "Confirm Delete", f"Are you sure you want to delete the profile '{current_item.text()}'?")
        if reply == QMessageBox.StandardButton.Yes:
            profiles = self.settings.value("profiles", {})
            profiles.pop(current_item.text(), None)
            self.settings.setValue("profiles", profiles)
            self.load_profiles()

    def get_selected_profile_name(self):
        return self.profile_list.currentItem().text() if self.profile_list.currentItem() else None

class SmartSortDialog(QDialog):
    """Preview and execute Smart Sort operations based on file extensions."""
    def __init__(self, parent, files, default_root, zoom_level=100):
        super().__init__(parent)
        self.setWindowTitle("Smart Sort")
        self.setMinimumSize(800, 500)
        self.parent_app = parent
        self.files = files  # list of file_info dicts
        self.default_root = default_root
        self.zoom_level = zoom_level

        main_layout = QVBoxLayout(self)

        # Controls
        ctrl_box = QGroupBox("Options")
        ctrl_layout = QGridLayout()
        ctrl_layout.addWidget(QLabel("Destination Root:"), 0, 0)
        self.root_edit = QLineEdit(self.default_root)
        browse_btn = QPushButton("Browse…")
        browse_btn.clicked.connect(self._browse_root)
        ctrl_layout.addWidget(self.root_edit, 0, 1)
        ctrl_layout.addWidget(browse_btn, 0, 2)
        self.copy_checkbox = QCheckBox("Copy instead of Move")
        ctrl_layout.addWidget(self.copy_checkbox, 1, 1)
        ctrl_box.setLayout(ctrl_layout)
        main_layout.addWidget(ctrl_box)

        # Table
        self.table = QTableWidget(0, 5)
        self.table.setHorizontalHeaderLabels(["Include", "File", "Current Path", "Ext", "Destination"])
        self.table.horizontalHeader().setSectionResizeMode(0, QHeaderView.ResizeMode.ResizeToContents)
        self.table.horizontalHeader().setSectionResizeMode(1, QHeaderView.ResizeMode.ResizeToContents)
        self.table.horizontalHeader().setSectionResizeMode(2, QHeaderView.ResizeMode.Stretch)
        self.table.horizontalHeader().setSectionResizeMode(3, QHeaderView.ResizeMode.ResizeToContents)
        self.table.horizontalHeader().setSectionResizeMode(4, QHeaderView.ResizeMode.Stretch)
        self._populate_table()
        main_layout.addWidget(self.table)

        # Buttons
        btn_box = QDialogButtonBox()
        self.exec_btn = btn_box.addButton("Execute Sort", QDialogButtonBox.ButtonRole.AcceptRole)
        self.cancel_btn = btn_box.addButton(QDialogButtonBox.StandardButton.Cancel)
        self.exec_btn.clicked.connect(self._execute)
        self.cancel_btn.clicked.connect(self.reject)
        main_layout.addWidget(btn_box)

        # Recompute destinations when root changes
        self.root_edit.textChanged.connect(self._recompute_destinations)
        
        # Apply zoom from parent window AFTER UI is created
        if zoom_level != 100:
            self._apply_zoom()
    
    def _apply_zoom(self):
        """Apply zoom level to all dialog elements"""
        # Calculate new font size based on zoom
        base_font_size = 9
        new_font_size = int(base_font_size * (self.zoom_level / 100))
        
        # Create font for the dialog
        app_font = self.font()
        app_font.setPointSize(new_font_size)
        
        # Apply to dialog and all widgets
        self.setFont(app_font)
        
        # Update all child widgets
        for widget in self.findChildren(QWidget):
            widget.setFont(app_font)
        
        # Special handling for tables to adjust row heights
        for table in self.findChildren(QTableWidget):
            table.verticalHeader().setDefaultSectionSize(int(25 * (self.zoom_level / 100)))

    def _browse_root(self):
        d = QFileDialog.getExistingDirectory(self, "Select Destination Root", self.root_edit.text() or str(Path.home()))
        if d:
            self.root_edit.setText(d)

    def _ext_folder(self, ext: str) -> str:
        if not ext:
            return "Unsorted"
        return ext.lstrip('.').upper()

    def _suggest_destination(self, file_info) -> Path:
        root = Path(self.root_edit.text())
        ext = file_info.get('extension', '')
        folder = self._ext_folder(ext)
        return root / folder / file_info['filename']

    def _populate_table(self):
        self.table.setRowCount(0)
        for fi in self.files:
            row = self.table.rowCount()
            self.table.insertRow(row)
            # Include checkbox
            include_item = QTableWidgetItem()
            include_item.setFlags(include_item.flags() | Qt.ItemFlag.ItemIsUserCheckable)
            include_item.setCheckState(Qt.CheckState.Checked)
            include_item.setData(Qt.ItemDataRole.UserRole, fi)
            self.table.setItem(row, 0, include_item)
            # File name
            self.table.setItem(row, 1, QTableWidgetItem(fi.get('filename', '')))
            # Current path
            self.table.setItem(row, 2, QTableWidgetItem(Path(fi['full_path']).parent.as_posix()))
            # Ext
            self.table.setItem(row, 3, QTableWidgetItem(fi.get('extension', '')))
            # Destination
            dest_path = self._suggest_destination(fi)
            self.table.setItem(row, 4, QTableWidgetItem(dest_path.as_posix()))

    def _recompute_destinations(self):
        for row in range(self.table.rowCount()):
            fi = self.table.item(row, 0).data(Qt.ItemDataRole.UserRole)
            dest_path = self._suggest_destination(fi)
            self.table.item(row, 4).setText(dest_path.as_posix())

    def keyPressEvent(self, event):
        # Space toggles include for selected rows
        if event.key() == Qt.Key.Key_Space:
            for idx in self.table.selectionModel().selectedRows():
                item = self.table.item(idx.row(), 0)
                if item.checkState() == Qt.CheckState.Checked:
                    item.setCheckState(Qt.CheckState.Unchecked)
                else:
                    item.setCheckState(Qt.CheckState.Checked)
            event.accept()
            return
        super().keyPressEvent(event)

    def _execute(self):
        root = Path(self.root_edit.text())
        if not root.exists():
            try:
                root.mkdir(parents=True, exist_ok=True)
            except Exception as e:
                QMessageBox.critical(self, "Error", f"Could not create destination root:\n{e}")
                return

        do_copy = self.copy_checkbox.isChecked()
        success = 0
        errors = []
        moved_paths = []
        for row in range(self.table.rowCount()):
            inc_item = self.table.item(row, 0)
            if inc_item.checkState() != Qt.CheckState.Checked:
                continue
            fi = inc_item.data(Qt.ItemDataRole.UserRole)
            src = Path(fi['full_path'])
            dest = Path(self.table.item(row, 4).text())
            try:
                dest.parent.mkdir(parents=True, exist_ok=True)
                final_dest = self.parent_app._generate_unique_dest_path(dest)
                if do_copy:
                    shutil.copy2(src, final_dest)
                else:
                    src.rename(final_dest)
                    moved_paths.append(src)
                success += 1
            except Exception as e:
                errors.append((str(src), str(e)))
            if success % 10 == 0:
                QApplication.processEvents()

        summary = f"Smart Sort complete. {'Copied' if do_copy else 'Moved'} {success} item(s)."
        if errors:
            summary += f" Failed: {len(errors)}."
        QMessageBox.information(self, "Smart Sort", summary)
        # Update parent table if we moved files
        if moved_paths:
            self.parent_app.remove_files_from_results(moved_paths)
        self.accept()

# Enhanced Preview Handler System
class PreviewHandler:
    """Base class for file preview handlers."""
    
    def __init__(self, name, extensions=None):
        self.name = name
        self.extensions = extensions or []
    
    def can_handle(self, file_path):
        """Check if this handler can preview the file."""
        file_path = Path(file_path)
        return file_path.suffix.lower() in self.extensions
    
    def generate_preview(self, file_path, max_size=1024*1024):
        """Generate preview content. Returns (content_type, data, metadata)."""
        raise NotImplementedError

class TextPreviewHandler(PreviewHandler):
    """Handler for text files with syntax highlighting."""
    
    def __init__(self):
        super().__init__("Text", ['.txt', '.log', '.cfg', '.ini', '.conf', '.md', '.rst'])
    
    def generate_preview(self, file_path, max_size=1024*1024):
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read(min(max_size, 5000))  # Read up to 5KB
                
                if PYGMENTS_AVAILABLE:
                    try:
                        lexer = get_lexer_for_filename(str(file_path))
                        formatter = HtmlFormatter(style='default', linenos=True)
                        html_content = highlight(content, lexer, formatter)
                        return ("html", html_content, {"syntax_highlighted": True})
                    except:
                        return ("text", content, {"syntax_highlighted": False})
                else:
                    return ("text", content, {"syntax_highlighted": False})
        except Exception as e:
            return ("error", str(e), {})

class CodePreviewHandler(PreviewHandler):
    """Handler for source code files with syntax highlighting."""
    
    def __init__(self):
        super().__init__("Code", [
            '.py', '.js', '.html', '.css', '.java', '.cpp', '.c', '.h', '.hpp',
            '.cs', '.php', '.rb', '.go', '.rs', '.swift', '.kt', '.scala',
            '.sh', '.bat', '.ps1', '.sql', '.xml', '.json', '.yaml', '.yml',
            '.toml', '.dockerfile', '.gitignore', '.env', '.htaccess'
        ])
    
    def generate_preview(self, file_path, max_size=1024*1024):
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read(min(max_size, 5000))
                
                if PYGMENTS_AVAILABLE:
                    try:
                        lexer = get_lexer_for_filename(str(file_path))
                        formatter = HtmlFormatter(style='monokai', linenos=True)
                        html_content = highlight(content, lexer, formatter)
                        return ("html", html_content, {"syntax_highlighted": True, "theme": "monokai"})
                    except:
                        return ("text", content, {"syntax_highlighted": False})
                else:
                    return ("text", content, {"syntax_highlighted": False})
        except Exception as e:
            return ("error", str(e), {})

class PDFPreviewHandler(PreviewHandler):
    """Handler for PDF files."""
    
    def __init__(self):
        super().__init__("PDF", ['.pdf'])
    
    def generate_preview(self, file_path, max_size=1024*1024):
        if not PYMUPDF_AVAILABLE:
            return ("error", "PyMuPDF not installed. Install with: pip install PyMuPDF", {})
        
        try:
            doc = fitz.open(str(file_path))
            if doc.page_count > 0:
                # Extract text from first few pages
                all_text = ""
                formatted_pages = []
                page_count = doc.page_count  # Store before closing
                
                for page_num in range(min(3, page_count)):  # First 3 pages max
                    page = doc[page_num]
                    text = page.get_text()
                    if text.strip():
                        all_text += f"\n--- Page {page_num + 1} ---\n{text}"
                        
                        # Create formatted version with page numbers
                        formatted_text = f"═══════════════════════════════════════\n"
                        formatted_text += f"          PAGE {page_num + 1} of {page_count}\n"
                        formatted_text += f"═══════════════════════════════════════\n\n"
                        formatted_text += text
                        formatted_pages.append(formatted_text)
                
                doc.close()
                
                if all_text.strip():
                    # Return both simple text and formatted document
                    simple_text = all_text[:3000]  # Limit for text tab
                    formatted_text = '\n\n'.join(formatted_pages)
                    
                    return ("pdf_dual", simple_text, {
                        "formatted": formatted_text,
                        "pages": page_count,
                        "preview_pages": min(3, page_count),
                        "file_path": str(file_path)
                    })
                else:
                    # Even if no text can be extracted, still return for visual viewing
                    return ("pdf_dual", "PDF contains no extractable text (image-based PDF)\n\nUse the PDF Viewer tab to see the visual content.", {
                        "formatted": "Image-based PDF - No extractable text available\n\nUse the PDF Viewer tab to see the visual content.",
                        "pages": page_count,
                        "preview_pages": min(3, page_count),
                        "file_path": str(file_path)
                    })
            else:
                doc.close()
                return ("error", "PDF appears to be empty", {})
        except Exception as e:
            return ("error", f"PDF Error: {e}", {})

class ExcelPreviewHandler(PreviewHandler):
    """Handler for Excel files."""
    
    def __init__(self):
        super().__init__("Excel", ['.xlsx', '.xlsm', '.xls'])
    
    def generate_preview(self, file_path, max_size=1024*1024):
        file_path = Path(file_path)
        extension = file_path.suffix.lower()
        
        try:
            if extension == '.xls':
                # Use xlrd for old .xls format
                if not XLRD_AVAILABLE:
                    return ("error", "xlrd not installed for .xls files. Install with: pip install xlrd", {})
                
                return self._read_xls_file(file_path)
            else:
                # Use openpyxl for .xlsx and .xlsm formats
                if not OPENPYXL_AVAILABLE:
                    return ("error", "openpyxl not installed for .xlsx/.xlsm files. Install with: pip install openpyxl", {})
                
                return self._read_xlsx_file(file_path)
                
        except Exception as e:
            return ("error", str(e), {})
    
    def _read_xls_file(self, file_path):
        """Read .xls file using xlrd."""
        try:
            wb = xlrd.open_workbook(str(file_path))
            sheet = wb.sheet_by_index(0)  # First sheet
            
            content = []
            # Read first 10 rows and 10 columns
            for row_idx in range(min(10, sheet.nrows)):
                row_data = []
                for col_idx in range(min(10, sheet.ncols)):
                    cell = sheet.cell(row_idx, col_idx)
                    
                    # Handle different cell types
                    if cell.ctype == xlrd.XL_CELL_NUMBER:
                        # Handle numbers (including dates)
                        if 0 < cell.value < 1 and sheet.book.datemode:
                            # Likely a date
                            try:
                                date_tuple = xlrd.xldate_as_tuple(cell.value, sheet.book.datemode)
                                if date_tuple[0] > 0:  # Valid date
                                    cell_value = f"{date_tuple[0]}-{date_tuple[1]:02d}-{date_tuple[2]:02d}"
                                else:
                                    cell_value = str(cell.value)
                            except:
                                cell_value = str(cell.value)
                        else:
                            cell_value = str(cell.value)
                    elif cell.ctype == xlrd.XL_CELL_TEXT:
                        cell_value = str(cell.value)
                    elif cell.ctype == xlrd.XL_CELL_BLANK:
                        cell_value = ''
                    else:
                        cell_value = str(cell.value)
                    
                    row_data.append(cell_value)
                
                if any(cell.strip() for cell in row_data if cell):
                    content.append('\t'.join(row_data))
            
            if content:
                return ("text", '\n'.join(content), {"sheets": wb.nsheets, "active_sheet": sheet.name, "format": "XLS"})
            else:
                return ("text", "Excel file appears to be empty", {"sheets": wb.nsheets, "format": "XLS"})
                
        except Exception as e:
            return ("error", f"Error reading .xls file: {e}", {})
    
    def _read_xlsx_file(self, file_path):
        """Read .xlsx/.xlsm file using openpyxl."""
        try:
            wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
            if wb.sheetnames:
                sheet = wb.active
                content = []
                # Read first 10 rows and 10 columns
                for row in sheet.iter_rows(max_row=10, max_col=10, values_only=True):
                    if any(cell is not None for cell in row):
                        content.append('\t'.join(str(cell) if cell is not None else '' for cell in row))
                
                wb.close()
                if content:
                    return ("text", '\n'.join(content), {"sheets": len(wb.sheetnames), "active_sheet": sheet.title, "format": "XLSX"})
                else:
                    return ("text", "Excel file appears to be empty", {"sheets": len(wb.sheetnames), "format": "XLSX"})
        except Exception as e:
            return ("error", f"Error reading .xlsx/.xlsm file: {e}", {})

class CSVPreviewHandler(PreviewHandler):
    """Handler for CSV files."""
    
    def __init__(self):
        super().__init__("CSV", ['.csv'])
    
    def generate_preview(self, file_path, max_size=1024*1024):
        try:
            import csv
            
            content = []
            row_count = 0
            max_rows = 20  # Show first 20 rows
            max_cols = 15  # Show first 15 columns
            
            with open(str(file_path), 'r', encoding='utf-8', errors='replace') as f:
                # Try to detect the delimiter
                sample = f.read(1024)
                f.seek(0)
                
                sniffer = csv.Sniffer()
                try:
                    delimiter = sniffer.sniff(sample).delimiter
                except:
                    delimiter = ','  # Default to comma
                
                reader = csv.reader(f, delimiter=delimiter)
                
                for row in reader:
                    if row_count >= max_rows:
                        break
                    
                    # Limit columns and format each cell
                    formatted_row = []
                    for i, cell in enumerate(row[:max_cols]):
                        # Clean up the cell content
                        cell_str = str(cell).strip()
                        if len(cell_str) > 50:
                            cell_str = cell_str[:47] + "..."
                        formatted_row.append(cell_str)
                    
                    # Add indicator if there are more columns
                    if len(row) > max_cols:
                        formatted_row.append("...")
                    
                    if any(cell.strip() for cell in formatted_row):  # Skip empty rows
                        content.append('\t'.join(formatted_row))
                        row_count += 1
            
            if content:
                # Add metadata about the CSV
                metadata = {
                    "delimiter": delimiter,
                    "rows_shown": row_count,
                    "max_columns": max_cols,
                    "format": "CSV"
                }
                
                # Add header info if we have data
                if content:
                    header_info = f"CSV Preview (delimiter: '{delimiter}') - First {row_count} rows, max {max_cols} columns\n"
                    header_info += "─" * 60 + "\n"
                    return ("text", header_info + '\n'.join(content), metadata)
                else:
                    return ("text", '\n'.join(content), metadata)
            else:
                return ("text", "CSV file appears to be empty or contains only empty rows", {"format": "CSV"})
                
        except Exception as e:
            return ("error", f"Error reading CSV file: {e}", {})

class WordPreviewHandler(PreviewHandler):
    """Handler for Word documents."""
    
    def __init__(self):
        super().__init__("Word", ['.docx'])
    
    def generate_preview(self, file_path, max_size=1024*1024):
        if not DOCX_AVAILABLE:
            return ("error", "python-docx not installed. Install with: pip install python-docx", {})
        
        try:
            doc = Document(str(file_path))
            paragraphs = []
            for para in doc.paragraphs[:20]:  # First 20 paragraphs
                if para.text.strip():
                    paragraphs.append(para.text)
            
            if paragraphs:
                return ("text", '\n\n'.join(paragraphs), {"paragraphs": len(doc.paragraphs)})
            else:
                return ("text", "Document appears to be empty", {"paragraphs": len(doc.paragraphs)})
        except Exception as e:
            return ("error", str(e), {})

class PowerPointPreviewHandler(PreviewHandler):
    """Handler for PowerPoint presentations."""
    
    def __init__(self):
        super().__init__("PowerPoint", ['.pptx'])
    
    def generate_preview(self, file_path, max_size=1024*1024):
        if not PPTX_AVAILABLE:
            return ("error", "python-pptx not installed. Install with: pip install python-pptx", {})
        
        try:
            prs = Presentation(str(file_path))
            slide_texts = []
            for i, slide in enumerate(prs.slides[:10]):  # First 10 slides
                text_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text)
                if text_content:
                    slide_texts.append(f"Slide {i+1}:\n" + '\n'.join(text_content))
            
            if slide_texts:
                return ("text", '\n\n---\n\n'.join(slide_texts), {"slides": len(prs.slides)})
            else:
                return ("text", f"Presentation has {len(prs.slides)} slides but no extractable text", {"slides": len(prs.slides)})
        except Exception as e:
            return ("error", str(e), {})

class AudioPreviewHandler(PreviewHandler):
    """Handler for audio files - shows metadata."""
    
    def __init__(self):
        super().__init__("Audio", ['.mp3', '.flac', '.ogg', '.wav', '.m4a', '.aac'])
    
    def generate_preview(self, file_path, max_size=1024*1024):
        if not MUTAGEN_AVAILABLE:
            return ("error", "mutagen not installed. Install with: pip install mutagen", {})
        
        try:
            audio = mutagen.File(str(file_path))
            if audio is not None:
                metadata = {}
                for key, value in audio.items():
                    if isinstance(value, list) and value:
                        metadata[key] = str(value[0])
                    else:
                        metadata[key] = str(value)
                
                # Format metadata as readable text
                info_lines = [f"Audio File Metadata:"]
                info_lines.extend([f"{k}: {v}" for k, v in metadata.items()])
                return ("text", '\n'.join(info_lines), {"metadata_keys": len(metadata)})
            else:
                return ("text", "Unable to read audio metadata", {})
        except Exception as e:
            return ("error", str(e), {})

class VideoPreviewHandler(PreviewHandler):
    """Handler for video files - shows basic info."""
    
    def __init__(self):
        super().__init__("Video", ['.mp4', '.avi', '.mkv', '.mov', '.wmv', '.flv', '.webm'])
    
    def generate_preview(self, file_path, max_size=1024*1024):
        try:
            file_path = Path(file_path)
            size_mb = file_path.stat().st_size / (1024 * 1024)
            
            info_lines = [
                f"Video File Information:",
                f"File: {file_path.name}",
                f"Size: {size_mb:.2f} MB",
                f"Extension: {file_path.suffix}",
                "",
                "Note: Detailed video metadata requires ffmpeg or similar tools."
            ]
            
            return ("text", '\n'.join(info_lines), {"size_mb": size_mb})
        except Exception as e:
            return ("error", str(e), {})

class ArchivePreviewHandler(PreviewHandler):
    """Handler for archive files - lists contents."""
    
    def __init__(self):
        super().__init__("Archive", ['.zip', '.rar', '.7z', '.tar', '.gz'])
    
    def generate_preview(self, file_path, max_size=1024*1024):
        try:
            if Path(file_path).suffix.lower() == '.zip':
                with zipfile.ZipFile(str(file_path), 'r') as zf:
                    file_list = zf.namelist()[:50]  # First 50 files
                    total_files = len(zf.namelist())
                    total_size = sum(f.file_size for f in zf.filelist)
                    
                    info_lines = [
                        f"ZIP Archive Contents:",
                        f"Total files: {total_files}",
                        f"Total size: {total_size / (1024*1024):.2f} MB",
                        "",
                        "First 50 files:"
                    ]
                    info_lines.extend(file_list)
                    
                    return ("text", '\n'.join(info_lines), {"total_files": total_files, "total_size_mb": total_size / (1024*1024)})
            else:
                return ("text", f"Archive preview supported for ZIP files only. File: {Path(file_path).name}", {})
        except Exception as e:
            return ("error", str(e), {})

class HexPreviewHandler(PreviewHandler):
    """Handler for binary files - shows hex dump."""
    
    def __init__(self):
        super().__init__("Hex", ['.exe', '.dll', '.bin', '.dat', '.img', '.iso'])
    
    def generate_preview(self, file_path, max_size=1024*1024):
        try:
            with open(file_path, 'rb') as f:
                data = f.read(min(max_size, 512))  # First 512 bytes
                
                # Create hex dump
                hex_lines = []
                for i in range(0, len(data), 16):
                    chunk = data[i:i+16]
                    hex_part = ' '.join(f'{b:02x}' for b in chunk)
                    ascii_part = ''.join(chr(b) if 32 <= b <= 126 else '.' for b in chunk)
                    hex_lines.append(f'{i:08x}  {hex_part:<48} |{ascii_part}|')
                
                return ("text", '\n'.join(hex_lines), {"bytes_previewed": len(data)})
        except Exception as e:
            return ("error", str(e), {})

class PreviewManager:
    """Manages all preview handlers."""
    
    def __init__(self):
        self.handlers = []
        self._register_default_handlers()
    
    def _register_default_handlers(self):
        """Register all built-in preview handlers."""
        self.handlers.extend([
            TextPreviewHandler(),
            CodePreviewHandler(),
            PDFPreviewHandler(),
            ExcelPreviewHandler(),
            CSVPreviewHandler(),
            WordPreviewHandler(),
            PowerPointPreviewHandler(),
            AudioPreviewHandler(),
            VideoPreviewHandler(),
            ArchivePreviewHandler(),
            HexPreviewHandler(),
        ])
    
    def get_handler(self, file_path):
        """Get appropriate handler for file."""
        for handler in self.handlers:
            if handler.can_handle(file_path):
                return handler
        return None
    
    def generate_preview(self, file_path, max_size=1024*1024):
        """Generate preview using appropriate handler."""
        handler = self.get_handler(file_path)
        if handler:
            return handler.generate_preview(file_path, max_size)
        return ("error", "No preview handler available for this file type", {})

class PDFViewerWidget(QWidget):
    """Full PDF viewer with pan, zoom, navigation, and search capabilities."""
    
    def __init__(self):
        super().__init__()
        self.current_pdf = None
        self.current_page = 0
        self.zoom_factor = 1.0
        self.search_results = []
        self.current_search_index = 0
        
        self.init_ui()
    
    def init_ui(self):
        """Initialize the PDF viewer UI."""
        layout = QVBoxLayout()
        
        # Toolbar with controls
        toolbar_layout = QHBoxLayout()
        
        # Navigation controls
        self.prev_btn = QPushButton("◀ Previous")
        self.prev_btn.clicked.connect(self.previous_page)
        self.prev_btn.setEnabled(False)
        
        self.next_btn = QPushButton("Next ▶")
        self.next_btn.clicked.connect(self.next_page)
        self.next_btn.setEnabled(False)
        
        self.page_label = QLabel("Page: 0 / 0")
        
        # Zoom controls
        self.zoom_out_btn = QPushButton("🔍−")
        self.zoom_out_btn.clicked.connect(self.zoom_out)
        
        self.zoom_label = QLabel("100%")
        
        self.zoom_in_btn = QPushButton("🔍+")
        self.zoom_in_btn.clicked.connect(self.zoom_in)
        
        self.fit_btn = QPushButton("Fit Width")
        self.fit_btn.clicked.connect(self.fit_width)
        
        # Pan controls
        self.pan_btn = QPushButton("✋ Drag Pan")
        self.pan_btn.setCheckable(True)
        self.pan_btn.setChecked(False)
        self.pan_btn.setToolTip("Enable mouse drag panning (click and drag to pan)")
        self.pan_btn.clicked.connect(self.toggle_pan_mode)
        
        self.reset_view_btn = QPushButton("⟲ Reset View")
        self.reset_view_btn.clicked.connect(self.reset_view)
        
        # Search controls
        self.search_input = QLineEdit()
        self.search_input.setPlaceholderText("Search in PDF...")
        self.search_input.returnPressed.connect(self.search_pdf)
        
        self.search_btn = QPushButton("🔍 Search")
        self.search_btn.clicked.connect(self.search_pdf)
        
        self.next_search_btn = QPushButton("Next ▶")
        self.next_search_btn.clicked.connect(self.next_search_result)
        self.next_search_btn.setEnabled(False)
        
        # Add controls to toolbar
        toolbar_layout.addWidget(self.prev_btn)
        toolbar_layout.addWidget(self.next_btn)
        toolbar_layout.addWidget(self.page_label)
        toolbar_layout.addStretch()
        toolbar_layout.addWidget(self.zoom_out_btn)
        toolbar_layout.addWidget(self.zoom_label)
        toolbar_layout.addWidget(self.zoom_in_btn)
        toolbar_layout.addWidget(self.fit_btn)
        toolbar_layout.addWidget(self.pan_btn)
        toolbar_layout.addWidget(self.reset_view_btn)
        toolbar_layout.addStretch()
        toolbar_layout.addWidget(self.search_input)
        toolbar_layout.addWidget(self.search_btn)
        toolbar_layout.addWidget(self.next_search_btn)
        
        # Scroll area for PDF display
        self.scroll_area = QScrollArea()
        self.scroll_area.setWidgetResizable(True)
        self.scroll_area.setVerticalScrollBarPolicy(Qt.ScrollBarPolicy.ScrollBarAsNeeded)
        self.scroll_area.setHorizontalScrollBarPolicy(Qt.ScrollBarPolicy.ScrollBarAsNeeded)
        
        # PDF display label with mouse tracking for panning
        self.pdf_label = QLabel()
        self.pdf_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self.pdf_label.setStyleSheet("background-color: #f0f0f0; border: 1px solid #ccc;")
        self.pdf_label.setText("No PDF loaded")
        self.pdf_label.setMouseTracking(True)  # Enable mouse tracking for pan hints
        
        self.scroll_area.setWidget(self.pdf_label)
        
        # Pan state for mouse drag panning
        self.pan_start_pos = None
        self.scroll_start_pos = None
        
        # Status bar
        self.status_label = QLabel("Ready")
        self.status_label.setStyleSheet("color: #666; font-size: 11px;")
        
        # Add everything to main layout
        layout.addLayout(toolbar_layout)
        layout.addWidget(self.scroll_area)
        layout.addWidget(self.status_label)
        
        self.setLayout(layout)
        
        # Enable mouse events for panning using event filter
        self.pdf_label.installEventFilter(self)
    
    def toggle_pan_mode(self):
        """Toggle drag pan mode on/off."""
        is_pan_mode = self.pan_btn.isChecked()
        if is_pan_mode:
            self.pdf_label.setCursor(Qt.CursorShape.OpenHandCursor)
            self.status_label.setText("Drag pan mode enabled - Click and drag to pan")
        else:
            self.pdf_label.setCursor(Qt.CursorShape.ArrowCursor)
            self.status_label.setText("Drag pan mode disabled")
    
    def reset_view(self):
        """Reset view to fit width and scroll to top."""
        if self.current_pdf:
            self.fit_width()
            self.scroll_area.verticalScrollBar().setValue(0)
            self.scroll_area.horizontalScrollBar().setValue(0)
            self.status_label.setText("View reset to fit width")
    
    def eventFilter(self, obj, event):
        """Handle events for mouse drag panning."""
        if obj is self.pdf_label and self.pan_btn.isChecked():
            if event.type() == QEvent.Type.MouseButtonPress:
                if self.current_pdf:
                    self.pan_start_pos = event.position()
                    self.scroll_start_pos = QPoint(
                        self.scroll_area.horizontalScrollBar().value(),
                        self.scroll_area.verticalScrollBar().value()
                    )
                    self.pdf_label.setCursor(Qt.CursorShape.ClosedHandCursor)
                    return True  # Event handled
                    
            elif event.type() == QEvent.Type.MouseMove:
                if self.pan_start_pos is not None:
                    # Calculate the delta movement
                    delta = event.position() - self.pan_start_pos
                    
                    # Update scroll position
                    new_h_scroll = self.scroll_start_pos.x() - delta.x()
                    new_v_scroll = self.scroll_start_pos.y() - delta.y()
                    
                    # Set new scroll position
                    self.scroll_area.horizontalScrollBar().setValue(int(new_h_scroll))
                    self.scroll_area.verticalScrollBar().setValue(int(new_v_scroll))
                    
                    return True  # Event handled
                    
            elif event.type() == QEvent.Type.MouseButtonRelease:
                self.pan_start_pos = None
                self.scroll_start_pos = None
                self.pdf_label.setCursor(Qt.CursorShape.OpenHandCursor)
                return True  # Event handled
        
        # Let parent class handle other events
        return super().eventFilter(obj, event)
    
    def load_pdf(self, file_path):
        """Load a PDF file for viewing."""
        try:
            if not PYMUPDF_AVAILABLE:
                self.status_label.setText("Error: PyMuPDF not installed")
                return
            
            self.current_pdf = fitz.open(str(file_path))
            self.current_page = 0
            self.zoom_factor = 1.0
            self.search_results = []
            self.current_search_index = 0
            
            self.update_page(auto_fit=True)  # Auto-fit on initial load
            self.update_controls()
            self.status_label.setText(f"Loaded: {Path(file_path).name}")
            
        except Exception as e:
            self.status_label.setText(f"Error loading PDF: {e}")
            self.pdf_label.setText("Failed to load PDF")
    
    def update_page(self, auto_fit=False, highlight_rects=None):
        """Update the current page display with optional highlighting."""
        if not self.current_pdf or self.current_page >= self.current_pdf.page_count:
            return
        
        try:
            page = self.current_pdf[self.current_page]
            
            # Only auto-fit if specifically requested or on initial load
            if auto_fit or self.zoom_factor == 1.0:
                widget_width = self.scroll_area.width() - 40  # Account for scrollbars
                if widget_width > 100:
                    page_width = page.rect.width
                    self.zoom_factor = widget_width / page_width
            
            # Render page with current zoom
            matrix = fitz.Matrix(self.zoom_factor, self.zoom_factor)
            pix = page.get_pixmap(matrix=matrix)
            
            # Convert to QImage and display
            img = QImage(pix.samples, pix.width, pix.height, pix.stride, QImage.Format.Format_RGB888)
            pixmap = QPixmap.fromImage(img)
            
            # Add highlighting if provided
            if highlight_rects:
                pixmap = self._add_highlighting(pixmap, highlight_rects)
            
            self.pdf_label.setPixmap(pixmap)
            self.pdf_label.setFixedSize(pixmap.size())
            
        except Exception as e:
            self.status_label.setText(f"Error rendering page: {e}")
    
    def _add_highlighting(self, pixmap, highlight_rects):
        """Add yellow highlighting to specified rectangles on the pixmap."""
        try:
            if not highlight_rects:
                return pixmap
            
            # Create a painter for the pixmap
            painter = QPainter(pixmap)
            
            if not painter.isActive():
                self.status_label.setText("Error: Painter not active")
                return pixmap
            
            # Set up highlight style (semi-transparent yellow)
            highlight_color = QColor(255, 255, 0, 100)  # Yellow with transparency
            painter.setBrush(QBrush(highlight_color))
            painter.setPen(Qt.PenStyle.NoPen)  # No border
            
            # Draw highlights for each rectangle
            highlight_count = 0
            for rect in highlight_rects:
                try:
                    # Scale rectangle coordinates to match zoom level
                    scaled_rect = QRect(
                        int(rect.x0 * self.zoom_factor),
                        int(rect.y0 * self.zoom_factor),
                        int((rect.x1 - rect.x0) * self.zoom_factor),
                        int((rect.y1 - rect.y0) * self.zoom_factor)
                    )
                    painter.drawRect(scaled_rect)
                    highlight_count += 1
                except Exception as e:
                    print(f"Error highlighting individual rect: {e}")
                    continue
            
            # Add red border for the current search result
            if hasattr(self, 'search_results') and self.search_results:
                if self.current_search_index < len(self.search_results):
                    current_result = self.search_results[self.current_search_index]
                    if current_result['page'] == self.current_page:
                        rect = current_result['rect']
                        scaled_rect = QRect(
                            int(rect.x0 * self.zoom_factor),
                            int(rect.y0 * self.zoom_factor),
                            int((rect.x1 - rect.x0) * self.zoom_factor),
                            int((rect.y1 - rect.y0) * self.zoom_factor)
                        )
                        painter.setPen(QPen(QColor(255, 0, 0), 3))  # Red border (thicker)
                        painter.setBrush(Qt.BrushStyle.NoBrush)
                        painter.drawRect(scaled_rect)
            
            painter.end()
            
            # Debug feedback
            if highlight_count > 0:
                print(f"Drew {highlight_count} highlights on page {self.current_page + 1}")
            
            return pixmap
            
        except Exception as e:
            import traceback
            self.status_label.setText(f"Highlighting error: {e}")
            print(f"Highlighting error details: {e}")
            traceback.print_exc()
            return pixmap
    
    def update_controls(self):
        """Update the state of navigation controls."""
        if self.current_pdf:
            total_pages = self.current_pdf.page_count
            self.page_label.setText(f"Page: {self.current_page + 1} / {total_pages}")
            
            self.prev_btn.setEnabled(self.current_page > 0)
            self.next_btn.setEnabled(self.current_page < total_pages - 1)
            
            self.zoom_label.setText(f"{int(self.zoom_factor * 100)}%")
        else:
            self.page_label.setText("Page: 0 / 0")
            self.prev_btn.setEnabled(False)
            self.next_btn.setEnabled(False)
            self.zoom_label.setText("100%")
    
    def previous_page(self):
        """Go to previous page."""
        if self.current_page > 0:
            self.current_page -= 1
            # Clear highlights when navigating normally (not via search)
            self.update_page(auto_fit=False, highlight_rects=None)
            self.update_controls()
    
    def next_page(self):
        """Go to next page."""
        if self.current_pdf and self.current_page < self.current_pdf.page_count - 1:
            self.current_page += 1
            # Clear highlights when navigating normally (not via search)
            self.update_page(auto_fit=False, highlight_rects=None)
            self.update_controls()
    
    def zoom_in(self):
        """Zoom in the PDF view."""
        self.zoom_factor = min(self.zoom_factor * 1.2, 5.0)
        # Preserve highlights if we're on a search result page
        highlight_rects = self._get_current_page_highlights() if self.search_results else None
        self.update_page(auto_fit=False, highlight_rects=highlight_rects)
        self.update_controls()
    
    def zoom_out(self):
        """Zoom out the PDF view."""
        self.zoom_factor = max(self.zoom_factor / 1.2, 0.1)
        # Preserve highlights if we're on a search result page
        highlight_rects = self._get_current_page_highlights() if self.search_results else None
        self.update_page(auto_fit=False, highlight_rects=highlight_rects)
        self.update_controls()
    
    def _get_current_page_highlights(self):
        """Get highlight rectangles for the current page if we have search results."""
        if not self.search_results:
            return None
        
        current_page_rects = []
        for search_result in self.search_results:
            if search_result['page'] == self.current_page:
                current_page_rects.append(search_result['rect'])
        
        return current_page_rects if current_page_rects else None
    
    def fit_width(self):
        """Fit PDF to widget width."""
        if self.current_pdf and self.current_page < self.current_pdf.page_count:
            try:
                page = self.current_pdf[self.current_page]
                widget_width = self.scroll_area.width() - 40
                if widget_width > 100:
                    page_width = page.rect.width
                    self.zoom_factor = widget_width / page_width
                    self.update_page(auto_fit=False)  # Don't auto-fit again
                    self.update_controls()
            except Exception as e:
                self.status_label.setText(f"Error fitting width: {e}")
    
    def search_pdf(self):
        """Search for text in the PDF with visual highlighting."""
        if not self.current_pdf or not self.search_input.text().strip():
            return
        
        search_text = self.search_input.text().strip()
        self.search_results = []
        
        try:
            # Search through all pages and collect all matches
            for page_num in range(self.current_pdf.page_count):
                page = self.current_pdf[page_num]
                text_instances = page.search_for(search_text)
                
                for inst in text_instances:
                    self.search_results.append({
                        'page': page_num,
                        'rect': inst,
                        'text': search_text
                    })
            
            if self.search_results:
                self.current_search_index = 0
                self.go_to_search_result(0)
                self.next_search_btn.setEnabled(len(self.search_results) > 1)
                self.status_label.setText(f"Found {len(self.search_results)} matches")
            else:
                # Clear any existing highlights
                self.update_page(auto_fit=False)
                self.status_label.setText("No matches found")
                self.next_search_btn.setEnabled(False)
                
        except Exception as e:
            self.status_label.setText(f"Search error: {e}")
    
    def next_search_result(self):
        """Go to next search result."""
        if self.search_results:
            self.current_search_index = (self.current_search_index + 1) % len(self.search_results)
            self.go_to_search_result(self.current_search_index)
    
    def go_to_search_result(self, index):
        """Navigate to a specific search result with visual highlighting."""
        if 0 <= index < len(self.search_results):
            result = self.search_results[index]
            
            # Go to the page with the match
            self.current_page = result['page']
            
            # Collect all rectangles for the current page to highlight them
            current_page_rects = []
            for search_result in self.search_results:
                if search_result['page'] == self.current_page:
                    current_page_rects.append(search_result['rect'])
            
            # Update page with highlighting
            self.update_page(auto_fit=False, highlight_rects=current_page_rects)
            self.update_controls()
            
            # Highlight the found text (visual feedback)
            self.status_label.setText(f"Match {index + 1}/{len(self.search_results)} on page {result['page'] + 1}")
    
    def resizeEvent(self, event):
        """Handle resize events to maintain zoom."""
        super().resizeEvent(event)
        if self.current_pdf and self.zoom_factor == 1.0:
            # Only auto-fit if still at default zoom
            QTimer.singleShot(100, self.fit_width)


class FileScoutApp(QMainWindow):
    """Main application window for File Scout."""
    def __init__(self, cli_args=None):
        super().__init__()
        self.settings = QSettings(SETTINGS_ORG, SETTINGS_APP)
        self.search_worker = None
        self.matching_files = []
        self.current_theme = 'light'
        self.duplicate_group_colors = [QColor("#FADBD8"), QColor("#D6EAF8"), QColor("#D1F2EB"), QColor("#FCF3CF"), QColor("#E8DAEF")]
        self.duplicate_group_counter = 0
        self.zoom_level = 100  # Default zoom percentage
        self.base_font_size = 9  # Base font size in points
        self.base_spacing = 10  # Base layout spacing in pixels
        self.base_groupbox_top_margin = 5  # Base top margin for QGroupBox content
        self.base_groupbox_margin_top = 10  # Base margin-top for QGroupBox in stylesheet
        self.base_groupbox_title_padding = 5  # Base title padding for QGroupBox
        
        # Result batching for performance
        self.result_batch = []
        self.batch_size = 50  # Add results in batches of 50
        self.display_limit = 5000  # Initial display limit
        self.hidden_results = []  # Results beyond display limit
        
        # Persistent File Audit Dialog (so results persist between opens)
        self.file_audit_dialog = None
        
        # System tray icon state
        self.tray_icon = None
        self.is_closing = False

        self._init_window()
        self._create_actions()
        self._create_menu_bar()
        self._create_main_widget()
        self._create_status_bar()
        self._create_system_tray()
        # Undo state for recycle-bin deletes
        self._undo_stack = []  # list of batches (list[str]) of original full paths
        self._last_delete_batch = None
        # Ensure button reflects initial state
        self._update_undo_button()
        
        self.load_settings()
        self.apply_theme(self.settings.value("theme", "light"))

        if cli_args:
            if cli_args.dir:
                # Set up directory and mode from command line
                QTimer.singleShot(100, lambda: self.setup_from_cli(cli_args))
            elif hasattr(cli_args, 'output') and cli_args.output:
                # Full CLI mode with output file
                QTimer.singleShot(100, lambda: self.run_from_cli(cli_args))

    def _init_window(self):
        self.setWindowTitle(f"{APP_NAME} v{APP_VERSION}")
        self.setMinimumSize(1000, 700)

    def _create_actions(self):
        self.save_profile_action = QAction("Save Search Profile...", self)
        self.save_profile_action.triggered.connect(self.save_search_profile)
        self.manage_profiles_action = QAction("Manage Profiles...", self)
        self.manage_profiles_action.triggered.connect(self.manage_profiles)
        # Smart Sort action
        self.smart_sort_action = QAction("Smart Sort...", self)
        self.smart_sort_action.triggered.connect(self.open_smart_sort_dialog)
        # File Audit action (Google Drive)
        self.file_audit_action = QAction("File Audit (Google Drive)...", self)
        self.file_audit_action.triggered.connect(self.open_file_audit)
        self.file_audit_action.setEnabled(FILE_AUDIT_AVAILABLE)
        # Zoom actions
        self.zoom_in_action = QAction("Zoom In", self)
        self.zoom_in_action.setShortcut("Ctrl++")
        self.zoom_in_action.triggered.connect(self.zoom_in)
        self.zoom_out_action = QAction("Zoom Out", self)
        self.zoom_out_action.setShortcut("Ctrl+-")
        self.zoom_out_action.triggered.connect(self.zoom_out)
        self.zoom_reset_action = QAction("Reset Zoom", self)
        self.zoom_reset_action.setShortcut("Ctrl+0")
        self.zoom_reset_action.triggered.connect(self.zoom_reset)
        self.exit_action = QAction("Exit", self)
        self.exit_action.triggered.connect(self.quit_application)
        
        # Minimize to tray action
        self.minimize_to_tray_action = QAction("Minimize to Tray", self)
        self.minimize_to_tray_action.setShortcut("Ctrl+H")
        self.minimize_to_tray_action.triggered.connect(self.hide)

    def _create_menu_bar(self):
        menu_bar = self.menuBar()
        file_menu = menu_bar.addMenu("&File")
        file_menu.addAction(self.save_profile_action)
        file_menu.addAction(self.manage_profiles_action)
        file_menu.addSeparator()
        file_menu.addAction(self.exit_action)

        settings_menu = menu_bar.addMenu("&Settings")
        theme_menu = settings_menu.addMenu("Theme")
        self.theme_action_group = QActionGroup(self)
        for theme_name in THEME_COLORS.keys():
            action = QAction(theme_name.replace('_', ' ').title(), self, checkable=True)
            action.triggered.connect(lambda checked, name=theme_name: self.apply_theme(name))
            self.theme_action_group.addAction(action)
            theme_menu.addAction(action)

        perf_menu = settings_menu.addMenu("Performance")
        self.count_files_checkbox = QCheckBox("Pre-count files for progress bar")
        self.count_files_checkbox.setChecked(True)
        count_action = QWidgetAction(self)
        count_action.setDefaultWidget(self.count_files_checkbox)
        perf_menu.addAction(count_action)
        
        # Tools menu
        tools_menu = menu_bar.addMenu("&Tools")
        tools_menu.addAction(self.smart_sort_action)
        tools_menu.addAction(self.file_audit_action)
        
        # View menu
        view_menu = menu_bar.addMenu("&View")
        view_menu.addAction(self.zoom_in_action)
        view_menu.addAction(self.zoom_out_action)
        view_menu.addAction(self.zoom_reset_action)
        view_menu.addSeparator()
        # Preset zoom levels
        zoom_presets_menu = view_menu.addMenu("Zoom Presets")
        for zoom_pct in [75, 100, 125, 150, 175, 200]:
            action = QAction(f"{zoom_pct}%", self)
            action.triggered.connect(lambda checked, z=zoom_pct: self.set_zoom(z))
            zoom_presets_menu.addAction(action)
        view_menu.addSeparator()
        view_menu.addAction(self.minimize_to_tray_action)

    def _create_main_widget(self):
        main_widget = QWidget()
        self.setCentralWidget(main_widget)
        main_layout = QHBoxLayout(main_widget)
        splitter = QSplitter(Qt.Orientation.Horizontal)
        main_layout.addWidget(splitter)
        search_panel = self._create_search_panel()
        results_panel = self._create_results_panel()
        splitter.addWidget(search_panel)
        splitter.addWidget(results_panel)
        splitter.setSizes([350, 650])

    def _create_search_panel(self):
        panel = QWidget()
        layout = QVBoxLayout(panel)
        layout.setSpacing(self.base_spacing)
        self.search_panel_layout = layout  # Store reference for zoom scaling

        mode_group = QGroupBox("Search Mode")
        mode_layout = QHBoxLayout()
        mode_layout.setContentsMargins(5, self.base_groupbox_top_margin, 5, 5)  # Add top margin
        self.search_mode_combo = QComboBox()
        self.search_mode_combo.addItems(["Find Files", "Find Duplicates"])
        self.search_mode_combo.setToolTip("Select 'Find Files' for normal search or 'Find Duplicates' to locate identical files.")
        self.search_mode_combo.currentIndexChanged.connect(self.toggle_search_mode_ui)
        mode_layout.addWidget(self.search_mode_combo)
        mode_group.setLayout(mode_layout)
        layout.addWidget(mode_group)

        dir_group = QGroupBox("Search Directory")
        dir_layout = QHBoxLayout()
        dir_layout.setContentsMargins(5, self.base_groupbox_top_margin, 5, 5)  # Add top margin
        self.dir_input = DropLineEdit()
        self.dir_input.setPlaceholderText("Drag & drop a folder here or click Browse...")
        browse_btn = QPushButton("Browse...")
        browse_btn.clicked.connect(self.browse_directory)
        dir_layout.addWidget(self.dir_input)
        dir_layout.addWidget(browse_btn)
        dir_group.setLayout(dir_layout)
        layout.addWidget(dir_group)

        self.keywords_group = QGroupBox("Keyword Filters")
        keywords_layout = QGridLayout()
        keywords_layout.setContentsMargins(5, self.base_groupbox_top_margin, 5, 5)  # Add top margin
        keywords_layout.addWidget(QLabel("Keywords:"), 0, 0)
        self.keywords_input = QLineEdit()
        self.keywords_input.setPlaceholderText("comma, separated, list")
        self.keywords_input.setToolTip("Find files with these words in the filename.")
        keywords_layout.addWidget(self.keywords_input, 0, 1)
        keywords_layout.addWidget(QLabel("Exclusions:"), 1, 0)
        self.exclusion_keywords_input = QLineEdit()
        self.exclusion_keywords_input.setPlaceholderText("ignore, these, words")
        self.exclusion_keywords_input.setToolTip("Exclude files with these words in the filename.")
        keywords_layout.addWidget(self.exclusion_keywords_input, 1, 1)
        
        self.regex_checkbox = QCheckBox("Use Regular Expressions")
        self.regex_checkbox.setToolTip("Treat keywords as powerful RegEx patterns.")
        self.whole_words_checkbox = QCheckBox("Match Whole Words Only")
        self.whole_words_checkbox.setToolTip("Ensure keywords match as a whole word (e.g., 'cat' won't match 'category').")
        self.regex_checkbox.stateChanged.connect(lambda state: self.whole_words_checkbox.setEnabled(not state))
        keywords_layout.addWidget(self.regex_checkbox, 2, 0, 1, 2)
        keywords_layout.addWidget(self.whole_words_checkbox, 3, 0, 1, 2)
        
        keywords_layout.addWidget(QLabel("Content Search:"), 4, 0)
        self.content_search_checkbox = QCheckBox("Search in file content")
        self.content_search_checkbox.setToolTip("Search for text within file contents (text files only).")
        keywords_layout.addWidget(self.content_search_checkbox, 4, 1)
        self.content_search_input = QLineEdit()
        self.content_search_input.setPlaceholderText("Text to find in files")
        self.content_search_input.setEnabled(False)
        keywords_layout.addWidget(self.content_search_input, 5, 0, 1, 2)
        self.content_search_checkbox.stateChanged.connect(lambda state: self.content_search_input.setEnabled(state))
        
        self.keywords_group.setLayout(keywords_layout)
        layout.addWidget(self.keywords_group)
        
        # --- MODIFIED: Added QLineEdit for custom extensions ---
        file_filter_group = QGroupBox("File Filters")
        file_filter_layout = QGridLayout()
        file_filter_layout.setContentsMargins(5, self.base_groupbox_top_margin, 5, 5)  # Add top margin
        file_filter_layout.addWidget(QLabel("Type Preset:"), 0, 0)
        self.file_types_combo = QComboBox()
        self.file_types_combo.addItems(FILE_TYPE_MAPPINGS.keys())
        file_filter_layout.addWidget(self.file_types_combo, 0, 1)
        file_filter_layout.addWidget(QLabel("Custom Exts:"), 1, 0)
        self.custom_exts_input = QLineEdit()
        self.custom_exts_input.setPlaceholderText(".ext1, .ext2, ext3")
        file_filter_layout.addWidget(self.custom_exts_input, 1, 1)
        file_filter_group.setLayout(file_filter_layout)
        layout.addWidget(file_filter_group)
        # --- END MODIFICATION ---

        size_group = QGroupBox("File Size Filter (KB)")
        size_layout = QHBoxLayout()
        size_layout.setContentsMargins(5, self.base_groupbox_top_margin, 5, 5)  # Add top margin
        size_layout.addWidget(QLabel("Min:"))
        self.min_size_input = QLineEdit()
        size_layout.addWidget(self.min_size_input)
        size_layout.addWidget(QLabel("Max:"))
        self.max_size_input = QLineEdit()
        size_layout.addWidget(self.max_size_input)
        size_group.setLayout(size_layout)
        layout.addWidget(size_group)

        date_group = QGroupBox("Date Filter")
        date_layout = QGridLayout()
        date_layout.setContentsMargins(5, self.base_groupbox_top_margin, 5, 5)  # Add top margin
        self.use_date_filter = QCheckBox("Enable Date Filter")
        date_layout.addWidget(self.use_date_filter, 0, 0, 1, 2)
        date_layout.addWidget(QLabel("Filter by:"), 1, 0)
        self.date_filter_type = QComboBox()
        self.date_filter_type.addItems(["Modified Date", "Creation Date"])
        date_layout.addWidget(self.date_filter_type, 1, 1)
        date_layout.addWidget(QLabel("From:"), 2, 0)
        self.min_date_edit = QDateEdit(calendarPopup=True, date=QDate.currentDate().addYears(-1))
        date_layout.addWidget(self.min_date_edit, 2, 1)
        date_layout.addWidget(QLabel("To:"), 3, 0)
        self.max_date_edit = QDateEdit(calendarPopup=True, date=QDate.currentDate())
        date_layout.addWidget(self.max_date_edit, 3, 1)
        date_group.setLayout(date_layout)
        layout.addWidget(date_group)

        exclude_group = QGroupBox("Excluded Directories")
        exclude_layout = QVBoxLayout()
        exclude_layout.setContentsMargins(5, self.base_groupbox_top_margin, 5, 5)  # Add top margin
        self.exclude_dirs_list = QListWidget()
        exclude_layout.addWidget(self.exclude_dirs_list)
        exclude_btn_layout = QHBoxLayout()
        add_dir_btn = QPushButton("Add")
        add_dir_btn.clicked.connect(self.add_excluded_dir)
        remove_dir_btn = QPushButton("Remove")
        remove_dir_btn.clicked.connect(self.remove_excluded_dir)
        exclude_btn_layout.addWidget(add_dir_btn)
        exclude_btn_layout.addWidget(remove_dir_btn)
        exclude_layout.addLayout(exclude_btn_layout)
        exclude_group.setLayout(exclude_layout)
        layout.addWidget(exclude_group)

        layout.addStretch()
        return panel

    def _create_results_panel(self):
        panel = QWidget()
        layout = QVBoxLayout(panel)
        layout.setSpacing(10)

        search_control_layout = QHBoxLayout()
        self.search_btn = QPushButton(" Search")
        self.search_btn.setIcon(self._get_icon(ICON_SEARCH))
        self.search_btn.clicked.connect(self.start_search)
        self.stop_btn = QPushButton(" Stop")
        self.stop_btn.setIcon(self._get_icon(ICON_STOP))
        self.stop_btn.clicked.connect(self.stop_search)
        self.stop_btn.setEnabled(False)
        search_control_layout.addWidget(self.search_btn)
        search_control_layout.addWidget(self.stop_btn)
        search_control_layout.addStretch()
        layout.addLayout(search_control_layout)

        self.progress_bar = QProgressBar()
        self.progress_bar.setProperty('theme-aware', True)  # Mark the progress bar for theme-aware styling
        layout.addWidget(self.progress_bar)

        results_group = QGroupBox("Results")
        results_layout = QVBoxLayout()
        count_layout = QHBoxLayout()
        count_layout.addWidget(QLabel("Files Found:"))
        self.file_count_label = QLabel("0")
        self.file_count_label.setStyleSheet("font-weight: bold;")
        count_layout.addWidget(self.file_count_label)
        count_layout.addStretch()
        results_layout.addLayout(count_layout)
        
        filter_layout = QHBoxLayout()
        filter_layout.addWidget(QLabel("Filter Results:"))
        self.results_filter_input = QLineEdit()
        self.results_filter_input.setPlaceholderText("Type to filter visible results...")
        self.results_filter_input.textChanged.connect(self.filter_results_table)
        filter_layout.addWidget(self.results_filter_input)
        self.export_btn = QPushButton(" Export")
        self.export_btn.setIcon(self._get_icon(ICON_EXPORT))
        self.export_btn.clicked.connect(self.show_export_options)
        self.export_btn.setEnabled(False)
        filter_layout.addWidget(self.export_btn)
        results_layout.addLayout(filter_layout)

        # Per-column filters row (built dynamically per mode)
        self.column_filters_widget = QWidget()
        self.column_filters_layout = QHBoxLayout(self.column_filters_widget)
        self.column_filters_layout.setContentsMargins(0, 0, 0, 0)
        self.column_filters_layout.setSpacing(6)
        self.column_filter_inputs = []
        results_layout.addWidget(self.column_filters_widget)

        results_splitter = QSplitter(Qt.Orientation.Vertical)
        self.results_table = QTableWidget()
        self.results_table.setSortingEnabled(True)
        self.results_table.setEditTriggers(QTableWidget.EditTrigger.NoEditTriggers)
        self.results_table.setSelectionBehavior(QTableWidget.SelectionBehavior.SelectRows)
        self.results_table.verticalHeader().setVisible(False)
        # Let specific columns stretch (Path), not just the last column
        self.results_table.horizontalHeader().setStretchLastSection(False)
        # Show middle of long paths by eliding in the center
        self.results_table.setTextElideMode(Qt.TextElideMode.ElideMiddle)
        self.results_table.setContextMenuPolicy(Qt.ContextMenuPolicy.CustomContextMenu)
        self.results_table.customContextMenuRequested.connect(self.show_results_context_menu)
        self.results_table.itemSelectionChanged.connect(self.update_selection_and_preview)
        # Track manual resizing of the Path column and respond to table resizes
        self._sizing_in_progress = False
        self.user_resized_path = False
        self.results_table.horizontalHeader().sectionResized.connect(self._on_header_section_resized)
        self.results_table.viewport().installEventFilter(self)
        results_splitter.addWidget(self.results_table)

        self.preview_widget = QWidget()
        preview_layout = QVBoxLayout(self.preview_widget)
        self.preview_tabs = QTabWidget()
        
        # Enhanced preview tabs
        self.image_preview = QLabel("Select a file to preview")
        self.image_preview.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self.image_preview.setWordWrap(True)
        
        self.text_preview = QTextEdit()
        self.text_preview.setReadOnly(True)
        self.text_preview.setFont(QFont("Consolas", 10))
        
        self.html_preview = QTextEdit()
        self.html_preview.setReadOnly(True)
        
        self.properties_preview = QTextEdit()
        self.properties_preview.setReadOnly(True)
        self.properties_preview.setFont(QFont("Consolas", 9))
        
        # Full PDF viewer widget with pan, zoom, navigation, and search
        self.pdf_viewer = PDFViewerWidget()
        
        self.preview_tabs.addTab(self.image_preview, "🖼️ Image")
        self.preview_tabs.addTab(self.text_preview, "📄 Text/Code")
        self.preview_tabs.addTab(self.html_preview, "🎨 Formatted")
        self.preview_tabs.addTab(self.pdf_viewer, "📋 PDF Viewer")
        self.preview_tabs.addTab(self.properties_preview, "ℹ️ Properties")
        
        # Initialize preview manager
        self.preview_manager = PreviewManager()
        
        preview_layout.addWidget(self.preview_tabs)
        results_splitter.addWidget(self.preview_widget)
        results_splitter.setSizes([500, 250])
        
        results_layout.addWidget(results_splitter)
        results_group.setLayout(results_layout)
        layout.addWidget(results_group)
        return panel

    def _create_status_bar(self):
        self.status_bar = QStatusBar()
        self.setStatusBar(self.status_bar)
        self.status_label = QLabel("Ready.")
        self.selection_label = QLabel("")
        self.status_bar.addWidget(self.status_label, 1)
        # Load More button for large result sets
        self.load_more_btn = QPushButton("📥 Load More")
        self.load_more_btn.setVisible(False)
        self.load_more_btn.setToolTip("Load next 1,000 results")
        self.load_more_btn.clicked.connect(self.load_more_results)
        self.status_bar.addPermanentWidget(self.load_more_btn)
        # Undo button for last recycle-bin delete
        self.undo_btn = QPushButton("Undo Delete")
        self.undo_btn.setEnabled(False)
        self.undo_btn.clicked.connect(self.undo_last_delete)
        self.status_bar.addPermanentWidget(self.undo_btn)
        self.status_bar.addPermanentWidget(self.selection_label)
    
    def _create_system_tray(self):
        """Create system tray icon with context menu."""
        # Create tray icon
        self.tray_icon = QSystemTrayIcon(self)
        
        # Try to load icon from file, otherwise create a simple one
        icon_path = Path(__file__).parent / "filescout.png"
        if icon_path.exists():
            icon = QIcon(str(icon_path))
        else:
            # Fallback: Create a simple blue square icon
            icon_pixmap = QPixmap(32, 32)
            icon_pixmap.fill(QColor(74, 144, 226))  # Blue color
            icon = QIcon(icon_pixmap)
        
        self.tray_icon.setIcon(icon)
        self.tray_icon.setToolTip(f"{APP_NAME} v{APP_VERSION}")
        
        # Create context menu
        tray_menu = QMenu()
        
        # Quick Search action
        search_action = QAction("🔍 Quick Search...", self)
        search_action.triggered.connect(self.show_and_focus)
        tray_menu.addAction(search_action)
        
        # Find Duplicates action
        duplicates_action = QAction("🔄 Find Duplicates", self)
        duplicates_action.triggered.connect(self.quick_duplicate_scan)
        tray_menu.addAction(duplicates_action)
        
        # Smart Sort action
        smart_sort_action = QAction("🗂️ Smart Sort...", self)
        smart_sort_action.triggered.connect(self.quick_smart_sort)
        tray_menu.addAction(smart_sort_action)
        
        tray_menu.addSeparator()
        
        # Recent Profiles submenu
        self.recent_profiles_menu = tray_menu.addMenu("📋 Recent Profiles")
        self._update_recent_profiles_menu()
        
        tray_menu.addSeparator()
        
        # Show/Hide window
        show_action = QAction("👁️ Show Window", self)
        show_action.triggered.connect(self.show_and_focus)
        tray_menu.addAction(show_action)
        
        hide_action = QAction("⬇️ Hide Window", self)
        hide_action.triggered.connect(self.hide)
        tray_menu.addAction(hide_action)
        
        tray_menu.addSeparator()
        
        # Exit action
        exit_action = QAction("🚪 Exit", self)
        exit_action.triggered.connect(self.quit_application)
        tray_menu.addAction(exit_action)
        
        self.tray_icon.setContextMenu(tray_menu)
        
        # Double-click to show/hide window
        self.tray_icon.activated.connect(self.tray_icon_activated)
        
        # Show the tray icon
        self.tray_icon.show()
    
    def _update_recent_profiles_menu(self):
        """Update the recent profiles submenu in system tray."""
        self.recent_profiles_menu.clear()
        profiles = self.settings.value("profiles", {})
        
        if not profiles:
            no_profiles_action = QAction("(No saved profiles)", self)
            no_profiles_action.setEnabled(False)
            self.recent_profiles_menu.addAction(no_profiles_action)
        else:
            for profile_name in sorted(profiles.keys())[:10]:  # Show max 10 profiles
                action = QAction(profile_name, self)
                action.triggered.connect(lambda checked, name=profile_name: self.load_profile_from_tray(name))
                self.recent_profiles_menu.addAction(action)
    
    def tray_icon_activated(self, reason):
        """Handle tray icon activation."""
        if reason == QSystemTrayIcon.ActivationReason.DoubleClick:
            if self.isVisible():
                self.hide()
            else:
                self.show_and_focus()
    
    def show_and_focus(self):
        """Show window and bring to front."""
        self.show()
        self.setWindowState(self.windowState() & ~Qt.WindowState.WindowMinimized | Qt.WindowState.WindowActive)
        self.activateWindow()
        self.raise_()
    
    def quick_duplicate_scan(self):
        """Quick duplicate scan from tray menu."""
        self.show_and_focus()
        self.search_mode_combo.setCurrentIndex(1)  # Set to Find Duplicates
        if self.dir_input.text() and Path(self.dir_input.text()).is_dir():
            self.start_search()
        else:
            self.browse_directory()
    
    def quick_smart_sort(self):
        """Quick smart sort from tray menu."""
        self.show_and_focus()
        if self.matching_files:
            self.open_smart_sort_dialog()
        else:
            self.show_tray_notification("No Results", "Please run a search first before using Smart Sort.", QSystemTrayIcon.MessageIcon.Warning)
    
    def load_profile_from_tray(self, profile_name):
        """Load a profile from the tray menu."""
        profiles = self.settings.value("profiles", {})
        if profile_data := profiles.get(profile_name):
            self.show_and_focus()
            self.apply_search_profile(profile_data)
            self.status_label.setText(f"Loaded profile: {profile_name}")
            self.show_tray_notification("Profile Loaded", f"Loaded search profile: {profile_name}")
    
    def show_tray_notification(self, title, message, icon=QSystemTrayIcon.MessageIcon.Information):
        """Show a system tray notification (toast)."""
        if self.tray_icon:
            self.tray_icon.showMessage(title, message, icon, 3000)  # 3 seconds
    
    def quit_application(self):
        """Properly quit the application."""
        self.is_closing = True
        if self.tray_icon:
            self.tray_icon.hide()
        QApplication.quit()

    def toggle_search_mode_ui(self, index):
        is_file_search = (index == 0)
        self.keywords_group.setEnabled(is_file_search)
        self.use_date_filter.setEnabled(is_file_search)
        self.date_filter_type.setEnabled(is_file_search)
        self.min_date_edit.setEnabled(is_file_search)
        self.max_date_edit.setEnabled(is_file_search)

    def browse_directory(self):
        directory = QFileDialog.getExistingDirectory(self, "Select Directory", self.dir_input.text())
        if directory: self.dir_input.setText(directory)

    def add_excluded_dir(self):
        directory = QFileDialog.getExistingDirectory(self, "Select Directory to Exclude")
        if directory: self.exclude_dirs_list.addItem(QListWidgetItem(Path(directory).as_posix()))

    def remove_excluded_dir(self):
        for item in self.exclude_dirs_list.selectedItems():
            self.exclude_dirs_list.takeItem(self.exclude_dirs_list.row(item))
            
    def start_search(self):
        if self.search_worker and self.search_worker.isRunning(): return
        search_dir = self.dir_input.text()
        if not search_dir or not Path(search_dir).is_dir():
            QMessageBox.warning(self, "Invalid Directory", "Please select a valid search directory.")
            return

        params = self.get_current_search_parameters(for_worker=True)
        if params is None: return
        
        self.search_btn.setEnabled(False)
        self.stop_btn.setEnabled(True)
        self.export_btn.setEnabled(False)
        self.results_table.setRowCount(0)
        self.matching_files.clear()
        self.duplicate_group_counter = 0
        # Clear result batching state for new search
        self.result_batch.clear()
        self.hidden_results.clear()
        self.display_limit = 5000
        self.preview_tabs.setTabVisible(0, False)
        self.preview_tabs.setTabVisible(1, False)

        self.search_worker = FileSearchWorker(params)
        self.search_worker.progress_update.connect(self.update_progress)
        self.search_worker.search_complete.connect(self.search_finished)
        
        if params['search_mode'] == 'duplicates':
            self.results_table.setColumnCount(4)
            self.results_table.setHorizontalHeaderLabels(["Group", "File", "Path", "Size (KB)"])
            self.results_table.setSortingEnabled(False)  # Disable during population
            self.search_worker.duplicate_group_found.connect(self.add_duplicate_group_to_table)
        else:
            self.results_table.setColumnCount(6)
            self.results_table.setHorizontalHeaderLabels(["File", "Path", "Ext", "Size (KB)", "Modified", "Created"])
            self.results_table.setSortingEnabled(False)  # Disable during population
            self.search_worker.live_result.connect(self.add_result_to_table)
        # Rebuild per-column filter inputs for current columns
        self._rebuild_column_filters()
        # Reset auto-sizing flag for a fresh layout per search
        self.user_resized_path = False
        # Apply column sizing so Path gets a sensible default width
        self._apply_results_column_sizes()
        
        self.search_worker.start()

    def stop_search(self):
        if self.search_worker and self.search_worker.isRunning():
            self.search_worker.stop()

    def update_progress_bar_style(self):
        # Get primary color from current theme for the progress bar
        if hasattr(self, 'current_theme') and hasattr(self, 'progress_bar'):
            colors = THEME_COLORS.get(self.current_theme, THEME_COLORS['light'])
            primary_color = colors.get('primary', '#3498db')  # Default to blue if primary not found
            self.progress_bar.setStyleSheet(f"QProgressBar::chunk {{ background-color: {primary_color}; }}")
    
    def update_progress(self, value, message):
        self.status_label.setText(message)
        if value >= 0:
            self.progress_bar.setRange(0, 100)
            self.progress_bar.setValue(value)
        else:
            self.progress_bar.setRange(0, 0)

    def add_result_to_table(self, file_info):
        """Add result with batching for better performance"""
        self.matching_files.append(file_info)
        self.result_batch.append(file_info)
        
        # Update count immediately
        total_count = len(self.matching_files)
        self.file_count_label.setText(str(total_count))
        
        # Batch insert for better performance
        if len(self.result_batch) >= self.batch_size:
            self._flush_result_batch()
    
    def _flush_result_batch(self):
        """Flush the result batch to the table"""
        if not self.result_batch:
            return
        
        # Disable sorting during batch insert for performance
        was_sorting_enabled = self.results_table.isSortingEnabled()
        self.results_table.setSortingEnabled(False)
        
        current_row_count = self.results_table.rowCount()
        
        for file_info in self.result_batch:
            # Check if we've hit the display limit
            if current_row_count >= self.display_limit:
                self.hidden_results.append(file_info)
                continue
            
            row = self.results_table.rowCount()
            self.results_table.insertRow(row)
            
            filename_item = QTableWidgetItem(file_info.get('filename', ''))
            filename_item.setData(Qt.ItemDataRole.UserRole, file_info)
            
            size_item = QTableWidgetItem()
            size_item.setData(Qt.ItemDataRole.EditRole, file_info.get('size_kb', 0))
            
            self.results_table.setItem(row, 0, filename_item)
            path = file_info.get('path', '')
            if not path and 'full_path' in file_info:
                path = str(Path(file_info['full_path']).parent)
            self.results_table.setItem(row, 1, QTableWidgetItem(path))
            self.results_table.setItem(row, 2, QTableWidgetItem(file_info.get('extension', '')))
            self.results_table.setItem(row, 3, size_item)
            self.results_table.setItem(row, 4, QTableWidgetItem(file_info['modified_date'].strftime('%Y-%m-%d %H:%M')))
            self.results_table.setItem(row, 5, QTableWidgetItem(file_info['created_date'].strftime('%Y-%m-%d %H:%M')))
            
            if row == 0:
                self.export_btn.setEnabled(True)
            
            current_row_count += 1
        
        # Clear batch
        self.result_batch.clear()
        
        # Restore sorting if it was enabled
        if was_sorting_enabled:
            self.results_table.setSortingEnabled(True)
        
        # Process events to keep UI responsive
        QApplication.processEvents()
    
    def load_more_results(self):
        """Load next batch of hidden results"""
        if not self.hidden_results:
            return
        
        # Load next 1000 results
        batch_to_load = self.hidden_results[:1000]
        self.hidden_results = self.hidden_results[1000:]
        
        # Temporarily increase display limit
        self.display_limit += len(batch_to_load)
        
        # Add batch to result_batch and flush
        self.result_batch.extend(batch_to_load)
        self._flush_result_batch()
        
        # Update status
        if self.hidden_results:
            self.status_label.setText(f"Loaded more results. {len(self.hidden_results)} still hidden. Click 'Load More' to see them.")
        else:
            self.status_label.setText("All results loaded.")

    def add_duplicate_group_to_table(self, file_group):
        if not file_group: return
        self.duplicate_group_counter += 1
        group_color = self.duplicate_group_colors[self.duplicate_group_counter % len(self.duplicate_group_colors)]

        for file_info in file_group:
            self.matching_files.append(file_info)
            row = self.results_table.rowCount()
            self.results_table.insertRow(row)

            group_item = QTableWidgetItem(str(self.duplicate_group_counter))
            filename_item = QTableWidgetItem(file_info.get('filename', ''))
            filename_item.setData(Qt.ItemDataRole.UserRole, file_info)
            path_item = QTableWidgetItem(file_info.get('path', ''))
            size_item = QTableWidgetItem()
            size_item.setData(Qt.ItemDataRole.EditRole, file_info.get('size_kb', 0))
            
            for item in [group_item, filename_item, path_item, size_item]:
                item.setBackground(QBrush(group_color))

            self.results_table.setItem(row, 0, group_item)
            self.results_table.setItem(row, 1, filename_item)
            self.results_table.setItem(row, 2, path_item)
            self.results_table.setItem(row, 3, size_item)

        if not self.export_btn.isEnabled(): self.export_btn.setEnabled(True)
        self.file_count_label.setText(str(len(self.matching_files)))
        QApplication.processEvents()

    def search_finished(self, success, message):
        self.search_btn.setEnabled(True)
        self.stop_btn.setEnabled(False)
        self.progress_bar.setValue(100 if success else 0)
        
        # Flush any remaining results in batch
        self._flush_result_batch()
        
        # Update status with hidden results info
        if self.hidden_results:
            hidden_count = len(self.hidden_results)
            message += f" | ⚠️ Showing first {self.display_limit:,} results. {hidden_count:,} more hidden (click 'Load More' below)."
            # Show load more button if it exists
            if hasattr(self, 'load_more_btn'):
                self.load_more_btn.setVisible(True)
                self.load_more_btn.setText(f"📥 Load More ({hidden_count:,} hidden)")
        else:
            # Hide load more button
            if hasattr(self, 'load_more_btn'):
                self.load_more_btn.setVisible(False)
        
        self.status_label.setText(message)
        
        # Re-enable sorting after population is complete
        self.results_table.setSortingEnabled(True)
        
        if success:
            self.save_last_search_parameters()
            # Show toast notification for completed search
            if self.matching_files:
                # Resize to contents, then set interactive widths with smart default for Path
                self.results_table.resizeColumnsToContents()
                self._apply_results_column_sizes()
                
                # Show notification if window is hidden
                if not self.isVisible():
                    count = len(self.matching_files)
                    search_type = "duplicate files" if self.search_mode_combo.currentIndex() == 1 else "files"
                    self.show_tray_notification(
                        "Search Complete",
                        f"Found {count:,} {search_type}",
                        QSystemTrayIcon.MessageIcon.Information
                    )

    def _apply_results_column_sizes(self):
        """Default: content-sizing for non-Path columns; Path is interactive and starts wide.

        Keeps Path user-resizable while giving it a sensible default width.
        """
        header = self.results_table.horizontalHeader()
        header.setStretchLastSection(False)
        header.setMinimumSectionSize(60)

        cols = self.results_table.columnCount()
        # Make all columns user-resizable (Interactive mode)
        for i in range(cols):
            header.setSectionResizeMode(i, QHeaderView.ResizeMode.Interactive)

        # Determine Path column index per mode
        path_col = 2 if self.search_mode_combo.currentText() == 'Find Duplicates' else 1
        # Provide a good starting width for Path column
        self._set_default_path_width(path_col)

    def _set_default_path_width(self, path_col: int):
        if self.user_resized_path:
            return  # Respect user's manual size
        self._sizing_in_progress = True
        try:
            table = self.results_table
            viewport_width = table.viewport().width()
            cols = table.columnCount()
            # Sum widths of non-path columns (after ResizeToContents)
            others = 0
            for i in range(cols):
                if i == path_col:
                    continue
                others += table.columnWidth(i)
            # Leave some margin for scrollbar and padding
            margin = 24
            target = max(150, viewport_width - others - margin)
            table.setColumnWidth(path_col, target)
        finally:
            # Delay unsetting sizing flag slightly to avoid treating our own set as user action
            self._sizing_in_progress = False

    def _on_header_section_resized(self, logicalIndex: int, oldSize: int, newSize: int):
        # If user drags the Path column, disable auto-sizing until next search
        path_col = 2 if self.search_mode_combo.currentText() == 'Find Duplicates' else 1
        if logicalIndex == path_col:
            if not self._sizing_in_progress:
                self.user_resized_path = True

    def eventFilter(self, obj, event):
        # Recalculate default Path width when the table viewport resizes (e.g., window/splitter changes)
        if obj is self.results_table.viewport() and event.type() == QEvent.Type.Resize:
            path_col = 2 if self.search_mode_combo.currentText() == 'Find Duplicates' else 1
            self._set_default_path_width(path_col)
        return False  # Don't filter the event, let it propagate normally

    def filter_results_table(self, text):
        # Combine global text and per-column filters
        self.apply_combined_filters(global_text=text)

    def _rebuild_column_filters(self):
        """Create one QLineEdit per column for per-column filtering."""
        # Clear previous inputs
        while self.column_filters_layout.count():
            item = self.column_filters_layout.takeAt(0)
            w = item.widget()
            if w:
                w.deleteLater()
        self.column_filter_inputs = []

        headers = [self.results_table.horizontalHeaderItem(i).text() for i in range(self.results_table.columnCount())]
        for idx, header in enumerate(headers):
            le = QLineEdit()
            le.setPlaceholderText(f"Filter {header}…")
            le.setToolTip(f"Filter by {header}")
            le.textChanged.connect(lambda _=None: self.apply_combined_filters())
            self.column_filters_layout.addWidget(le)
            self.column_filter_inputs.append(le)

    def _any_filters_active(self, global_text: str) -> bool:
        if global_text:
            return True
        return any(le.text() for le in self.column_filter_inputs)

    def apply_combined_filters(self, global_text: str | None = None):
        """Apply global and per-column filters together and update counts/visibility."""
        gtext = (global_text or self.results_filter_input.text()).lower().strip()
        col_filters = [le.text().lower().strip() for le in self.column_filter_inputs]

        visible_count = 0
        rows = self.results_table.rowCount()
        cols = self.results_table.columnCount()
        for i in range(rows):
            # Get column texts for this row once
            col_texts = [self.results_table.item(i, j).text().lower() if self.results_table.item(i, j) else "" for j in range(cols)]

            # Global filter: match if any column contains gtext
            global_match = True
            if gtext:
                global_match = any(gtext in t for t in col_texts)

            # Per-column filters: each non-empty filter must match its column
            percol_match = True
            for j, f in enumerate(col_filters):
                if f and f not in col_texts[j]:
                    percol_match = False
                    break

            match = global_match and percol_match
            self.results_table.setRowHidden(i, not match)
            if match:
                visible_count += 1

        # Update counts
        if self._any_filters_active(gtext):
            self.file_count_label.setText(f"{visible_count} / {len(self.matching_files)}")
        else:
            self.file_count_label.setText(str(len(self.matching_files)))

    def show_results_context_menu(self, pos):
        item = self.results_table.itemAt(pos)
        if not item: return

        is_dupe_mode = self.search_mode_combo.currentText() == 'Find Duplicates'
        file_item_col = 1 if is_dupe_mode else 0
        file_info = self.results_table.item(item.row(), file_item_col).data(Qt.ItemDataRole.UserRole)
        full_path = Path(file_info['full_path'])
        
        menu = QMenu()
        menu.addAction("Open File").triggered.connect(lambda: self._open_path(full_path))
        menu.addAction("Open Containing Folder").triggered.connect(lambda: self._open_path(full_path.parent))
        menu.addAction("Copy Full Path").triggered.connect(lambda: QApplication.clipboard().setText(str(full_path)))
        menu.addSeparator()
        # --- MODIFIED: Added "Copy" option ---
        menu.addAction("Copy Selected Files...").triggered.connect(self.copy_selected_files)
        menu.addAction("Move Selected Files...").triggered.connect(self.move_selected_files)
        menu.addAction("Delete Selected Files...").triggered.connect(self.delete_selected_files)
        menu.addSeparator()
        menu.addAction("Smart Sort...").triggered.connect(self.open_smart_sort_dialog)
    
        menu.exec(self.results_table.mapToGlobal(pos))
    
    def _open_path(self, path):
        try:
            if sys.platform == "win32": os.startfile(path)
            elif sys.platform == "darwin": subprocess.run(["open", path], check=True)
            else: subprocess.run(["xdg-open", path], check=True)
        except Exception as e:
            QMessageBox.critical(self, "Error", f"Could not open path: {e}")

    def get_selected_file_paths(self):
        paths = []
        is_dupe_mode = self.search_mode_combo.currentText() == 'Find Duplicates'
        file_item_col = 1 if is_dupe_mode else 0
        for index in self.results_table.selectionModel().selectedRows():
            # Ensure row is not hidden by the filter
            if not self.results_table.isRowHidden(index.row()):
                item = self.results_table.item(index.row(), file_item_col)
                if item and item.data(Qt.ItemDataRole.UserRole):
                    paths.append(Path(item.data(Qt.ItemDataRole.UserRole)['full_path']))
        return paths

    def delete_selected_files(self):
        paths_to_delete = self.get_selected_file_paths()
        if not paths_to_delete: return
        # Require send2trash to ensure Recycle Bin behavior
        if send2trash is None:
            QMessageBox.critical(self, "Missing Dependency", "The 'send2trash' package is required to send files to the Recycle Bin.\nInstall with: pip install send2trash")
            return
        reply = QMessageBox.question(self, "Send to Recycle Bin", f"Send {len(paths_to_delete)} file(s) to the Recycle Bin?\nYou can Undo this action from the status bar.", QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No, QMessageBox.StandardButton.No)
        if reply != QMessageBox.StandardButton.Yes: return

        rows_to_delete = sorted([index.row() for index in self.results_table.selectionModel().selectedRows() if not self.results_table.isRowHidden(index.row())], reverse=True)
        
        # We need to map visible rows to file infos correctly
        all_file_infos = [self.results_table.item(row, 1 if self.search_mode_combo.currentText() == 'Find Duplicates' else 0).data(Qt.ItemDataRole.UserRole) for row in rows_to_delete]
        
        batch = []
        errors = []
        for row, file_info in zip(rows_to_delete, all_file_infos):
            try:
                fp = str(Path(file_info['full_path']))
                send2trash(fp)
                batch.append(fp)
                self.results_table.removeRow(row)
                if file_info in self.matching_files:
                    self.matching_files.remove(file_info)
            except Exception as e:
                errors.append((file_info.get('filename', fp), str(e)))
                continue
        if batch:
            self._last_delete_batch = batch
            self._undo_stack.append(batch)
            self._update_undo_button()
            self.status_label.setText(f"Sent {len(batch)} item(s) to Recycle Bin. Click 'Undo Delete' to restore.")
        
        if errors:
            QMessageBox.warning(self, "Some Items Not Deleted", "One or more files could not be sent to the Recycle Bin.\n" + "\n".join(f"- {name}: {err}" for name, err in errors))
        
        self.file_count_label.setText(str(len(self.matching_files)))
        self.update_selection_and_preview()

    def _update_undo_button(self):
        if hasattr(self, 'undo_btn'):
            self.undo_btn.setEnabled(bool(self._undo_stack))

    def undo_last_delete(self):
        """Attempt to restore the last batch of files from the Recycle Bin.
        
        Uses PowerShell's Restore-RecycleBin if available. Provides clear feedback
        and opens Recycle Bin if automatic restore fails.
        """
        if not self._undo_stack:
            return
        batch = self._undo_stack.pop()
        restored, failed = self._restore_from_recycle_bin_powershell(batch)
        
        if restored:
            QMessageBox.information(self, "Undo Complete", f"Restored {restored} item(s) from the Recycle Bin.")
            self.status_label.setText(f"Undo restored {restored} item(s).")
        if failed:
            msg = "\n".join(f"- {p}: {err}" for p, err in failed)
            QMessageBox.warning(self, "Undo Partial", f"Some items could not be restored automatically:\n{msg}\n\nOpening Recycle Bin for manual review.")
            self.status_label.setText("Some items could not be restored. See dialog for details.")
            # Open Recycle Bin folder for manual restore
            try:
                if sys.platform == "win32":
                    os.startfile("shell:RecycleBinFolder")
            except Exception:
                pass
        # Update button state after operation
        self._update_undo_button()

    def _restore_from_recycle_bin_powershell(self, paths):
        """Restore given original paths from Recycle Bin using PowerShell.
        
        Returns (restored_count, failed_list[(path, error)])
        """
        restored = 0
        failed = []
        for p in paths:
            try:
                if sys.platform != "win32":
                    failed.append((p, 'Restore supported only on Windows'))
                    continue
                # Escape single quotes for PowerShell and normalize backslashes
                ps_path = str(Path(p)).replace("'", "''")
                ps_script = (
                    "$shell = New-Object -ComObject Shell.Application; "
                    "$rb = $shell.NameSpace(0xA); "
                    "$origIdx = -1; "
                    "for($i=0; $i -lt 256; $i++){ if($rb.GetDetailsOf($null, $i) -eq 'Original Location'){ $origIdx = $i; break } } "
                    "; if($origIdx -eq -1){ Write-Output 'ERR:Original Location column not found'; exit 1 } "
                    f"$target = '{ps_path}'; "
                    "$match = $null; foreach($it in $rb.Items()){ $orig = $rb.GetDetailsOf($it, $origIdx); $full = [System.IO.Path]::Combine($orig, $it.Name); if($full -ieq $target){ $match = $it; break } } "
                    "; if($match){ try{ $match.InvokeVerb('RESTORE') | Out-Null; Write-Output 'OK' } catch { Write-Output ('ERR:' + $_.Exception.Message) } } else { Write-Output 'MISS' }"
                )
                cmd = ["powershell", "-NoProfile", "-Command", ps_script]
                proc = subprocess.run(cmd, capture_output=True, text=True)
                out = (proc.stdout or '').strip()
                err = (proc.stderr or '').strip()
                if 'OK' in out and proc.returncode == 0:
                    restored += 1
                elif 'MISS' in out and proc.returncode == 0:
                    failed.append((p, 'Not found in Recycle Bin'))
                else:
                    failed.append((p, err or out or f'Exit code {proc.returncode}'))
            except Exception as ex:
                failed.append((p, str(ex)))
        return restored, failed

    def move_selected_files(self):
        paths_to_move = self.get_selected_file_paths()
        if not paths_to_move: return
        target_dir = QFileDialog.getExistingDirectory(self, "Select Destination Folder")
        if not target_dir: return
        
        target_path = Path(target_dir)
        rows_to_delete = sorted([index.row() for index in self.results_table.selectionModel().selectedRows() if not self.results_table.isRowHidden(index.row())], reverse=True)
        all_file_infos = [self.results_table.item(row, 1 if self.search_mode_combo.currentText() == 'Find Duplicates' else 0).data(Qt.ItemDataRole.UserRole) for row in rows_to_delete]

        for row, file_info in zip(rows_to_delete, all_file_infos):
            source_path = Path(file_info['full_path'])
            destination = target_path / source_path.name
            try:
                destination.parent.mkdir(parents=True, exist_ok=True)
                final_dest = self._generate_unique_dest_path(destination)
                source_path.rename(final_dest)
                self.results_table.removeRow(row)
                if file_info in self.matching_files:
                    self.matching_files.remove(file_info)
            except Exception as e:
                QMessageBox.critical(self, "Error", f"Could not move file '{file_info['filename']}':\n{e}")
                break
        self.file_count_label.setText(str(len(self.matching_files)))
        self.update_selection_and_preview()

    # --- NEW METHOD: Adapted from the second script's logic ---
    def _generate_unique_dest_path(self, original_dest_path: Path) -> Path:
        """Generates a unique path if original_dest_path exists by appending (1), (2), etc."""
        if not original_dest_path.exists():
            return original_dest_path

        p = original_dest_path
        base_name = p.stem
        ext = p.suffix
        parent_dir = p.parent
        counter = 1
        
        while True:
            new_name = f"{base_name} ({counter}){ext}"
            new_path = parent_dir / new_name
            if not new_path.exists():
                return new_path
            counter += 1

    # --- NEW METHOD: For handling the copy action ---
    def copy_selected_files(self):
        """Copies selected files to a destination, renaming conflicts."""
        paths_to_copy = self.get_selected_file_paths()
        if not paths_to_copy:
            return

        target_dir = QFileDialog.getExistingDirectory(self, "Select Destination Folder for Copy")
        if not target_dir:
            return

        target_path = Path(target_dir)
        success_count = 0
        error_count = 0
        
        self.status_label.setText(f"Copying {len(paths_to_copy)} files...")
        QApplication.processEvents()

        for source_path in paths_to_copy:
            try:
                original_destination = target_path / source_path.name
                final_destination = self._generate_unique_dest_path(original_destination)
                
                shutil.copy2(source_path, final_destination)
                success_count += 1
            except Exception as e:
                error_count += 1
                reply = QMessageBox.critical(self, "Copy Error", 
                                             f"Could not copy file '{source_path.name}':\n{e}\n\nContinue with other files?",
                                             QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No,
                                             QMessageBox.StandardButton.Yes)
                if reply == QMessageBox.StandardButton.No:
                    break
        
        summary_message = f"Copied {success_count} file(s)."
        if error_count > 0:
            summary_message += f" Failed to copy {error_count} file(s)."

        QMessageBox.information(self, "Copy Complete", summary_message)
        self.status_label.setText(summary_message)

    def open_smart_sort_dialog(self):
        """Open the Smart Sort dialog populated with the current visible results."""
        # Collect file_info dicts from visible rows in the results table
        files = []
        is_dupe_mode = self.search_mode_combo.currentText() == 'Find Duplicates'
        file_item_col = 1 if is_dupe_mode else 0
        for row in range(self.results_table.rowCount()):
            if self.results_table.isRowHidden(row):
                continue
            item = self.results_table.item(row, file_item_col)
            if not item:
                continue
            fi = item.data(Qt.ItemDataRole.UserRole)
            if fi:
                files.append(fi)

        if not files:
            QMessageBox.information(self, "Smart Sort", "No files available to sort. Run a search or clear filters and try again.")
            return

        default_root = self.dir_input.text() or str(Path.home())
        dlg = SmartSortDialog(self, files, default_root, zoom_level=self.zoom_level)
        dlg.exec()

    def open_file_audit(self):
        """Open the File Audit dialog for Google Drive inspection."""
        if not FILE_AUDIT_AVAILABLE:
            QMessageBox.warning(
                self,
                "Feature Unavailable",
                "File Audit feature is not available.\n\n"
                "The file_audit_dialog module could not be loaded."
            )
            return
        
        # Create dialog only once - results persist between opens
        if self.file_audit_dialog is None:
            self.file_audit_dialog = FileAuditDialog(self, theme=self.current_theme, zoom_level=self.zoom_level)
        else:
            # Update theme/zoom if changed since last open
            self.file_audit_dialog.theme = self.current_theme
            self.file_audit_dialog.zoom_level = self.zoom_level
            self.file_audit_dialog._apply_theme()
            if self.zoom_level != 100:
                self.file_audit_dialog._apply_zoom()
        
        self.file_audit_dialog.exec()

    def remove_files_from_results(self, removed_paths):
        """Remove rows and internal entries for the given moved/copied file paths.

        removed_paths: iterable of pathlib.Path or str
        """
        if not removed_paths:
            return

        # Normalize to lowercased absolute strings for robust matching
        norm_set = set()
        for p in removed_paths:
            try:
                s = str(Path(p).resolve())
            except Exception:
                s = str(p)
            norm_set.add(s.lower())

        is_dupe_mode = self.search_mode_combo.currentText() == 'Find Duplicates'
        file_item_col = 1 if is_dupe_mode else 0

        # Remove matching rows from the table, bottom-up to keep indices valid
        rows_removed = 0
        for row in range(self.results_table.rowCount() - 1, -1, -1):
            item = self.results_table.item(row, file_item_col)
            if not item:
                continue
            fi = item.data(Qt.ItemDataRole.UserRole)
            if not fi:
                continue
            try:
                fp = str(Path(fi['full_path']).resolve()).lower()
            except Exception:
                fp = str(fi.get('full_path', '')).lower()
            if fp in norm_set:
                self.results_table.removeRow(row)
                rows_removed += 1

        # Rebuild matching_files excluding removed ones for safety (sorting/filtering may de-sync indices)
        new_list = []
        for fi in getattr(self, 'matching_files', []):
            try:
                fp = str(Path(fi['full_path']).resolve()).lower()
            except Exception:
                fp = str(fi.get('full_path', '')).lower()
            if fp not in norm_set:
                new_list.append(fi)
        self.matching_files = new_list

        # Update UI
        self.file_count_label.setText(str(len(self.matching_files)))
        if self.results_table.rowCount() == 0:
            self.export_btn.setEnabled(False)
        self.update_selection_and_preview()

    def update_selection_and_preview(self):
        self.update_selection_status()
        self.update_preview()

    def update_selection_status(self):
        selected_rows = self.results_table.selectionModel().selectedRows()
        if not selected_rows:
            self.selection_label.setText("")
            return

        total_size_bytes = 0
        is_dupe_mode = self.search_mode_combo.currentText() == 'Find Duplicates'
        file_item_col = 1 if is_dupe_mode else 0
        
        visible_selected_count = 0
        for index in selected_rows:
            if not self.results_table.isRowHidden(index.row()):
                visible_selected_count += 1
                item = self.results_table.item(index.row(), file_item_col)
                if item and item.data(Qt.ItemDataRole.UserRole):
                    total_size_bytes += item.data(Qt.ItemDataRole.UserRole)['size_bytes']

        if visible_selected_count == 0:
            self.selection_label.setText("")
            return

        size_str = f"{total_size_bytes / 1024:.2f} KB"
        if total_size_bytes > 1024*1024: size_str = f"{total_size_bytes / (1024*1024):.2f} MB"
        if total_size_bytes > 1024*1024*1024: size_str = f"{total_size_bytes / (1024*1024*1024):.2f} GB"
        self.selection_label.setText(f"{visible_selected_count} items selected ({size_str})")

    def update_preview(self):
        selected_items = self.results_table.selectedItems()
        if not selected_items:
            # Hide all tabs when no selection
            for i in range(self.preview_tabs.count()):
                self.preview_tabs.setTabVisible(i, False)
            return
            
        row = selected_items[0].row()
        is_dupe_mode = self.search_mode_combo.currentText() == 'Find Duplicates'
        file_item_col = 1 if is_dupe_mode else 0
        file_info = self.results_table.item(row, file_item_col).data(Qt.ItemDataRole.UserRole)
        file_path = Path(file_info['full_path'])
        
        # Hide all tabs initially
        for i in range(self.preview_tabs.count()):
            self.preview_tabs.setTabVisible(i, False)
        
        # Show file properties in all cases
        self._show_file_properties(file_path)
        self.preview_tabs.setTabVisible(4, True)  # Properties tab
        
        # Check if it's an image first
        mime_type, _ = mimetypes.guess_type(str(file_path))
        if mime_type and mime_type.startswith('image/'):
            self._show_image_preview(file_path)
            self.preview_tabs.setCurrentIndex(0)  # Image tab
            return
            
        # Use enhanced preview system
        content_type, content, metadata = self.preview_manager.generate_preview(str(file_path))
        
        if content_type == "error":
            self.text_preview.setPlainText(f"Preview Error: {content}")
            self.preview_tabs.setTabVisible(1, True)  # Text tab
            self.preview_tabs.setCurrentIndex(1)
        elif content_type == "text":
            self.text_preview.setPlainText(content)
            self.preview_tabs.setTabVisible(1, True)  # Text tab
            if metadata.get("syntax_highlighted"):
                self.preview_tabs.setCurrentIndex(2)  # Formatted tab if syntax highlighted
            else:
                self.preview_tabs.setCurrentIndex(1)  # Text tab
        elif content_type == "pdf_dual":
            # Handle PDF with both text extraction and full PDF viewer
            self.text_preview.setPlainText(content)
            
            # Load PDF into the full viewer widget using the actual file path
            self.pdf_viewer.load_pdf(file_path)
            
            # Show tabs: Text, PDF Viewer, Properties
            self.preview_tabs.setTabVisible(1, True)  # Text tab
            self.preview_tabs.setTabVisible(3, True)  # PDF Viewer tab
            self.preview_tabs.setCurrentIndex(3)  # Start with PDF Viewer tab
            
        elif content_type == "html":
            self.html_preview.setHtml(content)
            self.preview_tabs.setTabVisible(1, True)  # Text tab
            self.preview_tabs.setTabVisible(2, True)  # Formatted tab
            self.preview_tabs.setCurrentIndex(2)  # Formatted tab

    def _show_image_preview(self, file_path):
        """Display image preview with metadata."""
        try:
            pixmap = QPixmap(str(file_path))
            if not pixmap.isNull():
                # Scale image to fit preview area
                scaled_pixmap = pixmap.scaled(400, 300, Qt.AspectRatioMode.KeepAspectRatio)
                self.image_preview.setPixmap(scaled_pixmap)
                
                # Add image info to the label
                img_size = pixmap.size()
                file_size = file_path.stat().st_size
                size_mb = file_size / (1024 * 1024)
                
                info_text = f"📐 {img_size.width()}×{img_size.height()} pixels\n"
                info_text += f"📦 {size_mb:.2f} MB\n"
                info_text += f"📁 {file_path.name}"
                
                self.image_preview.setToolTip(info_text)
                self.preview_tabs.setTabVisible(0, True)  # Image tab
            else:
                self.text_preview.setPlainText("Cannot load image file.")
                self.preview_tabs.setTabVisible(1, True)  # Text tab
                self.preview_tabs.setCurrentIndex(1)
        except Exception as e:
            self.text_preview.setPlainText(f"Error loading image: {e}")
            self.preview_tabs.setTabVisible(1, True)  # Text tab
            self.preview_tabs.setCurrentIndex(1)
    
    def _show_file_properties(self, file_path):
        """Display detailed file properties."""
        try:
            stat = file_path.stat()
            mime_type, _ = mimetypes.guess_type(str(file_path))
            
            # File size in different units
            size_bytes = stat.st_size
            if size_bytes < 1024:
                size_str = f"{size_bytes} bytes"
            elif size_bytes < 1024 * 1024:
                size_str = f"{size_bytes / 1024:.2f} KB"
            elif size_bytes < 1024 * 1024 * 1024:
                size_str = f"{size_bytes / (1024 * 1024):.2f} MB"
            else:
                size_str = f"{size_bytes / (1024 * 1024 * 1024):.2f} GB"
            
            properties = [
                "📁 File Properties",
                "=" * 40,
                f"Name: {file_path.name}",
                f"Path: {file_path.parent}",
                f"Extension: {file_path.suffix or 'None'}",
                f"Size: {size_str}",
                f"MIME Type: {mime_type or 'Unknown'}",
                "",
                "📅 Timestamps",
                "=" * 40,
                f"Created: {datetime.datetime.fromtimestamp(stat.st_ctime).strftime('%Y-%m-%d %H:%M:%S')}",
                f"Modified: {datetime.datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d %H:%M:%S')}",
                f"Accessed: {datetime.datetime.fromtimestamp(stat.st_atime).strftime('%Y-%m-%d %H:%M:%S')}",
                "",
                "🔐 Attributes",
                "=" * 40,
                f"Read-only: {'Yes' if not os.access(file_path, os.W_OK) else 'No'}",
                f"Hidden: {'Yes' if file_path.name.startswith('.') else 'No'}",
                f"Executable: {'Yes' if os.access(file_path, os.X_OK) else 'No'}",
            ]
            
            # Add file type specific info
            handler = self.preview_manager.get_handler(file_path)
            if handler:
                properties.extend([
                    "",
                    "🔍 Preview Handler",
                    "=" * 40,
                    f"Handler: {handler.name}",
                    f"Supported extensions: {', '.join(handler.extensions)}"
                ])
            
            # Add hash information for small files
            if size_bytes < 50 * 1024 * 1024:  # Files smaller than 50MB
                try:
                    with open(file_path, 'rb') as f:
                        file_hash = hashlib.md5(f.read(8192)).hexdigest()  # Hash first 8KB
                    properties.extend([
                        "",
                        "🔐 Hash (first 8KB)",
                        "=" * 40,
                        f"MD5: {file_hash}"
                    ])
                except:
                    pass
            
            self.properties_preview.setPlainText('\n'.join(properties))
            
        except Exception as e:
            self.properties_preview.setPlainText(f"Error reading file properties: {e}")

    def get_current_search_parameters(self, for_worker=False):
        params = {}
        params['search_mode'] = 'duplicates' if self.search_mode_combo.currentIndex() == 1 else 'files'
        params['search_dir'] = self.dir_input.text()
        params['keywords'] = self.keywords_input.text()
        params['exclusion_keywords'] = self.exclusion_keywords_input.text()
        params['use_regex'] = self.regex_checkbox.isChecked()
        params['whole_words'] = self.whole_words_checkbox.isChecked() and not params['use_regex']
        params['content_search'] = self.content_search_input.text() if self.content_search_checkbox.isChecked() else ""
        
        # --- MODIFIED: Custom extensions take priority over preset ---
        custom_exts_str = self.custom_exts_input.text().replace(' ', '').replace('.', '')
        if custom_exts_str:
            # If custom extensions are provided, use ONLY those (ignore preset)
            allowed_exts = set(ext.strip() for ext in custom_exts_str.split(',') if ext.strip())
        else:
            # Otherwise, use the preset
            allowed_exts = set(FILE_TYPE_MAPPINGS.get(self.file_types_combo.currentText(), []))
        params['allowed_extensions'] = list(allowed_exts)
        params['ui_file_type_preset'] = self.file_types_combo.currentText()
        params['ui_custom_exts'] = self.custom_exts_input.text()
        # --- END MODIFICATION ---

        try:
            params['min_size_kb'] = float(self.min_size_input.text()) if self.min_size_input.text() else None
            params['max_size_kb'] = float(self.max_size_input.text()) if self.max_size_input.text() else None
        except ValueError:
            QMessageBox.warning(self, "Invalid Input", "File size must be a valid number.")
            return None
        
        params['date_filter'] = self.use_date_filter.isChecked()
        if params['date_filter']:
            params['date_filter_type'] = 'modified' if self.date_filter_type.currentIndex() == 0 else 'created'
            min_qdate, max_qdate = self.min_date_edit.date(), self.max_date_edit.date()
            if for_worker:
                params['min_date'] = datetime.datetime(min_qdate.year(), min_qdate.month(), min_qdate.day(), 0, 0, 0)
                params['max_date'] = datetime.datetime(max_qdate.year(), max_qdate.month(), max_qdate.day(), 23, 59, 59)
            else:
                params['min_date_str'] = min_qdate.toString(Qt.DateFormat.ISODate)
                params['max_date_str'] = max_qdate.toString(Qt.DateFormat.ISODate)
        
        params['exclude_dirs'] = [self.exclude_dirs_list.item(i).text() for i in range(self.exclude_dirs_list.count())]
        if for_worker: 
            params['count_files'] = self.count_files_checkbox.isChecked()
            if params['search_mode'] == 'duplicates':
                 params['min_size_bytes'] = (params['min_size_kb'] or 0) * 1024
        return params

    def apply_search_profile(self, profile):
        self.search_mode_combo.setCurrentIndex(1 if profile.get('search_mode') == 'duplicates' else 0)
        self.dir_input.setText(profile.get('search_dir', ''))
        self.keywords_input.setText(profile.get('keywords', ''))
        self.exclusion_keywords_input.setText(profile.get('exclusion_keywords', ''))
        self.regex_checkbox.setChecked(profile.get('use_regex', False))
        self.whole_words_checkbox.setChecked(profile.get('whole_words', False))
        self.content_search_checkbox.setChecked(bool(profile.get('content_search', '')))
        self.content_search_input.setText(profile.get('content_search', ''))
        self.file_types_combo.setCurrentText(profile.get('ui_file_type_preset', 'All Files'))
        # --- MODIFIED: Load custom extensions from profile ---
        self.custom_exts_input.setText(profile.get('ui_custom_exts', ''))
        # --- END MODIFICATION ---
        min_size, max_size = profile.get('min_size_kb'), profile.get('max_size_kb')
        self.min_size_input.setText(str(min_size) if min_size is not None else '')
        self.max_size_input.setText(str(max_size) if max_size is not None else '')
        self.use_date_filter.setChecked(profile.get('date_filter', False))
        if profile.get('date_filter'):
            self.date_filter_type.setCurrentIndex(1 if profile.get('date_filter_type') == 'created' else 0)
            self.min_date_edit.setDate(QDate.fromString(profile.get('min_date_str'), Qt.DateFormat.ISODate))
            self.max_date_edit.setDate(QDate.fromString(profile.get('max_date_str'), Qt.DateFormat.ISODate))
        self.exclude_dirs_list.clear()
        self.exclude_dirs_list.addItems(profile.get('exclude_dirs', []))
        if 'theme' in profile:
            self.apply_theme(profile['theme'])
        if 'splitter_sizes' in profile:
            splitters = self.findChildren(QSplitter)
            for i, size in enumerate(profile['splitter_sizes']):
                if i < len(splitters):
                    splitters[i].setSizes([size, splitters[i].sizes()[1]])

    def save_search_profile(self):
        profile_name, ok = QInputDialog.getText(self, "Save Profile", "Enter profile name:")
        if not (ok and profile_name): return
        profiles = self.settings.value("profiles", {})
        if profile_name in profiles and QMessageBox.question(self, "Overwrite Profile", f"Profile '{profile_name}' already exists. Overwrite?") != QMessageBox.StandardButton.Yes:
            return
        profiles[profile_name] = {
            **self.get_current_search_parameters(for_worker=False),
            'theme': self.current_theme,
            'splitter_sizes': [splitter.sizes()[0] for splitter in self.findChildren(QSplitter)]
        }
        self.settings.setValue("profiles", profiles)
        QMessageBox.information(self, "Success", f"Profile '{profile_name}' saved.")
        # Update tray menu with new profile
        self._update_recent_profiles_menu()
        self.show_tray_notification("Profile Saved", f"Search profile '{profile_name}' saved successfully.")

    def manage_profiles(self):
        dialog = ProfileManagerDialog(self, zoom_level=self.zoom_level)
        if dialog.exec():
            profile_name = dialog.get_selected_profile_name()
            if profile_name and (profile_data := self.settings.value("profiles", {}).get(profile_name)):
                self.apply_search_profile(profile_data)
                self.status_label.setText(f"Loaded profile: {profile_name}")
                self.show_tray_notification("Profile Loaded", f"Loaded search profile: {profile_name}")
        # Update tray menu in case profiles were deleted
        self._update_recent_profiles_menu()
    
    def load_settings(self):
        self.restoreGeometry(self.settings.value("geometry", self.saveGeometry()))
        self.restoreState(self.settings.value("windowState", self.saveState()))
        self.count_files_checkbox.setChecked(self.settings.value("perf/count_files", True, type=bool))
        self.zoom_level = self.settings.value("zoom_level", 100, type=int)
        self.apply_zoom()
        self.load_last_search_parameters()

    def save_settings(self):
        self.settings.setValue("geometry", self.saveGeometry())
        self.settings.setValue("windowState", self.saveState())
        self.settings.setValue("theme", self.current_theme)
        self.settings.setValue("perf/count_files", self.count_files_checkbox.isChecked())
        self.settings.setValue("zoom_level", self.zoom_level)
        self.save_last_search_parameters()
    
    def save_last_search_parameters(self):
        if params := self.get_current_search_parameters(for_worker=False):
            self.settings.setValue("last_search_params", json.dumps(params))
            
    def load_last_search_parameters(self):
        if params_json := self.settings.value("last_search_params"):
            try:
                self.apply_search_profile(json.loads(params_json))
            except json.JSONDecodeError:
                pass
    
    def closeEvent(self, event):
        """Handle window close event - minimize to tray instead of closing."""
        if not self.is_closing:
            # Minimize to tray instead of closing
            event.ignore()
            self.hide()
            if self.tray_icon:
                self.show_tray_notification(
                    f"{APP_NAME} Minimized",
                    "File Scout is still running in the system tray. Double-click the tray icon to restore.",
                    QSystemTrayIcon.MessageIcon.Information
                )
        else:
            # Actually closing
            self.save_settings()
            if self.search_worker and self.search_worker.isRunning():
                self.search_worker.stop()
                self.search_worker.wait(1000)
            event.accept()
            super().closeEvent(event)
    
    def wheelEvent(self, event):
        """Handle Ctrl+mouse wheel for zooming"""
        if event.modifiers() == Qt.KeyboardModifier.ControlModifier:
            # Get the angle delta (positive = zoom in, negative = zoom out)
            delta = event.angleDelta().y()
            if delta > 0:
                self.zoom_in()
            elif delta < 0:
                self.zoom_out()
            event.accept()
        else:
            super().wheelEvent(event)
    
    def zoom_in(self):
        """Increase zoom level by 10%"""
        self.set_zoom(min(self.zoom_level + 10, 300))
    
    def zoom_out(self):
        """Decrease zoom level by 10%"""
        self.set_zoom(max(self.zoom_level - 10, 50))
    
    def zoom_reset(self):
        """Reset zoom to 100%"""
        self.set_zoom(100)
    
    def set_zoom(self, zoom_percent):
        """Set zoom to specific percentage"""
        self.zoom_level = zoom_percent
        self.apply_zoom()
        self.status_label.setText(f"Zoom: {self.zoom_level}%")
    
    def apply_zoom(self):
        """Apply current zoom level to all UI elements"""
        # Calculate new font size based on zoom
        new_font_size = int(self.base_font_size * (self.zoom_level / 100))
        
        # Create font for the application
        app_font = self.font()
        app_font.setPointSize(new_font_size)
        
        # Apply to main window and all widgets
        self.setFont(app_font)
        
        # Update all child widgets
        for widget in self.findChildren(QWidget):
            widget.setFont(app_font)
        
        # Special handling for tables to adjust row heights
        for table in self.findChildren(QTableWidget):
            table.verticalHeader().setDefaultSectionSize(int(25 * (self.zoom_level / 100)))
        
        # Scale layout spacing
        zoom_factor = self.zoom_level / 100
        scaled_spacing = int(self.base_spacing * zoom_factor)
        scaled_top_margin = int(self.base_groupbox_top_margin * zoom_factor)
        
        # Update search panel spacing
        if hasattr(self, 'search_panel_layout'):
            self.search_panel_layout.setSpacing(scaled_spacing)
        
        # Update all QGroupBox content margins
        for groupbox in self.findChildren(QGroupBox):
            layout = groupbox.layout()
            if layout:
                # Preserve left, right, bottom margins but scale the top margin
                margins = layout.contentsMargins()
                layout.setContentsMargins(
                    margins.left(),
                    scaled_top_margin,
                    margins.right(),
                    margins.bottom()
                )
        
        # Reapply theme to update QGroupBox margins in stylesheet
        self.apply_theme(self.current_theme)
        
        # Force update
        self.update()
    
    def _create_message_box(self, icon, title, text, buttons=QMessageBox.StandardButton.Ok):
        """Create a QMessageBox with proper zoom-scaled font"""
        msg_box = QMessageBox(icon, title, text, buttons, self)
        
        # Apply zoom to message box
        if self.zoom_level != 100:
            zoom_factor = self.zoom_level / 100
            font = msg_box.font()
            font.setPointSize(int(9 * zoom_factor))
            msg_box.setFont(font)
            
            # Also apply to all buttons
            for button in msg_box.findChildren(QPushButton):
                button.setFont(font)
        
        return msg_box
    
    def show_information(self, title, text):
        """Show information dialog with zoom-scaled font"""
        msg_box = self._create_message_box(QMessageBox.Icon.Information, title, text)
        return msg_box.exec()
    
    def show_warning(self, title, text):
        """Show warning dialog with zoom-scaled font"""
        msg_box = self._create_message_box(QMessageBox.Icon.Warning, title, text)
        return msg_box.exec()
    
    def show_question(self, title, text, buttons=QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No):
        """Show question dialog with zoom-scaled font"""
        msg_box = self._create_message_box(QMessageBox.Icon.Question, title, text, buttons)
        return msg_box.exec()
    
    def show_critical(self, title, text):
        """Show critical error dialog with zoom-scaled font"""
        msg_box = self._create_message_box(QMessageBox.Icon.Critical, title, text)
        return msg_box.exec()

    def _get_icon(self, icon_data):
        from PyQt6.QtCore import QByteArray
        pixmap = QPixmap()
        pixmap.loadFromData(QByteArray.fromBase64(icon_data.encode()))
        return QIcon(pixmap)

    def apply_theme(self, theme_name):
        self.current_theme = theme_name
        colors = THEME_COLORS.get(theme_name, THEME_COLORS['light'])
        for action in self.theme_action_group.actions():
            action.setChecked(action.text().lower().replace(' ', '_') == theme_name)
        
        # Calculate scaled values for QGroupBox based on zoom
        zoom_factor = self.zoom_level / 100
        groupbox_margin_top = int(self.base_groupbox_margin_top * zoom_factor)
        groupbox_title_padding = int(self.base_groupbox_title_padding * zoom_factor)
        
        style_sheet = f"""
        QWidget {{ background-color: {colors['background']}; color: {colors['text']}; }}
        QMainWindow, QDialog {{ background-color: {colors['background']}; }}
        QGroupBox {{ border: 1px solid {colors['border']}; border-radius: 6px; margin-top: {groupbox_margin_top}px; font-weight: bold; color: {colors['primary']}; padding-top: {groupbox_title_padding + 10}px; }}
        QGroupBox::title {{ subcontrol-origin: margin; subcontrol-position: top center; padding: {groupbox_title_padding}px {groupbox_title_padding}px; background-color: {colors['background']}; }}
        QTableWidget {{ border: 1px solid {colors['border']}; background-color: {colors['surface']}; gridline-color: {colors['grid']}; }}
        QTableWidget::item {{ border-bottom: 1px solid {colors['grid']}; padding: 5px; color: {colors['text']}; }}
        QTableWidget::item:selected {{ background-color: {colors['primary']}; color: {colors['surface']}; }}
        QHeaderView::section {{ background-color: {colors['header']}; color: {colors['text']}; padding: 6px; border: 1px solid {colors['grid']}; font-weight: bold; }}
        QLineEdit, QDateEdit, QComboBox {{ background-color: {colors['surface']}; color: {colors['text']}; border: 1px solid {colors['border']}; border-radius: 4px; padding: 5px; }}
        QLineEdit:focus, QDateEdit:focus, QComboBox:focus {{ border-color: {colors['primary']}; }}
        QListWidget {{ background-color: {colors['surface']}; color: {colors['text']}; border: 1px solid {colors['border']}; border-radius: 4px; }}
        QListWidget::item:selected {{ background-color: {colors['primary']}; color: {colors['surface']}; }}
        QPushButton {{ background-color: {colors['primary']}; color: {colors['surface']}; border: none; border-radius: 4px; padding: 6px 12px; font-weight: bold; }}
        QPushButton:hover {{ background-color: {colors['secondary']}; }}
        QPushButton:disabled {{ background-color: {colors['border']}; color: {colors['text_secondary']}; }}
        QStatusBar {{ color: {colors['text_secondary']}; }}
        QTabWidget::pane {{ border: 1px solid {colors['border']}; background-color: {colors['surface']}; }}
        QTabBar::tab {{ background: {colors['header']}; color: {colors['text']}; padding: 5px; }}
        QTabBar::tab:selected {{ background: {colors['primary']}; color: {colors['surface']}; }}
        QProgressBar {{ border: 1px solid {colors['border']}; border-radius: 4px; background-color: {colors['surface']}; text-align: center; color: {colors['text']}; }}
        QProgressBar::chunk {{ background-color: {colors['primary']}; border-radius: 3px; }}
        QCheckBox {{ color: {colors['text']}; spacing: 5px; }}
        QCheckBox::indicator {{ width: 16px; height: 16px; border: 1px solid {colors['text']}; border-radius: 3px; background-color: {colors['surface']}; }}
        QCheckBox::indicator:hover {{ border-color: {colors['primary']}; }}
        QCheckBox::indicator:checked {{ background-color: {colors['primary']}; border-color: {colors['primary']}; }}
        QCheckBox::indicator:checked:hover {{ background-color: {colors['secondary']}; border-color: {colors['secondary']}; }}
        """
        self.setStyleSheet(style_sheet)
    
    def show_export_options(self):
        menu = QMenu(self)
        menu.addAction("Export Visible to Excel...").triggered.connect(lambda: self.export_to_excel(visible_only=True))
        menu.addAction("Export All to Excel...").triggered.connect(lambda: self.export_to_excel(visible_only=False))
        menu.addSeparator()
        menu.addAction("Export Visible to CSV...").triggered.connect(lambda: self.export_to_csv(visible_only=True))
        menu.addAction("Export All to CSV...").triggered.connect(lambda: self.export_to_csv(visible_only=False))
        menu.addSeparator()
        menu.addAction("Clear Results").triggered.connect(self.clear_results)
        menu.exec(self.export_btn.mapToGlobal(self.export_btn.rect().bottomLeft()))

    def clear_results(self):
        """Clear all search results from the table."""
        if self.results_table.rowCount() == 0:
            return
        reply = QMessageBox.question(self, "Clear Results", "Clear all search results?", 
                                     QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No,
                                     QMessageBox.StandardButton.No)
        if reply == QMessageBox.StandardButton.Yes:
            self.results_table.setRowCount(0)
            self.matching_files.clear()
            self.file_count_label.setText("0")
            self.export_btn.setEnabled(False)
            self.status_label.setText("Results cleared.")
            self.preview_tabs.setTabVisible(0, False)
            self.preview_tabs.setTabVisible(1, False)

    def _get_export_data(self, visible_only=True):
        headers = [self.results_table.horizontalHeaderItem(i).text() for i in range(self.results_table.columnCount())]
        data = []
        rows_to_process = range(self.results_table.rowCount())
        if visible_only:
            rows_to_process = [r for r in rows_to_process if not self.results_table.isRowHidden(r)]
        
        for row in rows_to_process:
            row_data = {}
            for col in range(len(headers)):
                item = self.results_table.item(row, col)
                row_data[headers[col]] = item.text() if item else ""
            data.append(row_data)
        return headers, data

    def export_to_excel(self, visible_only=True):
        headers, data = self._get_export_data(visible_only)
        if not data: return
        file_path, _ = QFileDialog.getSaveFileName(self, "Save Excel File", f"file_scout_{datetime.datetime.now():%Y%m%d}.xlsx", "Excel Files (*.xlsx)")
        if not file_path: return
        try:
            ExcelExporter(theme=self.current_theme).export_data(file_path, headers, data)
            self._ask_to_open(file_path)
        except Exception as e:
            QMessageBox.critical(self, "Export Error", str(e))

    def export_to_csv(self, visible_only=True):
        headers, data = self._get_export_data(visible_only)
        if not data: return
        file_path, _ = QFileDialog.getSaveFileName(self, "Save CSV File", f"file_scout_{datetime.datetime.now():%Y%m%d}.csv", "CSV Files (*.csv)")
        if not file_path: return
        try:
            with open(file_path, 'w', newline='', encoding='utf-8') as f:
                writer = csv.DictWriter(f, fieldnames=headers)
                writer.writeheader()
                writer.writerows(data)
            self._ask_to_open(file_path)
        except Exception as e:
            QMessageBox.critical(self, "Export Error", f"Failed to export to CSV: {e}")
            
    def _ask_to_open(self, file_path):
        self.status_label.setText(f"Exported successfully to {Path(file_path).name}")
        if QMessageBox.question(self, 'Export Complete', 'Do you want to open the exported file?', QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No) == QMessageBox.StandardButton.Yes:
            self._open_path(file_path)
    
    def setup_from_cli(self, args):
        """Configure the GUI from command-line arguments (for shortcuts and context menu)."""
        # Set directory if provided
        if args.dir and os.path.isdir(args.dir):
            self.dir_input.setText(args.dir)
        
        # Set search mode
        if hasattr(args, 'duplicates') and args.duplicates:
            self.search_mode_combo.setCurrentIndex(1)  # Find Duplicates mode
        
        # Set file extensions if provided
        if args.exts:
            self.custom_exts_input.setText(args.exts)
        
        # Set size filters if provided
        if args.min_size is not None:
            self.min_size_input.setText(str(args.min_size))
        if args.max_size is not None:
            self.max_size_input.setText(str(args.max_size))
        
        # Set keywords if provided
        if args.keywords:
            self.keywords_input.setText(args.keywords)
        if hasattr(args, 'exclude_keywords') and args.exclude_keywords:
            self.exclusion_keywords_input.setText(args.exclude_keywords)
        
        # Show notification
        self.show_tray_notification(
            "File Scout Ready",
            f"Ready to search: {args.dir if args.dir else 'No directory set'}",
            QSystemTrayIcon.MessageIcon.Information
        )
        
        self.status_label.setText("Ready. Click 'Search' to begin.")

    def run_from_cli(self, args):
        print(f"--- {APP_NAME} v{APP_VERSION} CLI Mode ---")
        if args.profile:
            profiles = self.settings.value("profiles", {})
            if args.profile in profiles:
                params = profiles[args.profile]
                if args.dir: params['search_dir'] = args.dir
                if args.keywords: params['keywords'] = args.keywords
                if args.exclude_keywords: params['exclusion_keywords'] = args.exclude_keywords
                if args.exts: params['allowed_extensions'] = [e.strip().lstrip('.') for e in args.exts.split(',')]
                if args.min_size is not None: params['min_size_kb'] = args.min_size
                if args.max_size is not None: params['max_size_kb'] = args.max_size
                if args.min_date: params['min_date'] = datetime.datetime.strptime(args.min_date, '%Y-%m-%d')
                if args.max_date: params['max_date'] = datetime.datetime.strptime(args.max_date, '%Y-%m-%d')
                if args.date_type: params['date_filter_type'] = args.date_type
                params['date_filter'] = bool(args.min_date or args.max_date)
            else:
                print(f"Profile '{args.profile}' not found.", file=sys.stderr)
                sys.exit(1)
        else:
            params = {
                'search_dir': args.dir, 'keywords': args.keywords or "", 'exclusion_keywords': args.exclude_keywords or "",
                'allowed_extensions': [e.strip().lstrip('.') for e in args.exts.split(',')] if args.exts else [],
                'min_size_kb': args.min_size, 'max_size_kb': args.max_size,
                'use_regex': False, 'whole_words': False, 'date_filter': bool(args.min_date or args.max_date),
                'date_filter_type': args.date_type or 'modified',
                'min_date': datetime.datetime.strptime(args.min_date, '%Y-%m-%d') if args.min_date else None,
                'max_date': datetime.datetime.strptime(args.max_date, '%Y-%m-%d') if args.max_date else None,
                'exclude_dirs': [], 'count_files': True, 'content_search': args.content_search or ""
            }
        engine = SearchEngine()
        engine.progress_update.connect(lambda v, m: print(f"\r> {m.ljust(60)}", end=""))
        results = list(engine.find_files(params))
        print(f"\nSearch complete. Found {len(results)} files.")

        if args.output:
            try:
                with open(args.output, 'w', newline='', encoding='utf-8') as f:
                    writer = csv.DictWriter(f, fieldnames=results[0].keys() if results else [])
                    writer.writeheader()
                    writer.writerows(results)
                print(f"Results saved to {args.output}")
            except Exception as e:
                print(f"Error saving to output file: {e}", file=sys.stderr)
        else:
            for res in results: print(f"- {res['full_path']}")
        QApplication.quit()

def main():
    """Main entry point for the application."""
    parser = argparse.ArgumentParser(description=f"{APP_NAME} - File searching utility.", add_help=False)
    parser.add_argument('--dir', type=str, help="Directory to search")
    parser.add_argument('--keywords', type=str, help="Comma-separated keywords to include")
    parser.add_argument('--exclude-keywords', type=str, help="Comma-separated keywords to exclude")
    parser.add_argument('--exts', type=str, help="Comma-separated file extensions")
    parser.add_argument('--min-size', type=float, help="Minimum file size in KB")
    parser.add_argument('--max-size', type=float, help="Maximum file size in KB")
    parser.add_argument('--min-date', type=str, help="Minimum date (YYYY-MM-DD)")
    parser.add_argument('--max-date', type=str, help="Maximum date (YYYY-MM-DD)")
    parser.add_argument('--date-type', type=str, choices=['modified', 'created'], default='modified', help="Date type to filter")
    parser.add_argument('--content-search', type=str, help="Text to search within file contents")
    parser.add_argument('--profile', type=str, help="Load search parameters from a saved profile")
    parser.add_argument('--output', type=str, help="Output file path for results")
    parser.add_argument('--duplicates', action='store_true', help="Start in duplicate finding mode")
    args, _ = parser.parse_known_args()

    # Stabilize HiDPI rounding before creating the QApplication
    QGuiApplication.setHighDpiScaleFactorRoundingPolicy(Qt.HighDpiScaleFactorRoundingPolicy.RoundPreferFloor)

    app = QApplication(sys.argv)
    app.setStyle('Fusion')
    
    window = FileScoutApp(cli_args=args)
    if not (args.dir and os.path.isdir(args.dir)):
        window.show()

    sys.exit(app.exec())

if __name__ == "__main__":
    main()