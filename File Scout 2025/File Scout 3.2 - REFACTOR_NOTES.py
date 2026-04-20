"""
File Scout 3.2 - REFACTOR NOTES
================================
This file contains the original code with inline notes for refactoring into modular structure.
Notes are marked with === REFACTOR: ===
"""

import os
import sys
import re
import datetime
import json
import csv
import subprocess
import time
from pathlib import Path
import mimetypes
import hashlib
import argparse  # === REFACTOR: Unused import - can be removed ===
from collections import defaultdict
import concurrent.futures
import shutil
import zipfile
import xml.etree.ElementTree as ET
import base64
import struct

# === REFACTOR: Move all preview imports to features/preview/handlers.py ===
# Enhanced preview imports (lazy loaded)
try:
    import pygments
    from pygments import highlight
    from pygments.lexers import get_lexer_for_filename, TextLexer
    from pygments.formatters import HtmlFormatter
    PYGMENTS_AVAILABLE = True
except ImportError:
    PYGMENTS_AVAILABLE = False

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import mutagen
    MUTAGEN_AVAILABLE = True
except ImportError:
    MUTAGEN_AVAILABLE = False

# Core dependencies (always available)
try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    import xlrd
    XLRD_AVAILABLE = True
except ImportError:
    XLRD_AVAILABLE = False
try:
    from send2trash import send2trash
except Exception:
    send2trash = None

# === REFACTOR: Move to utils/config.py ===
# Quiet benign Qt logs and stabilize HiDPI behavior before importing PyQt6 modules
os.environ.setdefault("QT_LOGGING_RULES", "qt.qpa.screen=false;qt.svg.warning=false")
os.environ.setdefault("QT_SCALE_FACTOR_ROUNDING_POLICY", "RoundPreferFloor")

# === REFACTOR: Move to ui/__init__.py ===
from PyQt6.QtWidgets import (QApplication, QMainWindow, QWidget, QVBoxLayout,
                             QHBoxLayout, QLabel, QLineEdit, QPushButton,
                             QFileDialog, QProgressBar, QTableWidget, QComboBox,
                             QCheckBox, QGroupBox, QGridLayout, QMessageBox,
                             QListWidget, QSplitter, QDateEdit, QMenu,
                             QInputDialog, QTableWidgetItem, QHeaderView,
                             QDialog, QDialogButtonBox, QListWidgetItem, QWidgetAction,
                             QStatusBar, QTextEdit, QTabWidget, QSystemTrayIcon, QScrollArea,
                             QButtonGroup, QRadioButton)
from PyQt6.QtCore import Qt, QThread, pyqtSignal, QSize, QDate, QSettings, QObject, QTimer, QEvent, QPoint, QRect
from PyQt6.QtGui import QIcon, QColor, QAction, QActionGroup, QBrush, QPixmap, QGuiApplication, QFont, QImage, QPainter, QPen

# === REFACTOR: Already separate - keep as features/file_audit/google_drive.py ===
try:
    from file_audit_dialog import FileAuditDialog
    FILE_AUDIT_AVAILABLE = True
except ImportError:
    FILE_AUDIT_AVAILABLE = False

# === REFACTOR: Move to constants.py ===
# --- Constants and Utilities
APP_NAME = "File Scout"
APP_VERSION = "3.2"
SETTINGS_ORG = "WindsurfAI"
SETTINGS_APP = "FileScout"

# Performance limits for large directories
MAX_RESULTS = 50000
MAX_SCAN_FILES = 1000000
LARGE_DIR_THRESHOLD = 100000

# === REFACTOR: Move to ui/widgets.py ===
class DropLineEdit(QLineEdit):
    """A QLineEdit that accepts drag-and-drop for directory paths."""
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setAcceptDrops(True)

    def dragEnterEvent(self, event):
        if event.mimeData().hasUrls() and len(event.mimeData().urls()) == 1:
            url = event.mimeData().urls()[0]
            if url.isLocalFile() and os.path.isdir(url.toLocalFile()):
                event.acceptProposedAction()

    def dropEvent(self, event):
        url = event.mimeData().urls()[0]
        self.setText(url.toLocalFile())

# === REFACTOR: Move to utils/excel_exporter.py ===
class ExcelExporter:
    """Enhanced Excel exporter with formatting and multiple export options."""
    def __init__(self, theme='light'):
        self.theme = theme
        if theme == 'dark':
            self.header_bg = '1F1F1F'
            self.header_fg = 'FFFFFF'
            self.alt_row_bg = '2A2A2A'
        else:
            self.header_bg = '4472C4'
            self.header_fg = 'FFFFFF'
            self.alt_row_bg = 'E6E6E6'

    def export_data(self, file_path, headers, data, sheet_name="File Scout Results"):
        """Export data to Excel with formatting."""
        try:
            import openpyxl
            from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
            from openpyxl.utils import get_column_letter

            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = sheet_name

            for col, header in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col, value=header)
                cell.font = Font(bold=True, color=self.header_fg)
                cell.fill = PatternFill(start_color=self.header_bg, end_color=self.header_bg, fill_type="solid")
                cell.alignment = Alignment(horizontal='center', vertical='center')

            thin_border = Border(left=Side(style='thin'), right=Side(style='thin'), top=Side(style='thin'), bottom=Side(style='thin'))

            for row_idx, row_data in enumerate(data, 2):
                row_fill = PatternFill(start_color=self.alt_row_bg, end_color=self.alt_row_bg, fill_type="solid") if row_idx % 2 == 0 else None
                for col_idx, header in enumerate(headers, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=row_data.get(header, ""))
                    cell.border = thin_border
                    if row_fill:
                        cell.fill = row_fill

            for col in ws.columns:
                max_length = 0
                column_letter = get_column_letter(col[0].column)
                for cell in col:
                    if cell.value:
                        max_length = max(max_length, len(str(cell.value)))
                ws.column_dimensions[column_letter].width = min(max_length + 2, 50)

            ws.freeze_panes = 'A2'
            wb.save(file_path)
            return True
        except ImportError:
            raise ImportError("The 'openpyxl' library is required for Excel export.\nPlease install it: pip install openpyxl")
        except Exception as e:
            raise Exception(f"Failed to export to Excel: {str(e)}")

# === REFACTOR: Move to constants.py ===
ICON_SEARCH = "PHN2ZyB4bWxucz0iaHR0cDovL3d3dy53My5vcmcvMjAwMC9zdmciIHZpZXdCb3g9IjAgMCAyNCAyNCIgZmlsbD0iI2ZmZmZmZiI+PHBhdGggZD0iTTE1LjUgMTRoLS43OWwtLjI4LS4yN0E2LjQ3MSA2LjQ3MSAwIDAgMCAxNiA5LjVDMTYgNS45MSAxMy4wOSAzIDkuNSAzUzMgNS45MSAzIDkuNSA1LjkxIDE2IDkuNSAxNmMxLjQzIDAgMi43Ni0uNDcgMy45MS0xLjI1bC4yNy4yOHYuNzlsNSA0Ljk5TDIwLjQ5IDE5bC00Ljk5LTV6bS02IDBDNy4wMSAxNCA1IDExLjk5IDUgOS41UzcuMDEgNSA5LjUgNSAxNCA3LjAxIDE0IDkuNSAxMS45OSAxNCA5LjUgMTR6Ii8+PC9zdmc+"
ICON_STOP = "PHN2ZyB4bWxucz0iaHR0cDovL3d3dy53My5vcmcvMjAwMC9zdmciIHZpZXdCb3g9IjAgMCAyNCAyNCIgZmlsbD0iI2ZmZmZmZiI+PHBhdGggZD0iTTYgNmgxMnYxMkg2eiIvPjwvc3ZnPg=="
ICON_EXPORT = "PHN2ZyB4bWxucz0iaHR0cDovL3d3dy53My5vcmcvMjAwMC9zdmciIHZpZXdCb3g9IjAgMCAyNCAyNCIgZmlsbD0iI2ZmZmZmZiI+PHBhdGggZD0iTTE5IDEydjdINVYzdjloMlY5aDEwbDd2NGgtMnpNOCA1djhoOFY5SDh6bS0yIDE2aDEydjJIMFYxMmg0djlaIi8+PC9zdmc+"

# === REFACTOR: Move to utils/themes.py ===
THEME_COLORS = {
    'light': {
        'primary': '#4a90e2',
        'secondary': '#7ed321',
        'accent': '#f5a623',
        'background': '#f0f0f0',
        'surface': '#ffffff',
        'text': '#000000',
        'text_secondary': '#666666',
        'border': '#cccccc',
        'header': '#e8e8e8',
        'grid': '#dddddd',
        'menu': '#f8f8f8',
        'button': '#4a90e2',
        'button_hover': '#357abd',
        'success': '#7ed321',
        'warning': '#f5a623',
        'error': '#d0021b',
        'preview_bg': '#ffffff',
        'code_bg': '#f8f8f8'
    },
    'dark': {
        'primary': '#64b5f6',
        'secondary': '#81c784',
        'accent': '#ffb74d',
        'background': '#1e1e1e',
        'surface': '#2d2d30',
        'text': '#ffffff',
        'text_secondary': '#cccccc',
        'border': '#3e3e42',
        'header': '#252526',
        'grid': '#2d2d30',
        'menu': '#2d2d30',
        'button': '#64b5f6',
        'button_hover': '#42a5f5',
        'success': '#81c784',
        'warning': '#ffb74d',
        'error': '#e57373',
        'preview_bg': '#1e1e1e',
        'code_bg': '#2d2d30'
    },
    'high_contrast': {
        'primary': '#0000ff',
        'secondary': '#00ff00',
        'accent': '#ffff00',
        'background': '#000000',
        'surface': '#000000',
        'text': '#ffffff',
        'text_secondary': '#ffffff',
        'border': '#ffffff',
        'header': '#000000',
        'grid': '#ffffff',
        'menu': '#000000',
        'button': '#0000ff',
        'button_hover': '#6666ff',
        'success': '#00ff00',
        'warning': '#ffff00',
        'error': '#ff0000',
        'preview_bg': '#000000',
        'code_bg': '#000000'
    }
}

# === REFACTOR: Move to core/search_engine.py ===
class SearchEngine(QObject):
    """Core search engine with multi-threading support."""
    progress_update = pyqtSignal(int, str)
    finished = pyqtSignal(list)
    file_found = pyqtSignal(dict)
    error_occurred = pyqtSignal(str)
    
    def __init__(self):
        super().__init__()
        self.stopped = False
        self.keywords = []
        self.exclusion_keywords = []
        self.patterns = []
        self.exclusion_patterns = []
        
    # === REFACTOR: These methods belong in SearchEngine class ===
    def _compile_patterns(self, params):
        """Compile search patterns based on parameters."""
        self.keywords = [k.lower().strip() for k in params['keywords'].split(',') if k.strip()]
        self.exclusion_keywords = [k.lower().strip() for k in params['exclusion_keywords'].split(',') if k.strip()]
        if params['use_regex']:
            self.patterns = [re.compile(kw, re.IGNORECASE) for kw in self.keywords]
            self.exclusion_patterns = [re.compile(kw, re.IGNORECASE) for kw in self.exclusion_keywords]
        elif params['whole_words']:
            self.patterns = [re.compile(r'\b' + re.escape(kw) + r'\b', re.IGNORECASE) for kw in self.keywords]
            self.exclusion_patterns = [re.compile(r'\b' + re.escape(kw) + r'\b', re.IGNORECASE) for kw in self.exclusion_keywords]
        else:
            self.patterns = self.keywords
            self.exclusion_patterns = self.exclusion_keywords
    
    def _is_file_match(self, file_path, params):
        """Check if file matches search criteria."""
        if not self.keywords and not self.exclusion_keywords:
            return True
        
        # Check file extension
        if params['extensions']:
            file_ext = file_path.suffix.lower()
            if file_ext not in params['extensions']:
                return False
        
        # Check keywords
        if self.keywords:
            if not self._matches_keywords(file_path.stem):
                return False
        
        # Check exclusion keywords
        if self.exclusion_keywords:
            if self._matches_exclusion(file_path.stem):
                return False
        
        # Check content search
        if params.get('content_search'):
            if not self._matches_content(file_path, params['content_search']):
                return False
        
        return True
    
    def _matches_keywords(self, file_name):
        """Check if filename matches keywords."""
        if not self.keywords and not self.exclusion_keywords: 
            return True
        file_name_lower = file_name.lower()
        
        # Check exclusions first
        if self.exclusion_keywords:
            if self.exclusion_patterns and isinstance(self.exclusion_patterns[0], re.Pattern):
                if any(p.search(file_name_lower) for p in self.exclusion_patterns): 
                    return False
            else:
                if any(kw in file_name_lower for kw in self.exclusion_keywords): 
                    return False
        
        # Check inclusions
        if self.keywords:
            if self.patterns and isinstance(self.patterns[0], re.Pattern):
                return any(p.search(file_name_lower) for p in self.patterns)
            else:
                return any(kw in file_name_lower for kw in self.keywords)
        
        return True
    
    def _matches_exclusion(self, file_name):
        """Check if filename matches exclusion keywords."""
        file_name_lower = file_name.lower()
        if self.exclusion_patterns and isinstance(self.exclusion_patterns[0], re.Pattern):
            return any(p.search(file_name_lower) for p in self.exclusion_patterns)
        else:
            return any(kw in file_name_lower for kw in self.exclusion_keywords)
    
    def _matches_content(self, file_path, content_search):
        """Check if file content matches search term."""
        try:
            # Only search text files
            if not self._is_text_file(file_path):
                return False
            
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read(1024 * 10)  # Read first 10KB
                return content_search.lower() in content.lower()
        except Exception:
            return False
    
    def _is_text_file(self, file_path):
        """Check if file is a text file."""
        try:
            # Check by extension first
            text_extensions = {'.txt', '.log', '.csv', '.json', '.xml', '.py', '.js', '.html', '.css', '.md'}
            if file_path.suffix.lower() in text_extensions:
                return True
            
            # Try to read a small portion
            with open(file_path, 'r', encoding='utf-8') as f:
                f.read(1024)
            return True
        except Exception:
            return False
    
    # === REFACTOR: Move duplicate finding to core/duplicate_finder.py ===
    def find_duplicates(self, params):
        """Find duplicate files in directory."""
        self.progress_update.emit(0, "Stage 1/3: Grouping files by size...")
        size_map = defaultdict(list)
        total_files, processed_count = self._pre_scan(params)
        
        if self.stopped: 
            return
        
        # Group files by size
        for root, dirs, files in os.walk(params['search_dir']):
            if self.stopped: 
                return
            dirs[:] = [d for d in dirs if Path(root, d).as_posix() not in params['exclude_dirs']]
            for file in files:
                if self.stopped: 
                    return
                processed_count += 1
                if processed_count % 100 == 0:
                    progress = int((processed_count / total_files) * 30)
                    self.progress_update.emit(progress, f"Scanning files... {processed_count}/{total_files}")
                
                file_path = Path(root) / file
                try:
                    file_size = file_path.stat().st_size
                    if file_size > 0:  # Skip empty files
                        size_map[file_size].append(file_path)
                except (OSError, PermissionError):
                    continue
        
        self.progress_update.emit(30, "Stage 2/3: Computing file hashes...")
        
        # Find groups with multiple files of same size
        potential_duplicates = {size: files for size, files in size_map.items() if len(files) > 1}
        total_groups = len(potential_duplicates)
        processed_groups = 0
        
        duplicates = []
        
        # Compute hashes for potential duplicates
        for size, files in potential_duplicates.items():
            if self.stopped: 
                return
            
            processed_groups += 1
            progress = 30 + int((processed_groups / total_groups) * 60)
            self.progress_update.emit(progress, f"Checking group {processed_groups}/{total_groups}...")
            
            hash_map = defaultdict(list)
            
            for file_path in files:
                if self.stopped: 
                    return
                try:
                    file_hash = self._compute_file_hash(file_path)
                    hash_map[file_hash].append(file_path)
                except (OSError, PermissionError):
                    continue
            
            # Add actual duplicates to results
            for file_hash, dup_files in hash_map.items():
                if len(dup_files) > 1:
                    # Create file info for each duplicate
                    for file_path in dup_files:
                        try:
                            stat = file_path.stat()
                            file_info = {
                                'filename': file_path.name,
                                'full_path': str(file_path),
                                'extension': file_path.suffix.lower(),
                                'size': stat.st_size,
                                'modified': datetime.datetime.fromtimestamp(stat.st_mtime),
                                'hash': file_hash
                            }
                            duplicates.append(file_info)
                        except (OSError, PermissionError):
                            continue
        
        self.progress_update.emit(90, "Stage 3/3: Organizing results...")
        
        # Group duplicates by hash
        duplicate_groups = defaultdict(list)
        for file_info in duplicates:
            duplicate_groups[file_info['hash']].append(file_info)
        
        # Convert to list of groups
        results = list(duplicate_groups.values())
        
        self.progress_update.emit(100, f"Found {len(results)} duplicate groups")
        self.finished.emit(results)
    
    def _compute_file_hash(self, file_path):
        """Compute SHA256 hash of file."""
        hash_sha256 = hashlib.sha256()
        try:
            with open(file_path, 'rb') as f:
                # Read in chunks to handle large files
                for chunk in iter(lambda: f.read(4096), b''):
                    hash_sha256.update(chunk)
            return hash_sha256.hexdigest()
        except Exception:
            return None
    
    def _pre_scan(self, params):
        """Pre-scan directory to count total files."""
        count = 0
        for root, dirs, files in os.walk(params['search_dir']):
            dirs[:] = [d for d in dirs if Path(root, d).as_posix() not in params['exclude_dirs']]
            count += len(files)
            if count > MAX_SCAN_FILES:
                return count, count
        return count, 0
    
    # === REFACTOR: Move to core/file_scanner.py ===
    def search_files(self, params):
        """Search for files matching criteria."""
        self._compile_patterns(params)
        
        # Count total files for progress
        total_files = self._count_files(params)
        if total_files == -1:
            return
        
        self.progress_update.emit(0, f"Searching {total_files} files...")
        
        results = []
        processed = 0
        
        # Use multi-threading for large directories
        if total_files > LARGE_DIR_THRESHOLD and params.get('use_multithreading', True):
            results = list(self._search_multithreaded(params, total_files))
        else:
            results = list(self._search_singlethreaded(params, total_files, processed))
        
        # Filter results if needed
        if len(results) > MAX_RESULTS:
            results = results[:MAX_RESULTS]
            self.error_occurred.emit(f"Warning: Limited to first {MAX_RESULTS} results")
        
        self.progress_update.emit(100, f"Found {len(results)} matching files")
        self.finished.emit(results)
    
    def _search_singlethreaded(self, params, total_files, processed_count):
        """Original single-threaded implementation."""
        for root, dirs, files in os.walk(params['search_dir']):
            if self.stopped: 
                break
            dirs[:] = [d for d in dirs if Path(root, d).as_posix() not in params['exclude_dirs']]
            
            # Update progress with current folder name
            current_folder = Path(root).name or str(root)
            if processed_count % 50 == 0:
                progress = int((processed_count / total_files) * 100)
                self.progress_update.emit(progress, f"Searching in {current_folder}...")
            
            for file in files:
                if self.stopped: 
                    break
                processed_count += 1
                
                file_path = Path(root) / file
                
                # Skip if not matching
                if not self._is_file_match(file_path, params):
                    continue
                
                # Get file info
                try:
                    stat = file_path.stat()
                    file_info = {
                        'filename': file_path.name,
                        'full_path': str(file_path),
                        'extension': file_path.suffix.lower(),
                        'size': stat.st_size,
                        'modified': datetime.datetime.fromtimestamp(stat.st_mtime)
                    }
                    
                    # Add modifier if available
                    if sys.platform == "win32":
                        try:
                            import win32security
                            import win32api
                            import win32con
                            sd = win32security.GetFileSecurity(str(file_path), win32security.OWNER_SECURITY_INFORMATION)
                            owner_sid = sd.GetSecurityDescriptorOwner()
                            owner_name = win32security.LookupAccountSid(None, owner_sid)[0]
                            file_info['modifier'] = owner_name
                        except:
                            file_info['modifier'] = 'Unknown'
                    
                    results.append(file_info)
                    self.file_found.emit(file_info)
                    
                except (OSError, PermissionError):
                    continue
        
        yield from results
    
    def _search_multithreaded(self, params, total_files):
        """Multi-threaded search implementation."""
        import queue
        
        # Create a queue for directories to process
        dir_queue = queue.Queue()
        result_queue = queue.Queue()
        
        # Collect all directories first
        for root, dirs, files in os.walk(params['search_dir']):
            dirs[:] = [d for d in dirs if Path(root, d).as_posix() not in params['exclude_dirs']]
            if files:  # Only add non-empty directories
                dir_queue.put(root)
        
        # Worker function
        def scan_directory(dir_path):
            """Scan a single directory and its subdirectories."""
            local_results = []
            try:
                for root, dirs, files in os.walk(dir_path):
                    if self.stopped: 
                        break
                    
                    # Filter excluded dirs
                    dirs[:] = [d for d in dirs if Path(root, d).as_posix() not in params['exclude_dirs']]
                    
                    for file in files:
                        if self.stopped: 
                            break
                        
                        file_path = Path(root) / file
                        
                        # Skip if not matching
                        if not self._is_file_match(file_path, params):
                            continue
                        
                        # Get file info
                        try:
                            stat = file_path.stat()
                            file_info = {
                                'filename': file_path.name,
                                'full_path': str(file_path),
                                'extension': file_path.suffix.lower(),
                                'size': stat.st_size,
                                'modified': datetime.datetime.fromtimestamp(stat.st_mtime)
                            }
                            
                            # Add modifier if available
                            if sys.platform == "win32":
                                try:
                                    import win32security
                                    import win32api
                                    import win32con
                                    sd = win32security.GetFileSecurity(str(file_path), win32security.OWNER_SECURITY_INFORMATION)
                                    owner_sid = sd.GetSecurityDescriptorOwner()
                                    owner_name = win32security.LookupAccountSid(None, owner_sid)[0]
                                    file_info['modifier'] = owner_name
                                except:
                                    file_info['modifier'] = 'Unknown'
                            
                            local_results.append(file_info)
                            result_queue.put(file_info)
                            
                        except (OSError, PermissionError):
                            continue
            except Exception as e:
                self.error_occurred.emit(f"Error scanning {dir_path}: {str(e)}")
            
            return local_results
        
        # Start worker threads
        num_threads = min(4, os.cpu_count() or 2)
        futures = []
        
        with concurrent.futures.ThreadPoolExecutor(max_workers=num_threads) as executor:
            # Submit jobs
            for _ in range(num_threads):
                if not dir_queue.empty():
                    dir_path = dir_queue.get()
                    future = executor.submit(scan_directory, dir_path)
                    futures.append(future)
            
            # Collect results
            processed = 0
            while any(not f.done() for f in futures) or not result_queue.empty():
                try:
                    # Get results from queue
                    while not result_queue.empty():
                        result = result_queue.get_nowait()
                        processed += 1
                        self.file_found.emit(result)
                        
                        # Update progress
                        if processed % 100 == 0:
                            progress = min(95, int((processed / total_files) * 100))
                            self.progress_update.emit(progress, f"Found {processed} files...")
                        
                        yield result
                        
                except queue.Empty:
                    pass
                
                # Check for completed threads and submit new jobs
                for i, future in enumerate(futures):
                    if future.done() and not dir_queue.empty():
                        dir_path = dir_queue.get()
                        futures[i] = executor.submit(scan_directory, dir_path)
                
                time.sleep(0.01)  # Small delay to prevent busy waiting
    
    def _count_files(self, params):
        """Count total files to process."""
        self.progress_update.emit(0, "Counting total files...")
        count = 0
        for root, dirs, files in os.walk(params['search_dir']):
            if self.stopped: 
                return -1
            dirs[:] = [d for d in dirs if Path(root, d).as_posix() not in params['exclude_dirs']]
            count += len(files)
            if count > MAX_SCAN_FILES:
                return count
        self.progress_update.emit(0, f"Found {count} total files to process.")
        return count
    
    def stop(self):
        """Stop the search operation."""
        self.stopped = True

# === REFACTOR: Move to core/file_scanner.py ===
class FileSearchWorker(QThread):
    """Worker thread for file search operations."""
    progress_update = pyqtSignal(int, str)
    finished = pyqtSignal(list)
    file_found = pyqtSignal(dict)
    error_occurred = pyqtSignal(str)
    
    def __init__(self, search_params):
        super().__init__()
        self.search_params = search_params
        self.search_engine = SearchEngine()
        
        # Connect signals
        self.search_engine.progress_update.connect(self.progress_update)
        self.search_engine.finished.connect(self._on_finished)
        self.search_engine.file_found.connect(self.file_found)
        self.search_engine.error_occurred.connect(self.error_occurred)
    
    def run(self):
        """Run the search operation."""
        try:
            if self.search_params.get('search_mode') == 'duplicates':
                self.search_engine.find_duplicates(self.search_params)
            else:
                self.search_engine.search_files(self.search_params)
        except Exception as e:
            self.error_occurred.emit(f"Search error: {str(e)}")
            self.finished.emit([])
    
    def _on_finished(self, results):
        """Handle search completion."""
        self.finished.emit(results)
    
    def stop(self):
        """Stop the search."""
        self.search_engine.stop()

# === REFACTOR: Move to ui/profile_manager.py ===
class ProfileManager(QDialog):
    """Dialog for managing search profiles."""
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle("Manage Search Profiles")
        self.setMinimumSize(600, 400)
        self.parent = parent
        self.settings = QSettings(SETTINGS_ORG, SETTINGS_APP)
        
        layout = QVBoxLayout(self)
        
        # Profile list
        self.profile_list = QListWidget()
        layout.addWidget(QLabel("Saved Profiles:"))
        layout.addWidget(self.profile_list)
        
        # Buttons
        button_layout = QHBoxLayout()
        self.load_btn = QPushButton("Load")
        self.delete_btn = QPushButton("Delete")
        self.export_btn = QPushButton("Export")
        self.import_btn = QPushButton("Import")
        button_layout.addWidget(self.load_btn)
        button_layout.addWidget(self.delete_btn)
        button_layout.addWidget(self.export_btn)
        button_layout.addWidget(self.import_btn)
        layout.addLayout(button_layout)
        
        # Close button
        close_btn = QPushButton("Close")
        close_btn.clicked.connect(self.accept)
        layout.addWidget(close_btn)
        
        # Load profiles
        self.load_profiles()
        
        # Connect signals
        self.load_btn.clicked.connect(self.load_profile)
        self.delete_btn.clicked.connect(self.delete_profile)
        self.export_btn.clicked.connect(self.export_profile)
        self.import_btn.clicked.connect(self.import_profile)
    
    def load_profiles(self):
        """Load profiles from settings."""
        self.profile_list.clear()
        profiles = self.settings.value("profiles", {})
        for name in profiles.keys():
            self.profile_list.addItem(name)
    
    def load_profile(self):
        """Load selected profile."""
        current_item = self.profile_list.currentItem()
        if not current_item:
            return
        
        profile_name = current_item.text()
        profiles = self.settings.value("profiles", {})
        if profile_name in profiles:
            self.parent.apply_search_profile(profiles[profile_name])
            self.accept()
    
    def delete_profile(self):
        """Delete selected profile."""
        current_item = self.profile_list.currentItem()
        if not current_item:
            return
        
        reply = QMessageBox.question(
            self, "Delete Profile",
            f"Delete profile '{current_item.text()}'?",
            QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No
        )
        
        if reply == QMessageBox.StandardButton.Yes:
            profiles = self.settings.value("profiles", {})
            profiles.pop(current_item.text(), None)
            self.settings.setValue("profiles", profiles)
            self.load_profiles()
    
    def export_profile(self):
        """Export selected profile to file."""
        current_item = self.profile_list.currentItem()
        if not current_item:
            return
        
        profile_name = current_item.text()
        profiles = self.settings.value("profiles", {})
        if profile_name in profiles:
            file_path, _ = QFileDialog.getSaveFileName(
                self, "Export Profile",
                f"{profile_name}.json",
                "JSON Files (*.json)"
            )
            if file_path:
                with open(file_path, 'w') as f:
                    json.dump({profile_name: profiles[profile_name]}, f, indent=2)
                QMessageBox.information(self, "Export Complete", f"Profile exported to {file_path}")
    
    def import_profile(self):
        """Import profile from file."""
        file_path, _ = QFileDialog.getOpenFileName(
            self, "Import Profile",
            "",
            "JSON Files (*.json)"
        )
        if file_path:
            try:
                with open(file_path, 'r') as f:
                    imported_profiles = json.load(f)
                
                profiles = self.settings.value("profiles", {})
                profiles.update(imported_profiles)
                self.settings.setValue("profiles", profiles)
                self.load_profiles()
                QMessageBox.information(self, "Import Complete", "Profile imported successfully")
            except Exception as e:
                QMessageBox.critical(self, "Import Error", f"Failed to import profile: {str(e)}")
    
    def get_selected_profile_name(self):
        """Get selected profile name."""
        return self.profile_list.currentItem().text() if self.profile_list.currentItem() else None

# === REFACTOR: Move to features/smart_sort/smart_sort_dialog.py ===
# [SmartSortDialog class would go here - already extracted to separate file]

# === REFACTOR: Move to features/preview/handlers.py ===
# Enhanced Preview Handler System
class PreviewHandler:
    """Base class for file preview handlers."""
    
    def __init__(self, name, extensions=None):
        self.name = name
        self.extensions = extensions or []
    
    def can_handle(self, file_path):
        """Check if this handler can preview the file."""
        file_path = Path(file_path)
        return file_path.suffix.lower() in self.extensions
    
    def generate_preview(self, file_path, max_size=1024*1024):
        """Generate preview content. Returns (content_type, data, metadata)."""
        raise NotImplementedError

class TextPreviewHandler(PreviewHandler):
    """Handler for text files with syntax highlighting."""
    
    def __init__(self):
        super().__init__("Text", ['.txt', '.log', '.cfg', '.ini', '.conf', '.md', '.rst'])
    
    def generate_preview(self, file_path, max_size=1024*1024):
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read(min(max_size, 5000))  # Read up to 5KB
            
            if PYGMENTS_AVAILABLE:
                try:
                    lexer = get_lexer_for_filename(file_path)
                except:
                    lexer = TextLexer()
                
                formatter = HtmlFormatter(style='default', cssclass='source')
                highlighted = highlight(content, lexer, formatter)
                return 'html', highlighted, {'language': lexer.name}
            else:
                return 'text', content, {'language': 'text'}
        except Exception as e:
            return 'error', str(e), {}

class CodePreviewHandler(PreviewHandler):
    """Handler for code files with syntax highlighting."""
    
    def __init__(self):
        super().__init__("Code", [
            '.py', '.js', '.html', '.css', '.xml', '.json', '.yaml', '.yml',
            '.sql', '.sh', '.bat', '.ps1', '.java', '.cpp', '.c', '.h',
            '.cs', '.php', '.rb', '.go', '.rs', '.swift', '.kt'
        ])
    
    def generate_preview(self, file_path, max_size=1024*1024):
        if not PYGMENTS_AVAILABLE:
            return 'error', 'Pygments not available for code highlighting', {}
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read(min(max_size, 10000))  # Read up to 10KB
            
            try:
                lexer = get_lexer_for_filename(file_path)
            except:
                lexer = TextLexer()
            
            formatter = HtmlFormatter(style='default', cssclass='source', linenos='table')
            highlighted = highlight(content, lexer, formatter)
            return 'html', highlighted, {'language': lexer.name, 'lines': True}
        except Exception as e:
            return 'error', str(e), {}

class PDFPreviewHandler(PreviewHandler):
    """Handler for PDF files using PyMuPDF."""
    
    def __init__(self):
        super().__init__("PDF", ['.pdf'])
    
    def generate_preview(self, file_path, max_size=1024*1024):
        if not PYMUPDF_AVAILABLE:
            return 'error', 'PyMuPDF not available for PDF preview', {}
        
        try:
            doc = fitz.open(file_path)
            if doc.page_count > 0:
                page = doc[0]
                pix = page.get_pixmap()
                img_data = pix.tobytes("png")
                doc.close()
                return 'image', img_data, {'pages': doc.page_count}
            else:
                return 'error', 'PDF has no pages', {}
        except Exception as e:
            return 'error', str(e), {}

class ImagePreviewHandler(PreviewHandler):
    """Handler for image files."""
    
    def __init__(self):
        super().__init__("Image", [
            '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp',
            '.ico', '.svg'
        ])
    
    def generate_preview(self, file_path, max_size=1024*1024):
        try:
            from PyQt6.QtGui import QImage
            
            image = QImage(file_path)
            if image.isNull():
                return 'error', 'Failed to load image', {}
            
            # Scale image if too large
            if image.width() > 800 or image.height() > 600:
                image = image.scaled(800, 600, Qt.AspectRatioMode.KeepAspectRatio, Qt.TransformationMode.SmoothTransformation)
            
            # Convert to bytes
            buffer = QBuffer()
            buffer.open(QBuffer.OpenModeFlag.ReadWrite)
            image.save(buffer, "PNG")
            img_data = buffer.data()
            buffer.close()
            
            return 'image', img_data, {
                'width': image.width(),
                'height': image.height(),
                'format': image.format()
            }
        except Exception as e:
            return 'error', str(e), {}

class DocxPreviewHandler(PreviewHandler):
    """Handler for Word documents."""
    
    def __init__(self):
        super().__init__("Word", ['.docx'])
    
    def generate_preview(self, file_path, max_size=1024*1024):
        if not DOCX_AVAILABLE:
            return 'error', 'python-docx not available for Word preview', {}
        
        try:
            doc = Document(file_path)
            paragraphs = []
            for para in doc.paragraphs[:20]:  # First 20 paragraphs
                if para.text.strip():
                    paragraphs.append(para.text)
            
            content = '\n\n'.join(paragraphs)
            return 'text', content, {'paragraphs': len(doc.paragraphs)}
        except Exception as e:
            return 'error', str(e), {}

class PPTXPreviewHandler(PreviewHandler):
    """Handler for PowerPoint presentations."""
    
    def __init__(self):
        super().__init__("PowerPoint", ['.pptx'])
    
    def generate_preview(self, file_path, max_size=1024*1024):
        if not PPTX_AVAILABLE:
            return 'error', 'python-pptx not available for PowerPoint preview', {}
        
        try:
            prs = Presentation(file_path)
            slides_text = []
            for slide in prs.slides[:5]:  # First 5 slides
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                if slide_text:
                    slides_text.append(f"Slide {len(slides_text)+1}:\n" + "\n".join(slide_text))
            
            content = "\n\n---\n\n".join(slides_text)
            return 'text', content, {'slides': len(prs.slides)}
        except Exception as e:
            return 'error', str(e), {}

class AudioPreviewHandler(PreviewHandler):
    """Handler for audio files metadata."""
    
    def __init__(self):
        super().__init__("Audio", [
            '.mp3', '.wav', '.flac', '.ogg', '.m4a', '.aac', '.wma'
        ])
    
    def generate_preview(self, file_path, max_size=1024*1024):
        if not MUTAGEN_AVAILABLE:
            return 'error', 'mutagen not available for audio metadata', {}
        
        try:
            audio = mutagen.File(file_path)
            if audio is None:
                return 'error', 'Could not read audio file', {}
            
            metadata = {}
            for key, value in audio.items():
                if isinstance(value, list) and value:
                    metadata[key] = str(value[0])
                else:
                    metadata[key] = str(value)
            
            # Format metadata as text
            content = "Audio Metadata:\n\n"
            for key, value in metadata.items():
                content += f"{key}: {value}\n"
            
            return 'text', content, metadata
        except Exception as e:
            return 'error', str(e), {}

# === REFACTOR: Move to features/preview/preview_manager.py ===
class PreviewManager:
    """Manages file preview handlers."""
    
    def __init__(self):
        self.handlers = []
        self._register_default_handlers()
    
    def _register_default_handlers(self):
        """Register default preview handlers."""
        self.handlers = [
            TextPreviewHandler(),
            CodePreviewHandler(),
            PDFPreviewHandler(),
            ImagePreviewHandler(),
            DocxPreviewHandler(),
            PPTXPreviewHandler(),
            AudioPreviewHandler()
        ]
    
    def register_handler(self, handler):
        """Register a custom preview handler."""
        self.handlers.append(handler)
    
    def get_preview(self, file_path, max_size=1024*1024):
        """Get preview for a file."""
        for handler in self.handlers:
            if handler.can_handle(file_path):
                return handler.generate_preview(file_path, max_size)
        
        return 'error', 'No preview available for this file type', {}
    
    def can_preview(self, file_path):
        """Check if file can be previewed."""
        return any(handler.can_handle(file_path) for handler in self.handlers)

# === REFACTOR: Move to ui/main_window.py ===
# [This would be the main FileScoutApp class - too large to annotate here]

# === REFACTOR: Move to main.py ===
if __name__ == '__main__':
    # === REFACTOR: Move to utils/config.py ===
    # Set application ID for proper Windows integration
    if sys.platform == "win32":
        import ctypes
        myappid = 'windsurfai.filescout.3.2'
        ctypes.windll.shell32.SetCurrentProcessExplicitAppUserModelID(myappid)
    
    # === REFACTOR: Move to utils/config.py ===
    # Enable high DPI support
    QApplication.setHighDpiScaleFactorRoundingPolicy(
        Qt.HighDpiScaleFactorRoundingPolicy.RoundPreferFloor
    )
    
    app = QApplication(sys.argv)
    app.setApplicationName(APP_NAME)
    app.setApplicationVersion(APP_VERSION)
    app.setOrganizationName(SETTINGS_ORG)
    app.setOrganizationDomain(SETTINGS_APP)
    
    # === REFACTOR: Move to utils/themes.py ===
    # Set application style
    app.setStyle('Fusion')
    
    # Create and show main window
    window = FileScoutApp()
    window.show()
    
    sys.exit(app.exec())
